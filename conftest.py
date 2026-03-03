"""Shared pytest fixtures for Memo Automator tests."""
import os
import tempfile

import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def tmp_dir():
    """Provide a temporary directory that is cleaned up after the test."""
    with tempfile.TemporaryDirectory() as d:
        yield d


@pytest.fixture
def sample_pptx(tmp_dir):
    """Create a minimal PPTX with a 4x3 table for testing.

    Returns the path to the saved file.
    """
    path = os.path.join(tmp_dir, "test_memo.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a table: 4 rows x 3 columns
    table_shape = slide.shapes.add_table(
        4, 3, Inches(0.5), Inches(1.0), Inches(8.0), Inches(3.0),
    )
    table_shape.name = "UnitMixTable"
    table = table_shape.table

    # Header row
    table.cell(0, 0).text = "Unit Type"
    table.cell(0, 1).text = "Beds"
    table.cell(0, 2).text = "Rent"

    # Data rows
    table.cell(1, 0).text = "1BR"
    table.cell(1, 1).text = "120"
    table.cell(1, 2).text = "$1,825"

    table.cell(2, 0).text = "3BR"
    table.cell(2, 1).text = "200"
    table.cell(2, 2).text = "$1,345"

    table.cell(3, 0).text = "4BR/2BA"
    table.cell(3, 1).text = "150"
    table.cell(3, 2).text = "$1,250"

    # Narrative text shape for text update tests
    text_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(8.0), Inches(1.0)
    )
    text_box.name = "NarrativeBox"
    text_box.text_frame.text = "IRR is 5.0% and units are 120."

    prs.save(path)
    return path


@pytest.fixture
def sample_config(tmp_dir):
    """Create a minimal config.yaml for testing. Returns the path."""
    path = os.path.join(tmp_dir, "config.yaml")
    with open(path, "w") as f:
        f.write(
            "proforma:\n"
            "  tabs:\n"
            '    - "Executive Summary"\n'
            '    - "Cash Flow"\n'
            "  max_rows_per_tab: 100\n"
            "  max_cols_per_tab: 20\n"
            "memo:\n"
            '  pages: "all"\n'
            "claude:\n"
            '  model: "claude-sonnet-4-6"\n'
            "  max_tokens: 8000\n"
            "  temperature: 0\n"
        )
    return path


@pytest.fixture
def sample_proforma_xlsx(tmp_dir):
    """Create a synthetic proforma workbook with two relevant tabs."""
    path = os.path.join(tmp_dir, "proforma.xlsx")
    wb = openpyxl.Workbook()

    ws1 = wb.active
    ws1.title = "Executive Summary"
    ws1["A1"] = "Metric"
    ws1["B1"] = "Value"
    ws1["A2"] = "IRR"
    ws1["B2"] = "6.5%"
    ws1["A3"] = "Units"
    ws1["B3"] = 130

    ws2 = wb.create_sheet("Cash Flow")
    ws2["A1"] = "Year"
    ws2["B1"] = "NOI"
    ws2["A2"] = 2026
    ws2["B2"] = 1200000

    wb.save(path)
    wb.close()
    return path


@pytest.fixture
def empty_proforma_xlsx(tmp_dir):
    """Create a proforma with expected tab names but no data rows."""
    path = os.path.join(tmp_dir, "empty_proforma.xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Executive Summary"
    wb.create_sheet("Cash Flow")
    wb.save(path)
    wb.close()
    return path


@pytest.fixture
def layout_test_pptx(tmp_dir):
    """
    Create a deck with two content slides and intentionally misaligned
    elements so normalize_layout() has work to do.
    """
    path = os.path.join(tmp_dir, "layout_test.pptx")
    prs = Presentation()

    # Cover slide (ignored by normalize_layout title alignment)
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Cover"

    # Baseline content slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Table of Contents"
    slide1.placeholders[1].text = "Section A"

    # Outlier title position + out-of-margin shape
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Financials"
    title2 = slide2.shapes.title
    title2.left = title2.left + Inches(1.0)

    off_margin = slide2.shapes.add_textbox(
        Inches(0.1), Inches(0.1), Inches(4.0), Inches(1.0)
    )
    off_margin.name = "OffMargin"
    off_margin.text_frame.text = "This box starts outside configured margins."

    prs.save(path)
    return path
