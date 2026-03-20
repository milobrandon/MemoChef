#!/usr/bin/env python3
"""
Memo Automator
==============
Automatically updates an Investment Committee (IC) PowerPoint memo with
metrics from an Excel proforma.

Workflow
--------
1. Create a backup of the original memo.
2. Extract data from specified proforma tabs (openpyxl, data_only).
3. Extract text / tables from ALL slides in the memo (python-pptx).
4. Send both datasets to the Claude API, which identifies every metric that
   should be updated and returns structured JSON mappings.
5. Apply the text / table updates to the memo (python-pptx).
6. Save the updated memo and write a detailed change-log.

Usage
-----
    python memo_automator.py <memo.pptx> <proforma.xlsx>
    python memo_automator.py <memo.pptx> <proforma.xlsx> --config my_config.yaml
"""

# ============================================================================
# IMPORTS
# ============================================================================
import argparse
import errno
import json
import logging
import os
import re
import shutil
import sys
import time
import zipfile
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Literal

import anthropic
import openpyxl
import yaml
from dotenv import load_dotenv
from openpyxl.utils.exceptions import InvalidFileException
from pydantic import BaseModel, ConfigDict, Field, ValidationError, model_validator
from pptx import Presentation
from pptx.exc import PackageNotFoundError

# ============================================================================
# LOGGING SETUP
# ============================================================================
LOG_FMT = "%(asctime)s  %(levelname)-8s  %(message)s"
logging.basicConfig(level=logging.INFO, format=LOG_FMT)
log = logging.getLogger("memo_automator")

_AUTH_ERRORS = tuple(
    e for e in (
        getattr(anthropic, "AuthenticationError", None),
        getattr(anthropic, "PermissionDeniedError", None),
    ) if e is not None
)
_RATE_LIMIT_ERRORS = tuple(
    e for e in (getattr(anthropic, "RateLimitError", None),) if e is not None
)
_TIMEOUT_ERRORS = tuple(
    e for e in (
        getattr(anthropic, "APITimeoutError", None),
        getattr(anthropic, "APIConnectionError", None),
    ) if e is not None
)
_API_STATUS_ERRORS = tuple(
    e for e in (getattr(anthropic, "APIStatusError", None),) if e is not None
)
_ALL_API_ERRORS = _AUTH_ERRORS + _RATE_LIMIT_ERRORS + _TIMEOUT_ERRORS + _API_STATUS_ERRORS


def _is_api_error(err: Exception) -> bool:
    return bool(_ALL_API_ERRORS) and isinstance(err, _ALL_API_ERRORS)


def _exit_with_api_error(err: Exception):
    if _AUTH_ERRORS and isinstance(err, _AUTH_ERRORS):
        log.error(
            "Claude API authentication failed. Check ANTHROPIC_API_KEY in .env or "
            "your deployment secrets."
        )
    elif _RATE_LIMIT_ERRORS and isinstance(err, _RATE_LIMIT_ERRORS):
        log.error(
            "Claude API rate limit exceeded. Wait ~60 seconds and retry, or reduce "
            "batch size/token usage."
        )
    elif _TIMEOUT_ERRORS and isinstance(err, _TIMEOUT_ERRORS):
        log.error(
            "Claude API request timed out or connection failed. Retry and, for large "
            "decks, consider reducing pages per batch."
        )
    elif _API_STATUS_ERRORS and isinstance(err, _API_STATUS_ERRORS):
        status_code = getattr(err, "status_code", None)
        if status_code is None and hasattr(err, "response"):
            status_code = getattr(err.response, "status_code", None)
        if status_code:
            log.error("Claude API returned HTTP %s. Please retry.", status_code)
        else:
            log.error("Claude API returned an error. Please retry.")
    else:
        log.error("Claude API call failed: %s", err)
    sys.exit(1)


def _exit_with_os_error(err: OSError, action: str):
    if err.errno == errno.ENOSPC:
        log.error("Disk full while %s. Free up space and retry.", action)
    elif err.errno in (errno.EACCES, errno.EPERM):
        log.error("Permission denied while %s. Check file/folder permissions.", action)
    else:
        log.error("File system error while %s: %s", action, err)
    sys.exit(1)


def _load_presentation(memo_path: str):
    try:
        return Presentation(memo_path)
    except (PackageNotFoundError, zipfile.BadZipFile, KeyError) as e:
        raise ValueError(
            f"Unable to open memo PPTX '{memo_path}'. The file may be corrupt, "
            f"not a valid .pptx, or password-protected."
        ) from e


# ============================================================================
# 1. ARGUMENT PARSING
# ============================================================================
def parse_args():
    """Parse command-line arguments."""
    p = argparse.ArgumentParser(
        description="Update an IC memo with metrics from a proforma."
    )
    p.add_argument("memo", help="Path to the PowerPoint memo (.pptx)")
    p.add_argument("proforma", help="Path to the Excel proforma (.xlsx / .xlsm)")
    p.add_argument(
        "--config", "-c",
        default=os.path.join(os.path.dirname(__file__), "config.yaml"),
        help="Path to YAML config file (default: config.yaml beside this script)",
    )
    p.add_argument(
        "--output-dir", "-o",
        default=None,
        help="Directory for output artifacts (default: same folder as memo)",
    )
    p.add_argument(
        "--dry-run",
        action="store_true",
        help="Show what would change without modifying the memo",
    )
    p.add_argument(
        "--skip-validation",
        action="store_true",
        help="Skip the Claude validation pass for faster runs",
    )
    p.add_argument(
        "--property-name",
        default="",
        help="Property name as shown in the proforma (helps match rebranded names)",
    )
    verbosity = p.add_mutually_exclusive_group()
    verbosity.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Enable DEBUG-level logging for detailed output",
    )
    verbosity.add_argument(
        "--quiet", "-q",
        action="store_true",
        help="Suppress all output except warnings and errors",
    )
    return p.parse_args()


# ============================================================================
# 2. CONFIGURATION
# ============================================================================
class ProformaConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")
    tabs: list[str] = Field(
        default_factory=lambda: [
            "Executive Summary",
            "Development Summary",
            "Cash Flow",
            "Assumptions",
            "Proforma Comparison",
        ],
        min_length=1,
    )
    max_rows_per_tab: int = Field(default=250, ge=0)
    max_cols_per_tab: int = Field(default=30, ge=0)


class MemoConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")
    pages: Literal["all"] | list[int] = "all"


class ScheduleConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")
    max_tasks: int = Field(default=500, ge=0)


class BrandingConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")
    theme_path: str = ""
    heading_size_threshold: int = Field(default=18, ge=0)
    color_distance_threshold: float = Field(default=80, ge=0)


class LayoutConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")
    margin_left: float = Field(default=0.50, ge=0)
    margin_right: float = Field(default=0.50, ge=0)
    margin_top: float = Field(default=0.25, ge=0)
    margin_bottom: float = Field(default=0.50, ge=0)
    snap_tolerance: float = Field(default=0.05, ge=0)


class ClaudeConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")
    model: str = Field(default="claude-sonnet-4-6", min_length=1)
    validation_model: str | None = None
    max_tokens: int = Field(default=16000, ge=1)
    temperature: float = Field(default=0, ge=0, le=1)

    @model_validator(mode="after")
    def set_validation_model_default(self):
        if not self.validation_model:
            self.validation_model = self.model
        return self


class AppConfig(BaseModel):
    model_config = ConfigDict(extra="forbid")
    proforma: ProformaConfig = Field(default_factory=ProformaConfig)
    memo: MemoConfig = Field(default_factory=MemoConfig)
    schedule: ScheduleConfig = Field(default_factory=ScheduleConfig)
    branding: BrandingConfig = Field(default_factory=BrandingConfig)
    layout: LayoutConfig = Field(default_factory=LayoutConfig)
    claude: ClaudeConfig = Field(default_factory=ClaudeConfig)


def _format_validation_error(e: dict) -> str:
    location = ".".join(str(p) for p in e.get("loc", ()))
    msg = e.get("msg", "invalid value")
    return f"{location}: {msg}" if location else msg


def _validate_config(cfg: dict) -> list[str]:
    """Validate config schema. Returns a list of error messages (empty = valid)."""
    try:
        AppConfig.model_validate(cfg)
        return []
    except ValidationError as exc:
        return [_format_validation_error(e) for e in exc.errors()]


def load_config(config_path: str) -> dict:
    """Load and validate the YAML configuration file."""
    with open(config_path, "r", encoding="utf-8") as f:
        raw = yaml.safe_load(f)

    if raw is None:
        raw = {}
    if not isinstance(raw, dict):
        raise ValueError("Invalid config:\n  root: expected a YAML object")

    errors = _validate_config(raw)
    if errors:
        raise ValueError("Invalid config:\n  " + "\n  ".join(errors))

    cfg = AppConfig.model_validate(raw)
    return cfg.model_dump(mode="python")


# ============================================================================
# 3. BACKUP
# ============================================================================
def create_backup(memo_path: str, output_dir: str) -> str:
    """Copy the original memo to a timestamped backup file."""
    stem = Path(memo_path).stem
    ext = Path(memo_path).suffix
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_name = f"{stem}_BACKUP_{ts}{ext}"
    backup_path = os.path.join(output_dir, backup_name)
    shutil.copy2(memo_path, backup_path)
    log.info("Backup created: %s", backup_path)
    return backup_path


# ============================================================================
# 4. PROFORMA DATA EXTRACTION  (openpyxl, data_only=True)
# ============================================================================
def extract_proforma_data(proforma_path: str, cfg: dict) -> str:
    """
    Read the proforma workbook and return a compact text representation
    of every non-empty cell on the specified tabs.

    Uses data_only=True so formulas resolve to their cached values.
    """
    tabs = cfg["proforma"]["tabs"]
    max_rows = cfg["proforma"]["max_rows_per_tab"]
    max_cols = cfg["proforma"]["max_cols_per_tab"]

    log.info("Opening proforma (data_only): %s", proforma_path)
    try:
        wb = openpyxl.load_workbook(proforma_path, data_only=True)
    except (InvalidFileException, zipfile.BadZipFile) as e:
        raise ValueError(
            f"Unable to open proforma '{proforma_path}'. The file appears malformed "
            f"or is not a valid .xlsx/.xlsm workbook."
        ) from e
    log.info("Available sheets: %s", wb.sheetnames)

    lines = []
    found_tabs = 0
    data_rows = 0
    for tab_name in tabs:
        if tab_name not in wb.sheetnames:
            log.warning("Tab '%s' not found in proforma - skipping", tab_name)
            continue
        found_tabs += 1
        ws = wb[tab_name]
        lines.append(f"\n{'='*70}")
        lines.append(f"TAB: {tab_name}")
        lines.append(f"{'='*70}")

        end_row = ws.max_row if max_rows == 0 else min(ws.max_row, max_rows)
        end_col = ws.max_column if max_cols == 0 else min(ws.max_column, max_cols)

        for row in ws.iter_rows(
            min_row=1, max_row=end_row, max_col=end_col, values_only=False
        ):
            row_data = []
            for cell in row:
                if cell.value is not None:
                    row_data.append(str(cell.value))
            if row_data:
                lines.append(f"Row {row[0].row}:\t" + "\t".join(row_data))
                data_rows += 1

    wb.close()

    if found_tabs == 0:
        raise ValueError(
            f"No configured tabs found in proforma. "
            f"Expected tabs: {tabs}. Available sheets: {wb.sheetnames}"
        )

    if data_rows == 0:
        raise ValueError(
            "Proforma extraction found no non-empty values in configured tabs. "
            "If this workbook contains formulas, open it in Excel, let it "
            "recalculate, save, and retry so cached values are available."
        )

    proforma_text = "\n".join(lines)
    log.info(
        "Proforma extraction complete (%d lines, %d chars)",
        len(lines), len(proforma_text)
    )
    return proforma_text


# ============================================================================
# 4a. MARKET DATA EXTRACTION  (openpyxl, data_only=True)
# ============================================================================
_MARKET_DASHBOARD_TABS = [
    "Tables",
    "Comparison Graph",
    "Uncaptured Demand Comparison",
    "Rent Growth Comparison By Year",
    "Occupancy Comparison By Year",
    "Comp Set",
]


def extract_market_data(market_data_path: str, cfg: dict) -> str:
    """
    Read the RealPage market data workbook and return a compact text
    representation of the 6 dashboard tabs (ignoring back-end data tabs).

    Uses data_only=True so formulas resolve to their cached values.
    Returns empty string if no dashboard tabs are found (non-fatal).
    """
    max_rows = cfg["proforma"]["max_rows_per_tab"]
    max_cols = cfg["proforma"]["max_cols_per_tab"]

    log.info("Opening market data (data_only): %s", market_data_path)
    try:
        wb = openpyxl.load_workbook(market_data_path, data_only=True)
    except (InvalidFileException, zipfile.BadZipFile) as e:
        log.warning(
            "Unable to open market data '%s': %s. Continuing without market data.",
            market_data_path, e,
        )
        return ""
    except FileNotFoundError:
        log.warning("Market data file not found: %s", market_data_path)
        return ""
    log.info("Market data sheets: %s", wb.sheetnames)

    lines = [
        f"\n{'='*70}",
        "MARKET DATA (from RealPage)",
        f"{'='*70}",
    ]
    found_tabs = 0
    data_rows = 0
    for tab_name in _MARKET_DASHBOARD_TABS:
        if tab_name not in wb.sheetnames:
            log.warning("Market data tab '%s' not found - skipping", tab_name)
            continue
        found_tabs += 1
        ws = wb[tab_name]
        lines.append(f"\n{'='*70}")
        lines.append(f"TAB: {tab_name}")
        lines.append(f"{'='*70}")

        end_row = ws.max_row if max_rows == 0 else min(ws.max_row, max_rows)
        end_col = ws.max_column if max_cols == 0 else min(ws.max_column, max_cols)

        for row in ws.iter_rows(
            min_row=1, max_row=end_row, max_col=end_col, values_only=False
        ):
            row_data = []
            for cell in row:
                if cell.value is not None:
                    row_data.append(str(cell.value))
            if row_data:
                lines.append(f"Row {row[0].row}:\t" + "\t".join(row_data))
                data_rows += 1

    wb.close()

    if found_tabs == 0:
        # Fallback: scan ALL tabs for market-relevant data.
        # Pick tabs whose names contain keywords like rent, sales, comp,
        # pipeline, supply, demand, occupancy, market, land.
        _market_keywords = [
            "rent", "sale", "comp", "pipeline", "supply", "demand",
            "occupancy", "market", "land", "prelease", "sbys", "side",
            "growth", "rate", "unit mix", "taxes",
        ]
        fallback_tabs = []
        for sn in wb.sheetnames:
            sn_lower = sn.lower()
            if any(kw in sn_lower for kw in _market_keywords):
                fallback_tabs.append(sn)

        if fallback_tabs:
            log.info(
                "No configured dashboard tabs matched. Falling back to %d "
                "market-keyword tabs: %s",
                len(fallback_tabs), fallback_tabs[:10],
            )
            # Re-open since wb was already used
            wb2 = openpyxl.load_workbook(market_data_path, data_only=True)
            # Prioritize tabs with rent/comp/market keywords and recent dates.
            # Score each tab: high-value keywords + recency.
            _high_value = ["rent", "sbys", "side", "comp", "sale", "market", "unit mix"]
            def _tab_score(name: str) -> int:
                nl = name.lower()
                score = sum(2 for kw in _high_value if kw in nl)
                # Bonus for recent dates (2026, 2025)
                if "2026" in name:
                    score += 3
                elif "2025" in name:
                    score += 1
                return score
            fallback_tabs.sort(key=_tab_score, reverse=True)
            for tab_name in fallback_tabs[:4]:  # cap at 4 tabs
                ws = wb2[tab_name]
                found_tabs += 1
                lines.append(f"\n{'='*70}")
                lines.append(f"TAB: {tab_name}")
                lines.append(f"{'='*70}")

                end_row = ws.max_row if max_rows == 0 else min(ws.max_row, max_rows)
                end_col = ws.max_column if max_cols == 0 else min(ws.max_column, max_cols)

                for row in ws.iter_rows(
                    min_row=1, max_row=end_row, max_col=end_col, values_only=False
                ):
                    row_data = []
                    for cell in row:
                        if cell.value is not None:
                            row_data.append(str(cell.value))
                    if row_data:
                        lines.append(f"Row {row[0].row}:\t" + "\t".join(row_data))
                        data_rows += 1
            wb2.close()
        else:
            log.warning(
                "No dashboard tabs found in market data file. "
                "Expected tabs: %s. Available: %s. Skipping market data.",
                _MARKET_DASHBOARD_TABS, wb.sheetnames,
            )
            return ""

    if data_rows == 0:
        log.warning(
            "Market data extraction found no non-empty values. "
            "If this workbook contains formulas, open it in Excel, let it "
            "recalculate, save, and retry."
        )
        return ""

    market_text = "\n".join(lines)
    log.info(
        "Market data extraction complete (%d tabs, %d lines, %d chars)",
        found_tabs, len(lines), len(market_text),
    )
    return market_text

# ============================================================================
# 4b. SCHEDULE DATA EXTRACTION  (mpxj via jpype)
# ============================================================================
def _ensure_jvm():
    """Start the JVM once for mpxj access. No-op if already running."""
    import jpype
    if jpype.isJVMStarted():
        return

    # Auto-discover JAVA_HOME if not set
    if not os.environ.get("JAVA_HOME"):
        search_dirs = [
            os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "Microsoft"),
            os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "Java"),
            os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "Eclipse Adoptium"),
        ]
        for search_dir in search_dirs:
            if os.path.isdir(search_dir):
                for entry in sorted(os.listdir(search_dir), reverse=True):
                    if "jdk" in entry.lower():
                        candidate = os.path.join(search_dir, entry)
                        if os.path.isdir(candidate):
                            os.environ["JAVA_HOME"] = candidate
                            log.info("Auto-discovered JAVA_HOME: %s", candidate)
                            break
            if os.environ.get("JAVA_HOME"):
                break

    # Import mpxj for its classpath side effects before starting the JVM.
    import mpxj  # noqa: F401

    jpype.startJVM()
    log.info("JVM started with classpath from mpxj")


def extract_schedule_data(schedule_path: str, cfg: dict) -> str:
    """
    Read a Microsoft Project (.mpp) schedule and return a hierarchical text
    representation of tasks with dates and durations.

    Uses mpxj (via jpype) to parse the .mpp file.
    """
    _ensure_jvm()

    import jpype
    from java.io import File as JavaFile

    max_tasks = cfg.get("schedule", {}).get("max_tasks", 500)

    log.info("Opening schedule: %s", schedule_path)
    reader = jpype.JClass("org.mpxj.reader.UniversalProjectReader")()
    project = reader.read(JavaFile(schedule_path))

    lines = []
    lines.append(f"\n{'='*70}")
    lines.append("SCHEDULE DATA (from Microsoft Project)")
    lines.append(f"{'='*70}")

    task_count = 0
    for task in project.getTasks():
        if task_count >= max_tasks:
            lines.append(f"... (truncated at {max_tasks} tasks)")
            break

        name = str(task.getName()) if task.getName() else ""
        # Skip L0 unnamed separator tasks (grouping containers)
        outline_level = task.getOutlineLevel()
        if outline_level is not None:
            level = int(str(outline_level))
        else:
            level = 0
        if level == 0 and not name.strip():
            continue

        # Get dates and duration
        start = task.getStart()
        finish = task.getFinish()
        duration = task.getDuration()

        start_str = str(start).split("T")[0] if start else "N/A"
        finish_str = str(finish).split("T")[0] if finish else "N/A"

        if duration:
            dur_str = str(duration)
        else:
            dur_str = "0d"

        # Milestone detection
        is_milestone = task.getMilestone() if task.getMilestone() is not None else False
        milestone_tag = "  [MILESTONE]" if is_milestone else ""

        indent = "  " * max(level - 1, 0)
        level_tag = f"[L{level}]"

        lines.append(
            f"{indent}{level_tag} {name}{milestone_tag}  |  "
            f"Start: {start_str}  |  Finish: {finish_str}  |  Dur: {dur_str}"
        )
        task_count += 1

    schedule_text = "\n".join(lines)
    log.info("Schedule extraction complete (%d tasks, %d chars)",
             task_count, len(schedule_text))
    return schedule_text


# ============================================================================
# 5. MEMO CONTENT EXTRACTION  (python-pptx)
# ============================================================================
def extract_memo_content(memo_path: str, cfg: dict) -> str:
    """
    Read the PowerPoint memo and return a structured text representation
    of shapes (tables, text boxes) on the target pages.

    If cfg["memo"]["pages"] is "all", scans every slide in the deck.
    Otherwise expects a list of 1-based page numbers.
    """
    prs = _load_presentation(memo_path)
    total_slides = len(prs.slides)

    # Determine which pages to scan
    pages_cfg = cfg["memo"]["pages"]
    if pages_cfg == "all":
        page_numbers = list(range(1, total_slides + 1))
    else:
        page_numbers = [int(p) for p in pages_cfg]

    log.info("Scanning %d pages (out of %d total slides)", len(page_numbers), total_slides)

    lines = []
    for page_num in page_numbers:
        idx = page_num - 1
        if idx >= total_slides:
            log.warning("Page %d does not exist (only %d slides)", page_num, total_slides)
            continue

        slide = prs.slides[idx]
        lines.append(f"\n{'='*70}")
        lines.append(f"PAGE {page_num}  (slide index {idx})")
        lines.append(f"{'='*70}")

        for si, shape in enumerate(slide.shapes):
            lines.append(f"\n--- Shape {si}: type={shape.shape_type}, "
                         f"name='{shape.name}' ---")
            lines.append(f"    Position: left={shape.left}, top={shape.top}, "
                         f"width={shape.width}, height={shape.height}")

            # Text frames (text boxes, placeholders)
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    text = para.text.strip()
                    if text:
                        lines.append(f"    Para {pi}: '{text}'")

            # Tables
            if shape.has_table:
                tbl = shape.table
                lines.append(f"    Table: {len(tbl.rows)} rows x "
                             f"{len(tbl.columns)} cols")
                for ri, row in enumerate(tbl.rows):
                    cells = []
                    for ci, cell in enumerate(row.cells):
                        ct = cell.text.strip()
                        if ct:
                            cells.append(f"[{ci}]={ct}")
                    if cells:
                        lines.append(f"    Row {ri}: {' | '.join(cells)}")

            # Charts (embedded Excel chart objects)
            if shape.has_chart:
                try:
                    chart = shape.chart
                except (KeyError, AttributeError):
                    lines.append(f"    Chart: (relationship broken — skipped)")
                    continue
                chart_type_name = str(chart.chart_type) if chart.chart_type else "UNKNOWN"
                lines.append(f"    Chart type: {chart_type_name}")

                # Chart title
                if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
                    title_text = chart.chart_title.text_frame.text.strip()
                    lines.append(f"    Chart title: '{title_text}'")

                # Extract series and data
                try:
                    for s_idx, series in enumerate(chart.series):
                        s_name = ""
                        try:
                            s_name = series.name if series.name else f"Series {s_idx}"
                        except (AttributeError, IndexError):
                            s_name = f"Series {s_idx}"

                        # Extract series values
                        vals = []
                        try:
                            if series.values:
                                vals = [v for v in series.values if v is not None]
                        except Exception:
                            pass

                        lines.append(f"    Series {s_idx} ('{s_name}'): {vals[:20]}")
                except Exception as e:
                    lines.append(f"    (chart data extraction failed: {e})")

                # Extract category labels (x-axis)
                try:
                    plot = chart.plots[0]
                    if plot.categories:
                        cats = list(plot.categories)[:30]
                        lines.append(f"    Categories: {cats}")
                except Exception:
                    pass

    memo_text = "\n".join(lines)
    log.info("Memo extraction complete (%d lines, %d chars)",
             len(lines), len(memo_text))
    return memo_text


# ============================================================================
# 6. MEMO CONTENT CHUNKING
# ============================================================================
def chunk_memo_by_pages(memo_content: str, pages_per_chunk: int = 10) -> list:
    """
    Split memo content into chunks of up to pages_per_chunk pages each.

    Used when the full prompt would exceed the model's output token limit.
    Each chunk is processed in a separate API call and the results are merged.
    """
    # Each page block starts with the === PAGE N === header
    page_blocks = re.split(r"(?=\n={60,}\nPAGE \d+)", memo_content)
    chunks = []
    for i in range(0, len(page_blocks), pages_per_chunk):
        chunk = "".join(page_blocks[i:i + pages_per_chunk])
        if chunk.strip():
            chunks.append(chunk)
    return chunks


PROMPTS_DIR = Path(__file__).resolve().parent / "prompts"


def _load_prompt_template(filename: str) -> str:
    """
    Load a prompt template from prompts/filename.
    Raises FileNotFoundError if the file is missing or empty.
    """
    path = PROMPTS_DIR / filename
    text = path.read_text(encoding="utf-8").strip()
    if not text:
        raise FileNotFoundError(f"Prompt template is empty: {path}")
    return text


# ============================================================================
# 6b. FINAL REVIEW SIGN-OFF
# ============================================================================
FINAL_REVIEW_PROMPT = _load_prompt_template("final_review_v1.txt")


def run_final_review(
    client: "anthropic.Anthropic",
    proforma_data: str,
    memo_content: str,
    cfg: dict,
    max_rounds: int = 2,
) -> dict:
    """Final QA gate: Claude reviews the memo as a distributable document.

    Loops up to max_rounds: review → apply critical fixes → re-review.
    Returns the final review result with verdict, scores, and any
    remaining warnings.
    """
    model = cfg["claude"].get("validation_model", cfg["claude"]["model"])
    max_tokens = cfg["claude"]["max_tokens"]

    system_text = FINAL_REVIEW_PROMPT.format(
        proforma_data=proforma_data,
        memo_content="(see user message below)",
    )

    review_result = None
    for attempt in range(1, max_rounds + 1):
        user_text = f"## Updated Memo Content (final state)\n{memo_content}"
        if attempt > 1:
            user_text += (
                "\n\n## NOTE: This is review round {attempt}. Previous critical fixes "
                "have been applied. Re-evaluate the memo from scratch."
            )

        try:
            message = client.messages.create(
                model=model,
                max_tokens=max_tokens,
                temperature=0,
                system=[{
                    "type": "text",
                    "text": system_text,
                    "cache_control": {"type": "ephemeral"},
                }],
                messages=[{"role": "user", "content": user_text}],
            )

            raw = ""
            for block in message.content:
                if block.type == "text":
                    raw = block.text
                    break

            review_result = _parse_json_response(raw)
            if review_result is None:
                log.warning("Final review round %d: could not parse JSON", attempt)
                continue

            verdict = review_result.get("verdict", "REVISIONS_NEEDED")
            score = review_result.get("overall_score", 0)
            critical_fixes = review_result.get("critical_fixes", [])

            log.info(
                "Final review round %d: verdict=%s, score=%d, critical_fixes=%d",
                attempt, verdict, score, len(critical_fixes),
            )

            if verdict == "APPROVED" or not critical_fixes:
                return review_result

            # Return fixes for the caller to apply (pipeline handles apply + re-review)
            return review_result

        except Exception as e:
            log.warning("Final review round %d failed: %s", attempt, e)
            continue

    # All rounds exhausted
    return review_result or {
        "verdict": "REVISIONS_NEEDED",
        "overall_score": 0,
        "categories": {},
        "critical_fixes": [],
        "warnings": ["Final review could not be completed"],
        "summary": "Review failed after all attempts",
    }


# ============================================================================
# 6c. POST-APPLY CONSISTENCY CHECK
# ============================================================================
CONSISTENCY_PROMPT = _load_prompt_template("consistency_check_v1.txt")


def run_consistency_check(
    client: "anthropic.Anthropic",
    proforma_data: str,
    memo_content: str,
    changes: list[dict],
    cfg: dict,
    max_retries: int = 2,
) -> dict:
    """Re-read the updated memo and verify every metric ties out.

    Returns a dict with 'status' ('pass' or 'fail'), 'discrepancies' list,
    and 'fixes' list of auto-applicable corrections.
    """
    model = cfg["claude"].get("validation_model", cfg["claude"]["model"])
    max_tokens = cfg["claude"]["max_tokens"]

    # Build a compact summary of changes already applied
    changes_lines = []
    for c in changes[:50]:  # cap to avoid prompt bloat
        changes_lines.append(
            f"  page {c.get('page', '?')} [{c.get('type', '?')}]: "
            f"'{c.get('old', '')[:30]}' -> '{c.get('new', '')[:30]}'"
        )
    changes_summary = "\n".join(changes_lines) or "(no changes applied)"

    system_text = CONSISTENCY_PROMPT.format(
        proforma_data=proforma_data,
        memo_content=memo_content,
        changes_summary=changes_summary,
    )

    for attempt in range(1, max_retries + 1):
        try:
            message = client.messages.create(
                model=model,
                max_tokens=max_tokens,
                temperature=0,
                system=[{
                    "type": "text",
                    "text": system_text,
                    "cache_control": {"type": "ephemeral"},
                }],
                messages=[{"role": "user", "content": "Perform the consistency check now."}],
            )

            raw = ""
            for block in message.content:
                if block.type == "text":
                    raw = block.text
                    break

            result = _parse_json_response(raw)
            if result is None:
                if attempt < max_retries:
                    log.warning("Consistency check attempt %d: invalid JSON, retrying...", attempt)
                    continue
                return {"status": "error", "discrepancies": [], "summary": "Could not parse check results"}

            return result

        except Exception as e:
            if attempt < max_retries:
                log.warning("Consistency check attempt %d failed: %s", attempt, e)
                continue
            log.error("Consistency check failed after %d attempts: %s", max_retries, e)
            return {"status": "error", "discrepancies": [], "summary": str(e)}

    return {"status": "error", "discrepancies": [], "summary": "All attempts failed"}


# ============================================================================
# 7. CLAUDE API - METRIC MAPPING
# ============================================================================
MAPPING_PROMPT = _load_prompt_template("mapping_v1.txt")


def _salvage_truncated_json(raw: str) -> dict | None:
    """
    Attempt to recover valid mappings from a truncated JSON response.

    When Claude hits max_tokens mid-JSON, the response ends abruptly.
    This function tries to close the JSON by finding the last complete
    entry boundary (a '},' or '}]' pattern) and appending the missing
    closing brackets.

    Returns a valid mappings dict if salvageable, None otherwise.
    """
    text = raw.strip()
    if not text:
        return None

    # Strip markdown fences if present
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)

    # Must start with a JSON object
    start = text.find("{")
    if start == -1:
        return None
    text = text[start:]

    # Strategy: try a few common closing patterns to complete the JSON.
    # The response is typically:
    #   {"table_updates": [..., {last_complete}, {partial   <-- cut here
    # or:
    #   {"table_updates": [...], "text_updates": [..., {partial   <-- cut here

    # Find the last complete entry: look for the last '}' that's followed
    # by ',' or ']' or is at a natural boundary
    # Walk backwards to find last complete "}" entry
    last_good = -1
    depth = 0
    for i in range(len(text) - 1, -1, -1):
        ch = text[i]
        if ch == '}':
            if depth == 0:
                last_good = i
                break
            depth -= 1
        elif ch == '{':
            depth += 1

    if last_good == -1:
        return None

    truncated = text[:last_good + 1]

    # Count unmatched brackets to determine what closings are needed
    open_braces = truncated.count('{') - truncated.count('}')
    open_brackets = truncated.count('[') - truncated.count(']')

    if open_braces < 0 or open_brackets < 0:
        return None

    # Build closing sequence: close all open brackets then braces
    closing = ']' * open_brackets + '}' * open_braces
    candidate = truncated + closing

    try:
        mappings = json.loads(candidate)
        mappings.setdefault("table_updates", [])
        mappings.setdefault("text_updates", [])
        mappings.setdefault("row_inserts", [])
        mappings.setdefault("narrative_updates", [])
        mappings.setdefault("table_structure_updates", [])
        n = (len(mappings["table_updates"]) + len(mappings["text_updates"])
             + len(mappings["row_inserts"]) + len(mappings["narrative_updates"])
             + len(mappings["table_structure_updates"]))
        if n > 0:
            log.info("Salvaged %d updates from truncated response", n)
            return mappings
        return None
    except json.JSONDecodeError:
        pass

    # Fallback: try trimming back to last '},' boundary (complete array entry)
    last_comma_boundary = truncated.rfind('},')
    if last_comma_boundary == -1:
        return None

    truncated2 = truncated[:last_comma_boundary + 1]
    open_braces2 = truncated2.count('{') - truncated2.count('}')
    open_brackets2 = truncated2.count('[') - truncated2.count(']')
    if open_braces2 < 0 or open_brackets2 < 0:
        return None

    closing2 = ']' * open_brackets2 + '}' * open_braces2
    candidate2 = truncated2 + closing2

    try:
        mappings = json.loads(candidate2)
        mappings.setdefault("table_updates", [])
        mappings.setdefault("text_updates", [])
        mappings.setdefault("row_inserts", [])
        n = len(mappings["table_updates"]) + len(mappings["text_updates"]) + len(mappings["row_inserts"])
        if n > 0:
            log.info("Salvaged %d updates from truncated response (fallback)", n)
            return mappings
        return None
    except json.JSONDecodeError:
        log.debug("Could not salvage truncated JSON response")
        return None


def _parse_json_response(raw: str) -> dict | None:
    """
    Parse a JSON object from a Claude API response, handling common
    noise: markdown fences, trailing commentary, whitespace.

    Returns the parsed dict, or None if no valid JSON could be extracted.

    Parsing strategy (in order):
    1. Strip markdown fences and whitespace.
    2. Return None for empty / null / trivially empty responses.
    3. Try json.loads() on the full text.
    4. Fall back to raw_decode() at the first '{' - log at DEBUG.
    5. Fall back to brace-matching (first '{' to last '}') - log at WARNING.
    6. Return None on total failure.
    """
    text = raw.strip()

    # Strip markdown fences
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)

    # Empty / trivial
    if not text or text in ("{}", "[]", "null"):
        return None

    # Strategy 1: direct parse
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # Strategy 2: raw_decode from first '{'
    start = text.find("{")
    if start != -1:
        try:
            decoder = json.JSONDecoder()
            obj, _ = decoder.raw_decode(text, start)
            log.debug("Parsed JSON via raw_decode (response had extra data)")
            return obj
        except json.JSONDecodeError:
            pass

        # Strategy 3: brace-matching (first '{' to last '}')
        end = text.rfind("}")
        if end != -1 and end > start:
            try:
                obj = json.loads(text[start:end + 1])
                log.warning("Extracted JSON using brace matching")
                return obj
            except json.JSONDecodeError:
                pass

    log.error("Failed to parse JSON from Claude response. "
              "First 500 chars: %s", text[:500])
    return None


def _create_message(client: anthropic.Anthropic, **api_kwargs):
    """Call the Claude API using streaming to avoid the SDK timeout for large
    max_tokens values.  Returns a full Message object identical to what
    ``client.messages.create()`` would return."""
    with client.messages.stream(**api_kwargs) as stream:
        return stream.get_final_message()


def build_mapping_batch_requests(
    proforma_data: str,
    memo_chunks: list[str],
    cfg: dict,
    property_name: str = "",
    source_directives: list[dict] | None = None,
) -> list[dict]:
    """Build a list of batch API request dicts for the Message Batches API.

    Each request contains a system message (cached instructions + proforma)
    and a user message (one memo chunk).  Returns dicts compatible with
    ``anthropic.types.messages.batch_create_params.Request``.
    """
    from anthropic.types.message_create_params import MessageCreateParamsNonStreaming
    from anthropic.types.messages.batch_create_params import Request

    model = cfg["claude"]["model"]
    max_tokens = cfg["claude"]["max_tokens"]
    temperature = cfg["claude"]["temperature"]

    pn_section = _property_name_section(property_name, "mapping")
    directives_section = format_source_directives(source_directives or [], scope="mapping")
    system_text = MAPPING_PROMPT.format(
        proforma_data=proforma_data,
        memo_content="(see user message below)",
        property_name_section=pn_section,
        source_directives_section=directives_section,
    )

    requests = []
    for i, chunk in enumerate(memo_chunks):
        user_text = f"## Memo Content (from PowerPoint)\n{chunk}"
        requests.append(
            Request(
                custom_id=f"mapping-chunk-{i}",
                params=MessageCreateParamsNonStreaming(
                    model=model,
                    max_tokens=max_tokens,
                    temperature=temperature,
                    system=[{
                        "type": "text",
                        "text": system_text,
                        "cache_control": {"type": "ephemeral", "ttl": "1h"},
                    }],
                    messages=[{"role": "user", "content": user_text}],
                ),
            )
        )
    log.info("Built %d batch requests (system_prefix=%d chars)", len(requests), len(system_text))
    return requests


def submit_and_poll_batch(
    client: anthropic.Anthropic,
    requests: list,
    poll_interval: int = 30,
    timeout: int = 3600,
) -> dict[str, dict]:
    """Submit a Message Batch, poll until done, return results keyed by custom_id.

    Returns a dict mapping custom_id -> parsed JSON mappings dict.
    Raises RuntimeError if the batch fails or times out.
    """
    batch = client.messages.batches.create(requests=requests)
    batch_id = batch.id
    log.info("Batch submitted: %s (%d requests)", batch_id, len(requests))

    start = time.time()
    while True:
        batch = client.messages.batches.retrieve(batch_id)
        status = batch.processing_status
        counts = batch.request_counts
        log.info("Batch %s: status=%s, succeeded=%d, errored=%d, expired=%d",
                 batch_id, status, counts.succeeded, counts.errored, counts.expired)

        if status == "ended":
            break
        if time.time() - start > timeout:
            raise RuntimeError(f"Batch {batch_id} timed out after {timeout}s")
        time.sleep(poll_interval)

    # Collect results
    results = {}
    for result in client.messages.batches.results(batch_id):
        cid = result.custom_id
        if result.result.type == "succeeded":
            message = result.result.message
            raw = ""
            for block in message.content:
                if block.type == "text":
                    raw = block.text
                    break
            mappings = _parse_json_response(raw)
            if mappings is None:
                log.warning("Batch result %s: unparseable response", cid)
                mappings = {"table_updates": [], "text_updates": [], "row_inserts": []}
            results[cid] = mappings
        else:
            log.warning("Batch result %s: %s", cid, result.result.type)
            results[cid] = {"table_updates": [], "text_updates": [], "row_inserts": []}

    log.info("Batch %s complete: %d results collected", batch_id, len(results))
    return results


def _property_name_section(property_name: str, purpose: str = "mapping") -> str:
    """Build the property-name targeting section shared by mapping & validation."""
    if not property_name:
        return ""
    if purpose == "validation":
        return (
            f"\n## Property Name - CRITICAL TARGETING CHECK\n"
            f"The proforma data corresponds to **\"{property_name}\"** in the memo. "
            f"The proforma's own internal name may differ (old/rebranded name).\n\n"
            f"Verify that:\n"
            f"- All updates target the \"{property_name}\" column/row, NOT other "
            f"  properties' columns/rows.\n"
            f"- The name \"{property_name}\" is NOT renamed to the proforma's "
            f"  internal name. Flag any mapping that renames it as REJECT.\n"
            f"- If any mapping targets a different property's column/row, REJECT it.\n"
        )
    return (
        f"\n## Property Name - CRITICAL TARGETING OVERRIDE\n"
        f"The proforma data corresponds to the property named **\"{property_name}\"** "
        f"in the memo. The proforma's own internal name may differ (e.g. an old or "
        f"rebranded name) - IGNORE the proforma's internal project name for targeting "
        f"purposes.\n\n"
        f"Apply all proforma data to the column, row, or section labeled "
        f"\"{property_name}\" in the memo:\n"
        f"- In **side-by-side comparison tables** (column-oriented), update the "
        f"  \"{property_name}\" column, NOT any other property's column.\n"
        f"- In **row-oriented tables** (comp summary, pipeline), update the row "
        f"  for \"{property_name}\".\n"
        f"- In **narrative text**, update metrics that describe \"{property_name}\".\n"
        f"- Do NOT rename \"{property_name}\" to the proforma's internal name. "
        f"  Keep the memo's name as-is.\n"
        f"- Do NOT update columns/rows for other properties (those belong to "
        f"  different proformas).\n"
    )


def format_source_directives(directives: list[dict], scope: str = "both") -> str:
    """Format user source directives into a prompt section.

    Parameters
    ----------
    directives:
        List of dicts with keys: source_id, source_type, directive, scope.
    scope:
        Filter to only include directives matching this scope ("mapping",
        "slide_generation", or "both").  Directives with scope "both"
        are always included.
    """
    relevant = [
        d for d in directives
        if d.get("directive", "").strip()
        and d.get("scope", "both") in (scope, "both")
    ]
    if not relevant:
        return ""
    lines = ["\n## Source Directives — FOLLOW THESE USER INSTRUCTIONS"]
    lines.append(
        "The user has provided specific instructions for how to use certain sources. "
        "You MUST follow these directives precisely. If a directive says to ignore a "
        "source or limit its use to specific sections, obey that constraint."
    )
    for d in relevant:
        src_label = d.get("source_id", d.get("source_type", "unknown"))
        lines.append(f"- **{src_label}**: {d['directive']}")
    return "\n".join(lines) + "\n"


def get_metric_mappings(
    client: anthropic.Anthropic,
    proforma_data: str,
    memo_content: str,
    cfg: dict,
    property_name: str = "",
    telemetry: dict | None = None,
    source_directives: list[dict] | None = None,
) -> dict:
    """
    Send proforma data + memo content to Claude and receive structured
    JSON describing every metric update.

    Uses prompt caching: the instructions + proforma data go in a cached
    system message (identical across batches), while only the memo chunk
    varies in the user message.
    """
    model = cfg["claude"]["model"]
    max_tokens = cfg["claude"]["max_tokens"]
    temperature = cfg["claude"]["temperature"]
    use_thinking = "opus" in model.lower()

    pn_section = _property_name_section(property_name, "mapping")
    directives_section = format_source_directives(source_directives or [], scope="mapping")

    # Split prompt into cached system prefix (instructions + proforma)
    # and varying user message (memo chunk only).
    system_text = MAPPING_PROMPT.format(
        proforma_data=proforma_data,
        memo_content="(see user message below)",
        property_name_section=pn_section,
        source_directives_section=directives_section,
    )
    user_text = f"## Memo Content (from PowerPoint)\n{memo_content}"

    total_chars = len(system_text) + len(user_text)
    log.info("Calling Claude API (model=%s, thinking=%s, prompt=%d chars, "
             "cached_prefix=%d chars)...",
             model, use_thinking, total_chars, len(system_text))

    api_kwargs = dict(
        model=model,
        max_tokens=max_tokens,
        system=[{
            "type": "text",
            "text": system_text,
            "cache_control": {"type": "ephemeral"},
        }],
        messages=[{"role": "user", "content": user_text}],
    )
    if use_thinking:
        api_kwargs["thinking"] = {"type": "adaptive"}
    else:
        api_kwargs["temperature"] = temperature

    if telemetry is not None:
        telemetry["mapping_api_calls"] = telemetry.get("mapping_api_calls", 0) + 1
    message = _create_message(client, **api_kwargs)

    # Log cache performance
    usage = message.usage
    cache_read = getattr(usage, "cache_read_input_tokens", 0) or 0
    cache_write = getattr(usage, "cache_creation_input_tokens", 0) or 0
    log.info("Token usage: input=%d, cache_read=%d, cache_write=%d, output=%d",
             usage.input_tokens, cache_read, cache_write, usage.output_tokens)
    if telemetry is not None:
        telemetry["cache_read_tokens"] = telemetry.get("cache_read_tokens", 0) + cache_read
        telemetry["cache_write_tokens"] = telemetry.get("cache_write_tokens", 0) + cache_write

    # Extract text from response (skip thinking blocks)
    raw = ""
    for block in message.content:
        if block.type == "text":
            raw = block.text
            break
    log.info("Claude response received (%d chars, %s stop_reason)",
             len(raw), message.stop_reason)

    if message.stop_reason == "max_tokens":
        log.warning(
            "Claude's response was cut off (hit max_tokens=%d). "
            "Attempting to salvage partial entries...",
            max_tokens,
        )
        salvaged = _salvage_truncated_json(raw)
        if salvaged is not None:
            n = len(salvaged["table_updates"]) + len(salvaged["text_updates"]) + len(salvaged["row_inserts"])
            log.info("Salvaged %d updates from truncated response", n)
            salvaged["_truncated"] = True
            return salvaged
        else:
            log.warning("Could not salvage truncated response - "
                        "caller will retry with smaller chunks")
            return {"table_updates": [], "text_updates": [], "row_inserts": [], "_truncated": True}

    # Parse JSON using consolidated helper
    empty_mappings = {"table_updates": [], "text_updates": [], "row_inserts": []}
    mappings = _parse_json_response(raw)
    if mappings is None:
        # Retry once - Claude sometimes returns analysis text instead of JSON
        log.warning("Claude returned non-JSON response - retrying with stricter prompt...")
        retry_suffix = (
            "\n\nIMPORTANT: You MUST respond with ONLY the JSON object. "
            "Do NOT include any analysis, explanation, or reasoning. "
            "Start your response with { and end with }."
        )
        api_kwargs["messages"] = [{"role": "user", "content": user_text + retry_suffix}]
        if telemetry is not None:
            telemetry["mapping_api_calls"] = telemetry.get("mapping_api_calls", 0) + 1
        retry_msg = _create_message(client, **api_kwargs)
        retry_raw = ""
        for block in retry_msg.content:
            if block.type == "text":
                retry_raw = block.text
                break
        log.info("Retry response received (%d chars, %s stop_reason)",
                 len(retry_raw), retry_msg.stop_reason)
        mappings = _parse_json_response(retry_raw)
        if mappings is None:
            log.info("Retry also returned unparseable response - no updates for this batch")
            return empty_mappings

    # Ensure expected keys exist
    mappings.setdefault("table_updates", [])
    mappings.setdefault("text_updates", [])
    mappings.setdefault("row_inserts", [])
    mappings.setdefault("narrative_updates", [])
    mappings.setdefault("table_structure_updates", [])

    n_table = len(mappings["table_updates"])
    n_text = len(mappings["text_updates"])
    n_row_ins = len(mappings["row_inserts"])
    n_narrative = len(mappings["narrative_updates"])
    n_structure = len(mappings["table_structure_updates"])
    log.info("Parsed mappings: %d table, %d text, %d row inserts, %d narrative, %d structure",
             n_table, n_text, n_row_ins, n_narrative, n_structure)
    return mappings


# ============================================================================
# 7. CLAUDE API - VALIDATION PASS
# ============================================================================
VALIDATION_PROMPT = _load_prompt_template("validation_v1.txt")


def _call_validation_api(
    client: anthropic.Anthropic,
    indexed_mappings: dict,
    proforma_data: str,
    memo_content: str,
    cfg: dict,
    property_name: str = "",
    telemetry: dict | None = None,
    source_directives: list[dict] | None = None,
) -> dict:
    """
    Single validation API call. Returns the parsed JSON result from Claude.
    Extracted as a helper so validate_mappings can batch multiple calls.
    Uses validation_model (defaults to same as mapping model if not set).

    Uses prompt caching: memo content + proforma + instructions go in a
    cached system message, while only the mappings JSON varies per batch.
    """
    model = cfg["claude"].get("validation_model", cfg["claude"]["model"])
    max_tokens = cfg["claude"]["max_tokens"]
    temperature = cfg["claude"]["temperature"]
    use_thinking = "opus" in model.lower()

    pn_section = _property_name_section(property_name, "validation")
    directives_section = format_source_directives(source_directives or [], scope="mapping")

    # Split: static context (instructions + memo + proforma) cached,
    # varying part (mappings JSON) in user message.
    system_text = VALIDATION_PROMPT.format(
        mappings_json="(see user message below)",
        memo_content=memo_content,
        proforma_data=proforma_data,
        property_name_section=pn_section,
        source_directives_section=directives_section,
    )
    user_text = (
        "## Proposed Changes (JSON, each entry has an \"idx\" field)\n"
        + json.dumps(indexed_mappings, indent=2)
    )

    total_chars = len(system_text) + len(user_text)
    log.info("Calling Claude API for validation (model=%s, thinking=%s, "
             "prompt=%d chars, cached_prefix=%d chars)...",
             model, use_thinking, total_chars, len(system_text))

    api_kwargs = dict(
        model=model,
        max_tokens=max_tokens,
        system=[{
            "type": "text",
            "text": system_text,
            "cache_control": {"type": "ephemeral"},
        }],
        messages=[{"role": "user", "content": user_text}],
    )
    if use_thinking:
        api_kwargs["thinking"] = {"type": "adaptive"}
    else:
        api_kwargs["temperature"] = temperature

    if telemetry is not None:
        telemetry["validation_api_calls"] = telemetry.get("validation_api_calls", 0) + 1
    message = _create_message(client, **api_kwargs)

    # Log cache performance
    usage = message.usage
    cache_read = getattr(usage, "cache_read_input_tokens", 0) or 0
    cache_write = getattr(usage, "cache_creation_input_tokens", 0) or 0
    log.info("Token usage: input=%d, cache_read=%d, cache_write=%d, output=%d",
             usage.input_tokens, cache_read, cache_write, usage.output_tokens)
    if telemetry is not None:
        telemetry["cache_read_tokens"] = telemetry.get("cache_read_tokens", 0) + cache_read
        telemetry["cache_write_tokens"] = telemetry.get("cache_write_tokens", 0) + cache_write

    # Extract text from response (skip thinking blocks)
    raw = ""
    for block in message.content:
        if block.type == "text":
            raw = block.text
            break
    log.info("Validation response received (%d chars, %s stop_reason)",
             len(raw), message.stop_reason)

    if message.stop_reason == "max_tokens":
        log.warning(
            "Claude's validation response was cut off (hit max_tokens=%d). "
            "Marking as truncated for re-batching.",
            max_tokens,
        )
        return {"rejected": [], "corrections": [], "missed": [], "_truncated": True}

    # Parse JSON using consolidated helper
    empty_result = {"rejected": [], "corrections": [], "missed": []}
    result = _parse_json_response(raw)
    if result is None:
        log.info("Validation returned empty/unparseable response - all entries pass")
        return empty_result

    return result


def validate_mappings(
    client: anthropic.Anthropic,
    mappings: dict,
    proforma_data: str,
    memo_content: str,
    cfg: dict,
    property_name: str = "",
    telemetry: dict | None = None,
    source_directives: list[dict] | None = None,
) -> dict:
    """
    Second Claude API call - validates the proposed mappings by cross-checking
    old values against the memo and new values against the proforma.

    The prompt asks Claude to return ONLY rejections, corrections, and missed
    entries (not all valid entries), keeping the response compact and well
    within token limits. Valid entries are inferred by exclusion.

    For large decks, batches the validation by page groups (same threshold
    as the mapping step) so the prompt stays within model limits.

    Returns a validated/cleaned version of the mappings with rejected entries
    removed and any missed metrics flagged.
    """
    BATCH_THRESHOLD = 80_000  # chars; same as mapping step
    RATE_LIMIT_INTERVAL = 5  # seconds between API calls

    # Add idx to each entry so Claude can reference them by index
    table_updates = mappings.get("table_updates", [])
    text_updates = mappings.get("text_updates", [])
    row_inserts = mappings.get("row_inserts", [])
    indexed_mappings = {
        "table_updates": [
            {**entry, "idx": i} for i, entry in enumerate(table_updates)
        ],
        "text_updates": [
            {**entry, "idx": i} for i, entry in enumerate(text_updates)
        ],
        "row_inserts": [
            {**entry, "idx": i} for i, entry in enumerate(row_inserts)
        ],
    }

    prompt_size = (len(proforma_data) + len(memo_content)
                   + len(json.dumps(indexed_mappings)))

    if prompt_size > BATCH_THRESHOLD:
        # Batch by page groups - split memo into chunks and only send
        # the mappings relevant to each chunk's pages.
        log.info("Large validation prompt (%d chars) - batching by page groups",
                 prompt_size)
        memo_chunks = chunk_memo_by_pages(memo_content, pages_per_chunk=5)

        # Determine which pages each chunk covers
        merged_result = {"rejected": [], "corrections": [], "missed": []}
        last_api_call = 0
        for ci, chunk in enumerate(memo_chunks, 1):
            # Extract page numbers from this chunk
            chunk_pages = set(
                int(m) for m in re.findall(r"PAGE (\d+)", chunk)
            )

            # Filter mappings to only entries for pages in this chunk
            chunk_indexed = {
                "table_updates": [
                    e for e in indexed_mappings["table_updates"]
                    if e.get("page") in chunk_pages
                ],
                "text_updates": [
                    e for e in indexed_mappings["text_updates"]
                    if e.get("page") in chunk_pages
                ],
                "row_inserts": [
                    e for e in indexed_mappings["row_inserts"]
                    if e.get("page") in chunk_pages
                ],
            }

            n_entries = (len(chunk_indexed["table_updates"])
                         + len(chunk_indexed["text_updates"])
                         + len(chunk_indexed["row_inserts"]))
            if n_entries == 0:
                log.info("Validation batch %d/%d: no mappings for pages %s - skipping",
                         ci, len(memo_chunks), sorted(chunk_pages))
                continue

            if ci > 1 and last_api_call > 0:
                elapsed = time.time() - last_api_call
                wait = RATE_LIMIT_INTERVAL - elapsed
                if wait > 0:
                    log.info("Rate limit: waiting %.0f seconds...", wait)
                    time.sleep(wait)

            log.info("Validation batch %d/%d (%d entries, pages %s)...",
                     ci, len(memo_chunks), n_entries, sorted(chunk_pages))
            last_api_call = time.time()
            batch_result = _call_validation_api(
                client, chunk_indexed, proforma_data, chunk, cfg,
                property_name=property_name,
                telemetry=telemetry,
                source_directives=source_directives,
            )

            if batch_result.pop("_truncated", False):
                # Re-batch with single-page sub-chunks
                log.warning(
                    "Validation batch %d/%d truncated — retrying with "
                    "single-page sub-chunks for pages %s",
                    ci, len(memo_chunks), sorted(chunk_pages),
                )
                sub_chunks = chunk_memo_by_pages(chunk, pages_per_chunk=1)
                for si, sub_chunk in enumerate(sub_chunks, 1):
                    sub_pages = set(
                        int(m) for m in re.findall(r"PAGE (\d+)", sub_chunk)
                    )
                    sub_indexed = {
                        "table_updates": [
                            e for e in chunk_indexed["table_updates"]
                            if e.get("page") in sub_pages
                        ],
                        "text_updates": [
                            e for e in chunk_indexed["text_updates"]
                            if e.get("page") in sub_pages
                        ],
                        "row_inserts": [
                            e for e in chunk_indexed["row_inserts"]
                            if e.get("page") in sub_pages
                        ],
                    }
                    n_sub = (len(sub_indexed["table_updates"])
                             + len(sub_indexed["text_updates"])
                             + len(sub_indexed["row_inserts"]))
                    if n_sub == 0:
                        continue
                    if last_api_call > 0:
                        elapsed = time.time() - last_api_call
                        wait = RATE_LIMIT_INTERVAL - elapsed
                        if wait > 0:
                            time.sleep(wait)
                    last_api_call = time.time()
                    sub_result = _call_validation_api(
                        client, sub_indexed, proforma_data, sub_chunk, cfg,
                        property_name=property_name,
                        telemetry=telemetry,
                        source_directives=source_directives,
                    )
                    if sub_result.pop("_truncated", False):
                        log.warning(
                            "  Validation sub-chunk %d (pages %s) still "
                            "truncated — these entries pass through "
                            "UNVALIDATED",
                            si, sorted(sub_pages),
                        )
                        unvalidated_pages = merged_result.setdefault(
                            "_unvalidated_pages", [])
                        unvalidated_pages.extend(sorted(sub_pages))
                    else:
                        merged_result["rejected"].extend(
                            sub_result.get("rejected", []))
                        merged_result["corrections"].extend(
                            sub_result.get("corrections", []))
                        merged_result["missed"].extend(
                            sub_result.get("missed", []))
            else:
                merged_result["rejected"].extend(batch_result.get("rejected", []))
                merged_result["corrections"].extend(batch_result.get("corrections", []))
                merged_result["missed"].extend(batch_result.get("missed", []))

        result = merged_result
    else:
        result = _call_validation_api(
            client, indexed_mappings, proforma_data, memo_content, cfg,
            property_name=property_name,
            telemetry=telemetry,
            source_directives=source_directives,
        )

    # Reconstruct validated mappings: start with originals, remove rejections,
    # apply corrections
    rejected_table_idxs = set()
    rejected_text_idxs = set()
    rejected_row_insert_idxs = set()
    correction_table = {}
    correction_text = {}
    correction_row_insert = {}

    for rej in result.get("rejected", []):
        idx = rej.get("idx")
        if idx is not None:
            if rej.get("type") == "text":
                rejected_text_idxs.add(idx)
            elif rej.get("type") == "row_insert":
                rejected_row_insert_idxs.add(idx)
            else:
                rejected_table_idxs.add(idx)

    for cor in result.get("corrections", []):
        idx = cor.get("idx")
        if idx is not None:
            if cor.get("type") == "text":
                correction_text[idx] = cor["corrected_entry"]
            elif cor.get("type") == "row_insert":
                correction_row_insert[idx] = cor["corrected_entry"]
            else:
                correction_table[idx] = cor["corrected_entry"]

    valid_table = []
    for i, entry in enumerate(table_updates):
        if i in rejected_table_idxs:
            continue
        if i in correction_table:
            valid_table.append(correction_table[i])
        else:
            valid_table.append(entry)

    valid_text = []
    for i, entry in enumerate(text_updates):
        if i in rejected_text_idxs:
            continue
        if i in correction_text:
            valid_text.append(correction_text[i])
        else:
            valid_text.append(entry)

    valid_row_inserts = []
    for i, entry in enumerate(row_inserts):
        if i in rejected_row_insert_idxs:
            continue
        if i in correction_row_insert:
            valid_row_inserts.append(correction_row_insert[i])
        else:
            valid_row_inserts.append(entry)

    n_rejected = len(result.get("rejected", []))
    n_corrections = len(result.get("corrections", []))
    n_missed = len(result.get("missed", []))
    log.info("Validation: %d passed, %d rejected, %d corrected, %d missed",
             len(valid_table) + len(valid_text) + len(valid_row_inserts),
             n_rejected, n_corrections, n_missed)

    if n_rejected > 0:
        for rej in result["rejected"]:
            log.warning("  REJECTED idx=%s: %s", rej.get("idx", "?"),
                        rej.get("reason", "unknown"))
    if n_corrections > 0:
        for cor in result["corrections"]:
            log.warning("  CORRECTED idx=%s: %s", cor.get("idx", "?"),
                        cor.get("reason", "unknown"))
    if n_missed > 0:
        for miss in result["missed"]:
            log.warning("  MISSED: page %s - %s", miss.get("page", "?"),
                        miss.get("description", ""))

    # Build rejected list for change log (include full original entries)
    rejected_entries = []
    for rej in result.get("rejected", []):
        idx = rej.get("idx")
        entry_type = rej.get("type", "table")
        original = {}
        if entry_type == "text" and idx is not None and idx < len(text_updates):
            original = text_updates[idx]
        elif entry_type == "row_insert" and idx is not None and idx < len(row_inserts):
            original = row_inserts[idx]
        elif idx is not None and idx < len(table_updates):
            original = table_updates[idx]
        rejected_entries.append({
            "original": original,
            "reason": rej.get("reason", "unknown"),
        })

    validated = {
        "table_updates": valid_table,
        "text_updates": valid_text,
        "row_inserts": valid_row_inserts,
        "rejected": rejected_entries,
        "missed": result.get("missed", []),
    }
    # Propagate unvalidated pages so callers can warn users
    unvalidated = result.get("_unvalidated_pages")
    if unvalidated:
        validated["_unvalidated_pages"] = sorted(set(unvalidated))
    return validated


# ============================================================================
# 8a. FORMAT VALIDATION (detect style mismatches)
# ============================================================================

def _detect_number_format(value: str) -> dict:
    """Detect the formatting style of a numeric string.

    Returns a dict with keys: has_commas, has_dollar, decimal_places, has_percent.
    """
    # Strip non-numeric suffixes (like %) before counting decimal places
    numeric_part = re.sub(r"[%$,\s]", "", value)
    if "." in numeric_part:
        decimal_places = len(numeric_part.split(".")[-1])
    else:
        decimal_places = 0
    return {
        "has_dollar": "$" in value,
        "has_commas": bool(re.search(r"\d{1,3}(,\d{3})+", value)),
        "has_percent": "%" in value,
        "decimal_places": decimal_places,
    }


def _format_matches(old_val: str, new_val: str) -> tuple[bool, str]:
    """Check if new_val preserves old_val's formatting style.

    Returns (is_ok, reason) — reason is empty if ok.
    """
    old_fmt = _detect_number_format(old_val)
    new_fmt = _detect_number_format(new_val)

    # Dollar sign consistency
    if old_fmt["has_dollar"] and not new_fmt["has_dollar"]:
        return False, f"missing dollar sign: '{old_val}' -> '{new_val}'"
    if not old_fmt["has_dollar"] and new_fmt["has_dollar"]:
        return False, f"unexpected dollar sign: '{old_val}' -> '{new_val}'"

    # Comma consistency (only check if value is large enough to need commas)
    # Strip dollar/percent for numeric comparison
    new_digits = re.sub(r"[^\d.]", "", new_val)
    if old_fmt["has_commas"] and not new_fmt["has_commas"]:
        try:
            if float(new_digits) >= 1000:
                return False, f"missing commas: '{old_val}' -> '{new_val}'"
        except ValueError:
            pass

    # Percentage consistency
    if old_fmt["has_percent"] and not new_fmt["has_percent"]:
        return False, f"missing percent sign: '{old_val}' -> '{new_val}'"

    # Decimal precision
    if old_fmt["has_percent"] and new_fmt["has_percent"]:
        if old_fmt["decimal_places"] != new_fmt["decimal_places"]:
            return False, (
                f"decimal precision mismatch: '{old_val}' ({old_fmt['decimal_places']}dp) "
                f"-> '{new_val}' ({new_fmt['decimal_places']}dp)"
            )

    return True, ""


def _auto_fix_format(old_val: str, new_val: str) -> str:
    """Attempt to auto-correct formatting in new_val to match old_val's style."""
    old_fmt = _detect_number_format(old_val)

    # Extract the raw number from new_val
    raw = re.sub(r"[^\d.\-]", "", new_val)
    if not raw:
        return new_val

    try:
        num = float(raw)
    except ValueError:
        return new_val

    # Build formatted string
    result = ""

    if old_fmt["has_dollar"]:
        result += "$"

    if old_fmt["has_commas"]:
        if old_fmt["decimal_places"] > 0:
            result += f"{num:,.{old_fmt['decimal_places']}f}"
        else:
            result += f"{int(num):,}"
    else:
        if old_fmt["decimal_places"] > 0:
            result += f"{num:.{old_fmt['decimal_places']}f}"
        else:
            result += str(int(num))

    if old_fmt["has_percent"]:
        result += "%"

    return result


def validate_mapping_formats(mappings: dict) -> dict:
    """Check and auto-fix formatting consistency between old and new values.

    Returns the mappings dict with formatting issues fixed in-place and
    unfixable mismatches added to the rejected list.
    """
    fixed_count = 0
    format_warnings = []

    for upd in mappings.get("table_updates", []):
        old_val = upd.get("old_value", "")
        new_val = upd.get("new_value", "")
        if not old_val or not new_val:
            continue
        ok, reason = _format_matches(old_val, new_val)
        if not ok:
            corrected = _auto_fix_format(old_val, new_val)
            if corrected != new_val:
                upd["new_value"] = corrected
                fixed_count += 1
            else:
                format_warnings.append(
                    f"page {upd.get('page')}: {reason}"
                )

    for upd in mappings.get("text_updates", []):
        old_txt = upd.get("old_text", "")
        new_txt = upd.get("new_text", "")
        if not old_txt or not new_txt:
            continue
        # Extract embedded numbers for format checking
        old_nums = re.findall(r"\$[\d,]+\.?\d*%?|\d[\d,]*\.?\d*%", old_txt)
        new_nums = re.findall(r"\$[\d,]+\.?\d*%?|\d[\d,]*\.?\d*%", new_txt)
        for old_n, new_n in zip(old_nums, new_nums):
            ok, reason = _format_matches(old_n, new_n)
            if not ok:
                corrected = _auto_fix_format(old_n, new_n)
                if corrected != new_n:
                    upd["new_text"] = upd["new_text"].replace(new_n, corrected, 1)
                    fixed_count += 1

    if fixed_count > 0:
        log.info("Format validation auto-corrected %d values", fixed_count)
    if format_warnings:
        for w in format_warnings[:5]:
            log.warning("Format mismatch (unfixable): %s", w)

    return mappings


# ============================================================================
# 8b. PRE-VALIDATION (local Python check)
# ============================================================================
def pre_validate_mappings(mappings: dict, memo_content: str) -> dict:
    """
    Quick local check: verify each old_value / old_text actually exists in the
    memo content. Reject entries that don't match. This catches the most common
    Claude errors instantly without an API call.

    Returns a new mappings dict with non-matching entries moved to 'rejected'.
    """
    def _split_memo_by_page(content: str) -> dict:
        """
        Build a {page_number: page_text_block} map from extracted memo content.
        """
        blocks = {}
        page_headers = list(re.finditer(
            r"={60,}\nPAGE\s+(\d+)[^\n]*\n={60,}",
            content,
        ))
        if not page_headers:
            return blocks

        for i, m in enumerate(page_headers):
            page_num = int(m.group(1))
            start = m.start()
            end = page_headers[i + 1].start() if i + 1 < len(page_headers) else len(content)
            blocks[page_num] = content[start:end]
        return blocks

    page_blocks = _split_memo_by_page(memo_content)
    valid_table = []
    valid_text = []
    rejected = list(mappings.get("rejected", []))

    for upd in mappings.get("table_updates", []):
        page = upd.get("page")
        haystack = page_blocks.get(page, memo_content)
        old_value = upd.get("old_value", "")
        if old_value and old_value in haystack:
            valid_table.append(upd)
        else:
            rejected.append({
                "original": upd,
                "reason": (
                    f"old_value not found on page {page}: '{old_value}'"
                    if page in page_blocks
                    else f"old_value not found in memo: '{old_value}'"
                ),
            })

    for upd in mappings.get("text_updates", []):
        page = upd.get("page")
        haystack = page_blocks.get(page, memo_content)
        old_text = upd.get("old_text", "")
        if old_text and old_text in haystack:
            valid_text.append(upd)
        else:
            rejected.append({
                "original": upd,
                "reason": (
                    f"old_text not found on page {page}: '{old_text}'"
                    if page in page_blocks
                    else f"old_text not found in memo: '{old_text}'"
                ),
            })

    valid_narrative = []
    for upd in mappings.get("narrative_updates", []):
        page = upd.get("page")
        haystack = page_blocks.get(page, memo_content)
        old_narrative = upd.get("old_narrative", "")
        if old_narrative and old_narrative in haystack:
            valid_narrative.append(upd)
        else:
            rejected.append({
                "original": upd,
                "reason": (
                    f"old_narrative not found on page {page}: '{old_narrative[:60]}...'"
                    if page in page_blocks
                    else f"old_narrative not found in memo: '{old_narrative[:60]}...'"
                ),
            })

    n_rejected_new = (len(rejected) - len(mappings.get("rejected", [])))
    if n_rejected_new > 0:
        log.warning("Pre-validation rejected %d entries (old value not in memo)",
                    n_rejected_new)

    return {
        "table_updates": valid_table,
        "text_updates": valid_text,
        "row_inserts": mappings.get("row_inserts", []),
        "narrative_updates": valid_narrative,
        "table_structure_updates": mappings.get("table_structure_updates", []),
        "rejected": rejected,
        "missed": mappings.get("missed", []),
    }


# ============================================================================
# 9. APPLY TEXT / TABLE UPDATES  (python-pptx)
# ============================================================================
def _replace_in_para(para, old_text: str, new_text: str) -> bool:
    """
    Replace old_text with new_text in a paragraph, handling the common
    case where a value is split across multiple XML runs.

    Strategy:
    1. Try a direct single-run replacement first (fastest, preserves all formatting).
    2. Fall back to a full-paragraph merge: concatenate all run texts, do the
       replacement, write the result into the first run, and clear the rest.
       This loses per-run character formatting within the cell but preserves
       paragraph-level formatting (alignment, spacing), which is acceptable
       for financial table values that are typically uniform within a cell.
    """
    # Pass 1: single-run replacement
    for run in para.runs:
        if old_text in run.text:
            run.text = run.text.replace(old_text, new_text)
            return True

    # Pass 2: cross-run replacement (format-preserving)
    full_text = "".join(r.text for r in para.runs)
    if old_text not in full_text:
        return False

    if para.runs:
        # Calculate character positions for each run
        run_starts = []
        pos = 0
        for run in para.runs:
            run_starts.append(pos)
            pos += len(run.text)

        # Find where old_text starts and ends
        match_start = full_text.index(old_text)
        match_end = match_start + len(old_text)

        # Find first and last runs that overlap with the match
        first_run_idx = None
        last_run_idx = None
        for i, run in enumerate(para.runs):
            run_end = run_starts[i] + len(run.text)
            if first_run_idx is None and run_end > match_start:
                first_run_idx = i
            if run_starts[i] < match_end:
                last_run_idx = i

        if first_run_idx is not None and last_run_idx is not None:
            # Preserve text before match in first overlapping run
            pre = full_text[run_starts[first_run_idx]:match_start]
            # Preserve text after match in last overlapping run
            last_run_end = run_starts[last_run_idx] + len(para.runs[last_run_idx].text)
            post = full_text[match_end:last_run_end]

            # Save formatting from the first overlapping run so we can
            # re-apply it after writing new text.  The run's XML element
            # keeps its rPr (run properties) automatically — we just need
            # to make sure the cleared runs don't leave orphan formatting.
            first_run = para.runs[first_run_idx]

            # Write replacement into first overlapping run
            first_run.text = pre + new_text + post

            # For cleared runs: preserve the run element (keeps XML valid)
            # but empty the text.  If the last run has text AFTER the
            # match region, move it there to preserve its formatting.
            for i in range(first_run_idx + 1, last_run_idx + 1):
                para.runs[i].text = ""

            return True

    return False


def _replace_in_cell(cell, old_text: str, new_text: str) -> bool:
    """Replace old_text with new_text inside a table cell, preserving runs.
    Falls back to Unicode-normalized matching if exact match fails."""
    # Pass 1: exact match
    for para in cell.text_frame.paragraphs:
        if _replace_in_para(para, old_text, new_text):
            return True
    # Pass 2: Unicode-normalized match (non-breaking spaces, smart quotes, etc.)
    norm_old = _normalize_unicode(old_text)
    for para in cell.text_frame.paragraphs:
        para_text = para.text
        norm_para = _normalize_unicode(para_text)
        if norm_old not in norm_para:
            continue
        idx = norm_para.find(norm_old)
        actual_old = para_text[idx:idx + len(norm_old)]
        if _replace_in_para(para, actual_old, new_text):
            return True
    return False


def _normalize_unicode(text: str) -> str:
    """
    Normalize Unicode characters that PowerPoint often substitutes:
    non-breaking spaces -> spaces, smart quotes -> ASCII quotes,
    en/em dashes -> hyphens, etc.
    """
    return (text
            .replace("\u00a0", " ")    # non-breaking space
            .replace("\u2013", "-")    # en dash
            .replace("\u2014", "-")    # em dash
            .replace("\u2018", "'")    # left single quote
            .replace("\u2019", "'")    # right single quote
            .replace("\u201c", '"')    # left double quote
            .replace("\u201d", '"')    # right double quote
            .replace("\u2502", "|")    # box drawing vertical
            .replace("\uff5c", "|"))   # fullwidth vertical line


def _replace_in_shape(shape, old_text: str, new_text: str) -> bool:
    """Replace old_text in any text-frame shape, preserving formatting.
    Falls back to Unicode-normalized matching if exact match fails."""
    if not shape.has_text_frame:
        return False
    # Pass 1: exact match
    for para in shape.text_frame.paragraphs:
        if old_text not in para.text:
            continue
        if _replace_in_para(para, old_text, new_text):
            return True
    # Pass 2: normalized match - find the actual text in the shape that
    # corresponds to the normalized old_text, then replace that
    norm_old = _normalize_unicode(old_text)
    for para in shape.text_frame.paragraphs:
        para_text = para.text
        norm_para = _normalize_unicode(para_text)
        if norm_old not in norm_para:
            continue
        # Find the actual substring in the original para text
        idx = norm_para.find(norm_old)
        actual_old = para_text[idx:idx + len(norm_old)]
        if _replace_in_para(para, actual_old, new_text):
            return True
    return False


def _normalize_for_match(text: str) -> str:
    """Normalize text for robust table/row matching."""
    return re.sub(r"\s+", " ", (text or "").strip()).lower()


def _normalize_ws(text: str) -> str:
    """Normalize all whitespace (including newlines from multi-paragraph cells) to single spaces."""
    return re.sub(r"\s+", " ", (text or "").strip())


def _strip_to_core(text: str) -> str:
    """
    Strip a label down to its alphanumeric core for bedroom-type matching.
    '1BR/1BA' -> '1br1ba', '1 BR' -> '1br', '4BR/2BA' -> '4br2ba'
    """
    return re.sub(r"[^a-z0-9]", "", text.lower())


def _loose_match(expected: str, actual: str) -> bool:
    """
    True when expected roughly matches actual (exact or containment after
    whitespace/case normalization), or when both have the same full bedroom
    label (e.g. '1BR/1BA' matches '1 BR / 1 BA').

    The previous implementation matched on leading digit only ('1br' == '1br'),
    which caused '1BR/1BA' to match '1BR/2BA'.  Now requires the full
    alphanumeric core to match.
    """
    e = _normalize_for_match(expected)
    a = _normalize_for_match(actual)
    if not e:
        return True
    if not a:
        return False
    if e == a or e in a or a in e:
        return True
    # Bedroom-label fallback: require FULL alphanumeric core to match
    # e.g. '1br1ba' == '1br1ba', NOT '1br1ba' == '1br2ba'
    ec = _strip_to_core(expected)
    ac = _strip_to_core(actual)
    if ec and ac and ec == ac:
        return True
    return False


def _find_table_target(slide, table_name: str, row_label: str,
                       col_idx: int, old_value: str):
    """
    Find the best target cell for a table update.
    Preference order:
    1) Tables matching table_name + rows matching row_label + cell has old_value
    1b) Same as 1 but with Unicode-normalized old_value
    2) Any table + row_label match + old_value at col_idx (ignore table_name)
    3) Any table + old_value at col_idx in any row (ignore row_label) — WARNING logged

    Pass 3 (any cell, any column) has been **removed** because it silently
    matches the wrong table/row when common values (e.g. "$1,200") appear in
    multiple places.

    Returns:
      (shape, row_label_actual, cell, match_pass) on success, or
      (None, row_label_actual, cell_text, 0) for diagnostics.
    """
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return None, None, None, 0

    if table_name:
        preferred = [s for s in tables if _loose_match(table_name, s.name)]
        fallback = [s for s in tables if s not in preferred]
        ordered_groups = [preferred, fallback] if preferred else [tables]
    else:
        ordered_groups = [tables]

    diagnostic_row = None
    diagnostic_cell = None

    norm_old = _normalize_unicode(old_value)
    # Whitespace-normalized old_value for multi-paragraph cell matching
    ws_old = _normalize_ws(old_value)

    # Pass 1: row_label match + old_value at col_idx (exact, then Unicode/ws-normalized)
    for group in ordered_groups:
        for shape in group:
            for row in shape.table.rows:
                if col_idx >= len(row.cells):
                    continue
                row_head = row.cells[0].text.strip() if row.cells else ""
                if row_label and not _loose_match(row_label, row_head):
                    continue

                cell = row.cells[col_idx]
                cell_text = cell.text or ""
                ws_cell = _normalize_ws(cell_text)
                if old_value in cell_text or ws_old in ws_cell:
                    return shape, row_head, cell, 1

                # Unicode-normalized fallback within Pass 1
                if norm_old != old_value and norm_old in _normalize_unicode(cell_text):
                    return shape, row_head, cell, 1

                if diagnostic_row is None:
                    diagnostic_row = row_head
                    diagnostic_cell = cell_text

    # Pass 2: ignore row_label, find old_value at col_idx in any row
    for group in ordered_groups:
        for shape in group:
            for row in shape.table.rows:
                if col_idx >= len(row.cells):
                    continue
                cell = row.cells[col_idx]
                cell_text = cell.text or ""
                ws_cell = _normalize_ws(cell_text)
                matched = old_value in cell_text or ws_old in ws_cell
                if not matched and norm_old != old_value:
                    matched = norm_old in _normalize_unicode(cell_text)
                if matched:
                    row_head = row.cells[0].text.strip() if row.cells else ""
                    log.warning(
                        "Table match degraded to Pass 2 (row_label mismatch): "
                        "expected row '%s', matched row '%s' col %d for '%s'",
                        row_label, row_head, col_idx, old_value,
                    )
                    return shape, row_head, cell, 2

    # Pass 3: Search ALL columns for old_value (for side-by-side comp tables
    # where the subject property column is col 2, 3, or 4).
    # Only enabled when old_value is "specific enough" to avoid false matches:
    # at least 4 chars, contains a dollar sign/percent, looks like a unit type
    # pattern like "4BR / 4BA (53)", or is a pure integer >= 2 digits (e.g. "141").
    is_specific = (
        len(old_value) >= 4
        or "$" in old_value
        or "%" in old_value
        or re.search(r"\dBR\s*/\s*\d", old_value)
        or re.match(r"^\d{2,}$", old_value.replace(",", "").strip())
    )
    if is_specific:
        for group in ordered_groups:
            for shape in group:
                for row in shape.table.rows:
                    for ci, cell in enumerate(row.cells):
                        if ci == col_idx:
                            continue  # already checked in Pass 2
                        cell_text = cell.text or ""
                        ws_cell = _normalize_ws(cell_text)
                        matched = old_value in cell_text or ws_old in ws_cell
                        if not matched and norm_old != old_value:
                            matched = norm_old in _normalize_unicode(cell_text)
                        if matched:
                            row_head = row.cells[0].text.strip() if row.cells else ""
                            log.warning(
                                "Table match degraded to Pass 3 (cross-column): "
                                "expected col %d, matched col %d for '%s' in row '%s'",
                                col_idx, ci, old_value, row_head,
                            )
                            return shape, row_head, cell, 3

    return None, diagnostic_row, diagnostic_cell, 0


def _find_row_by_label(table, row_label: str) -> int | None:
    """
    Find a row index in a python-pptx table by matching column-0 text
    using _loose_match. Returns the 0-based row index, or None.
    """
    for idx, row in enumerate(table.rows):
        cell_text = row.cells[0].text.strip() if row.cells else ""
        if _loose_match(row_label, cell_text):
            return idx
    return None


def _add_table_row(table, reference_row_idx: int, cell_values: list):
    """
    Clone the XML of the row at *reference_row_idx*, clear all cell text,
    populate with *cell_values*, and insert the new row immediately after
    the reference row.

    Uses lxml operations (deepcopy / addnext) on the underlying DrawingML
    XML, which is the only reliable way to add rows to a python-pptx table.
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tbl_xml = table._tbl  # lxml element <a:tbl>

    # Get all <a:tr> elements
    rows = tbl_xml.findall(f"{{{ns}}}tr")
    if reference_row_idx < 0 or reference_row_idx >= len(rows):
        log.warning("_add_table_row: reference_row_idx %d out of range (table has %d rows)",
                    reference_row_idx, len(rows))
        return

    ref_tr = rows[reference_row_idx]
    new_tr = deepcopy(ref_tr)

    # Clear text in every cell of the cloned row and set new values
    new_cells = new_tr.findall(f"{{{ns}}}tc")
    for ci, tc in enumerate(new_cells):
        # Clear all text runs inside the cell
        for p in tc.findall(f".//{{{ns}}}p"):
            for r in p.findall(f"{{{ns}}}r"):
                t = r.find(f"{{{ns}}}t")
                if t is not None:
                    t.text = ""
        # Set new value in the first run of the first paragraph
        first_p = tc.find(f".//{{{ns}}}p")
        if first_p is not None:
            first_r = first_p.find(f"{{{ns}}}r")
            if first_r is not None:
                t = first_r.find(f"{{{ns}}}t")
                if t is not None:
                    t.text = cell_values[ci] if ci < len(cell_values) else ""

    ref_tr.addnext(new_tr)


def _add_table_column(table, after_col: int, header: str, values: list):
    """Add a column to a table after the given column index.

    Uses lxml operations on DrawingML XML to add <a:tc> elements to each row.
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tbl_xml = table._tbl
    rows = tbl_xml.findall(f"{{{ns}}}tr")

    for row_idx, tr in enumerate(rows):
        cells = tr.findall(f"{{{ns}}}tc")
        if not cells:
            continue

        # Clone the cell at after_col to preserve formatting
        ref_idx = min(after_col, len(cells) - 1)
        new_tc = deepcopy(cells[ref_idx])

        # Clear text and set new value
        for p in new_tc.findall(f".//{{{ns}}}p"):
            for r in p.findall(f"{{{ns}}}r"):
                t = r.find(f"{{{ns}}}t")
                if t is not None:
                    t.text = ""

        # Set the value
        first_p = new_tc.find(f".//{{{ns}}}p")
        if first_p is not None:
            first_r = first_p.find(f"{{{ns}}}r")
            if first_r is not None:
                t = first_r.find(f"{{{ns}}}t")
                if t is not None:
                    if row_idx == 0:
                        t.text = header
                    else:
                        val_idx = row_idx - 1
                        t.text = values[val_idx] if val_idx < len(values) else ""

        # Insert after the reference cell
        if ref_idx < len(cells) - 1:
            cells[ref_idx].addnext(new_tc)
        else:
            tr.append(new_tc)

    # Update gridCol count
    tbl_grid = tbl_xml.find(f"{{{ns}}}tblGrid")
    if tbl_grid is not None:
        grid_cols = tbl_grid.findall(f"{{{ns}}}gridCol")
        if grid_cols:
            new_gc = deepcopy(grid_cols[-1])
            tbl_grid.append(new_gc)


def _remove_table_row(table, row_idx: int):
    """Remove a row from a table by index."""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tbl_xml = table._tbl
    rows = tbl_xml.findall(f"{{{ns}}}tr")
    if 0 <= row_idx < len(rows):
        tbl_xml.remove(rows[row_idx])


def _reorder_table_rows(table, new_order: list[str]):
    """Reorder table rows to match the given label order.

    new_order is a list of row labels (column 0 text) in desired order.
    The header row (index 0) is always kept first.
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tbl_xml = table._tbl
    rows = tbl_xml.findall(f"{{{ns}}}tr")
    if len(rows) < 2:
        return

    header_tr = rows[0]
    data_rows = rows[1:]

    # Build label -> row element map
    label_map: dict[str, Any] = {}
    for tr in data_rows:
        cells = tr.findall(f"{{{ns}}}tc")
        if cells:
            first_p = cells[0].find(f".//{{{ns}}}p")
            if first_p is not None:
                first_r = first_p.find(f"{{{ns}}}r")
                if first_r is not None:
                    t = first_r.find(f"{{{ns}}}t")
                    if t is not None and t.text:
                        label_map[t.text.strip()] = tr

    # Remove all data rows
    for tr in data_rows:
        tbl_xml.remove(tr)

    # Re-add in new order
    for label in new_order:
        tr = label_map.pop(label, None)
        if tr is not None:
            tbl_xml.append(tr)

    # Append any remaining rows not in new_order
    for tr in label_map.values():
        tbl_xml.append(tr)


def global_property_rename(memo_path: str, old_name: str, new_name: str) -> int:
    """
    Replace ALL occurrences of *old_name* with *new_name* across every slide
    in the PPTX (text frames **and** table cells).  This is a mechanical
    find-replace that runs **before** the AI mapping/validation passes so
    that Claude sees the already-corrected property name.

    Returns the number of replacements made.
    """
    prs = _load_presentation(memo_path)
    count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # --- Text frames (titles, narrative, text boxes) ---
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if _replace_in_para(para, old_name, new_name):
                        count += 1
                    else:
                        # Try Unicode-normalized match
                        norm_old = _normalize_unicode(old_name)
                        norm_para = _normalize_unicode(para.text)
                        if norm_old in norm_para:
                            idx = norm_para.find(norm_old)
                            actual_old = para.text[idx:idx + len(norm_old)]
                            if _replace_in_para(para, actual_old, new_name):
                                count += 1

            # --- Table cells ---
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            if _replace_in_para(para, old_name, new_name):
                                count += 1
                            else:
                                norm_old = _normalize_unicode(old_name)
                                norm_para = _normalize_unicode(para.text)
                                if norm_old in norm_para:
                                    idx = norm_para.find(norm_old)
                                    actual_old = para.text[idx:idx + len(norm_old)]
                                    if _replace_in_para(para, actual_old, new_name):
                                        count += 1

    if count > 0:
        prs.save(memo_path)
        log.info("Global property rename: '%s' -> '%s' (%d replacements)",
                 old_name, new_name, count)
    else:
        log.warning("Global property rename: '%s' not found in memo", old_name)

    return count


def apply_updates(memo_path: str, mappings: dict, dry_run: bool = False) -> list:
    """
    Open the memo, apply every table_update and text_update from the
    Claude-validated mappings, and save. Returns a list of change records.
    """
    prs = _load_presentation(memo_path)
    changes = []

    # --- Table updates ---
    for upd in mappings.get("table_updates", []):
        # Skip misclassified row_inserts (have 'cells' but no 'old_value')
        if "old_value" not in upd:
            if "cells" in upd:
                log.warning("Skipping misclassified row_insert in table_updates: %s", upd)
            else:
                log.warning("Skipping table_update missing 'old_value': %s", upd)
            continue
        page = upd["page"]
        tbl_name = upd.get("table_name", "")
        col_idx = upd.get("column_index", 1)
        old_val = upd["old_value"]
        new_val = upd["new_value"]
        source = upd.get("source", "")
        row_label = upd.get("row_label", "")

        try:
            slide = prs.slides[page - 1]
        except IndexError:
            log.warning("Table update SKIPPED: page %d does not exist", page)
            continue
        shape, matched_row_label, cell_or_text, match_pass = _find_table_target(
            slide, tbl_name, row_label, col_idx, old_val
        )
        if shape is not None:
            if not dry_run:
                _replace_in_cell(cell_or_text, old_val, new_val)
            location_table = shape.name or tbl_name or "<unnamed table>"
            location_row = matched_row_label or row_label or "<unknown row>"
            change_record = {
                "page": page, "type": "table",
                "location": f"{location_table} / {location_row} / col {col_idx}",
                "old": old_val, "new": new_val, "source": source,
            }
            if match_pass >= 2:
                change_record["match_quality"] = f"degraded_pass_{match_pass}"
                change_record["attempted_row"] = row_label
                change_record["actual_row"] = matched_row_label
            changes.append(change_record)
            continue

        # Check if the update was already applied (cell already has new_val)
        already_applied = new_val and cell_or_text and (
            new_val in (cell_or_text or "")
            or _normalize_ws(new_val) in _normalize_ws(cell_or_text or "")
        )
        if already_applied:
            log.info(
                "Table update ALREADY APPLIED: page %d, '%s' already present "
                "(skipping '%s' -> '%s')",
                page, new_val, old_val, new_val,
            )
        elif matched_row_label is not None:
            log.warning(
                "Table update NOT FOUND: page %d, '%s' -> '%s' "
                "(closest row '%s' col %d has '%s')",
                page, old_val, new_val, matched_row_label, col_idx, cell_or_text
            )
        else:
            log.warning(
                "Table update NOT FOUND: page %d, '%s' -> '%s' "
                "(no matching table/row for table_name='%s' row_label='%s')",
                page, old_val, new_val, tbl_name, row_label
            )

    # --- Text (narrative) updates ---
    for upd in mappings.get("text_updates", []):
        page = upd["page"]
        old_txt = upd["old_text"]
        new_txt = upd["new_text"]
        source = upd.get("source", "")

        try:
            slide = prs.slides[page - 1]
        except IndexError:
            log.warning("Text update SKIPPED: page %d does not exist", page)
            continue
        found = False
        for shape in slide.shapes:
            if not dry_run:
                if _replace_in_shape(shape, old_txt, new_txt):
                    changes.append({
                        "page": page, "type": "text",
                        "location": shape.name,
                        "old": old_txt, "new": new_txt, "source": source,
                    })
                    found = True
                    break
            else:
                # Dry-run: check if the text exists without modifying
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if old_txt in para.text:
                            changes.append({
                                "page": page, "type": "text",
                                "location": shape.name,
                                "old": old_txt, "new": new_txt, "source": source,
                            })
                            found = True
                            break
                if found:
                    break

        if not found:
            log.warning("Text update NOT FOUND: page %d, '%s' -> '%s'",
                        page, old_txt, new_txt)

    # --- Row inserts (capped per table to prevent catastrophic insertions) ---
    _MAX_ROW_INSERTS_PER_TABLE = 6
    _row_insert_counts: dict[tuple[int, str], int] = {}  # (page, table_name) -> count

    for ins in mappings.get("row_inserts", []):
        page = ins["page"]
        tbl_name = ins.get("table_name", "")
        ref_label = ins.get("insert_after_row_label", "")
        cell_values = ins.get("cells", [])
        source = ins.get("source", "")

        # Guard: cap row_inserts per table to prevent table structure corruption
        table_key = (page, tbl_name)
        prior = _row_insert_counts.get(table_key, 0)
        if prior >= _MAX_ROW_INSERTS_PER_TABLE:
            log.warning(
                "Row insert SKIPPED (cap %d reached): page %d, table '%s', ref '%s'",
                _MAX_ROW_INSERTS_PER_TABLE, page, tbl_name, ref_label,
            )
            continue
        _row_insert_counts[table_key] = prior + 1

        try:
            slide = prs.slides[page - 1]
        except IndexError:
            log.warning("Row insert SKIPPED: page %d does not exist", page)
            continue

        # Find the target table
        tables = [s for s in slide.shapes if s.has_table]
        target_shape = None
        if tbl_name:
            for s in tables:
                if _loose_match(tbl_name, s.name):
                    target_shape = s
                    break
        # Fallback: find any table containing the reference row label
        if target_shape is None:
            for s in tables:
                if _find_row_by_label(s.table, ref_label) is not None:
                    target_shape = s
                    break

        if target_shape is None:
            log.warning("Row insert NOT FOUND: page %d, table '%s', "
                        "ref_label '%s'", page, tbl_name, ref_label)
            continue

        table = target_shape.table
        ref_idx = _find_row_by_label(table, ref_label)
        if ref_idx is None:
            log.warning("Row insert ref row NOT FOUND: page %d, label '%s'",
                        page, ref_label)
            continue

        # Pad or truncate cells to match column count
        n_cols = len(table.rows[0].cells) if table.rows else 0
        if len(cell_values) < n_cols:
            cell_values = cell_values + [""] * (n_cols - len(cell_values))
        elif len(cell_values) > n_cols:
            cell_values = cell_values[:n_cols]

        if not dry_run:
            _add_table_row(table, ref_idx, cell_values)

        location_table = target_shape.name or tbl_name or "<unnamed table>"
        changes.append({
            "page": page, "type": "row_insert",
            "location": f"{location_table} / after '{ref_label}'",
            "old": "(new row)", "new": " | ".join(cell_values),
            "source": source,
        })

    # --- Narrative updates (paragraph-level rewrites) ---
    for upd in mappings.get("narrative_updates", []):
        page = upd.get("page")
        old_narrative = upd.get("old_narrative", "")
        new_narrative = upd.get("new_narrative", "")
        source = upd.get("source", "")

        if not old_narrative or not new_narrative:
            continue

        try:
            slide = prs.slides[page - 1]
        except (IndexError, TypeError):
            log.warning("Narrative update SKIPPED: page %s does not exist", page)
            continue

        found = False
        for shape in slide.shapes:
            if not dry_run:
                if _replace_in_shape(shape, old_narrative, new_narrative):
                    changes.append({
                        "page": page, "type": "narrative",
                        "location": shape.name,
                        "old": old_narrative[:80] + ("..." if len(old_narrative) > 80 else ""),
                        "new": new_narrative[:80] + ("..." if len(new_narrative) > 80 else ""),
                        "source": source,
                    })
                    found = True
                    break
            else:
                if shape.has_text_frame:
                    full_text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                    if old_narrative in full_text:
                        changes.append({
                            "page": page, "type": "narrative",
                            "location": shape.name,
                            "old": old_narrative[:80] + ("..." if len(old_narrative) > 80 else ""),
                            "new": new_narrative[:80] + ("..." if len(new_narrative) > 80 else ""),
                            "source": source,
                        })
                        found = True
                        break
        if not found:
            log.warning("Narrative update NOT FOUND: page %d, '%s...'",
                        page, old_narrative[:60])

    # --- Table structure updates (add_column, remove_row, reorder_rows) ---
    for upd in mappings.get("table_structure_updates", []):
        page = upd.get("page")
        tbl_name = upd.get("table_name", "")
        action = upd.get("action", "")
        source = upd.get("source", "")

        try:
            slide = prs.slides[page - 1]
        except (IndexError, TypeError):
            log.warning("Table structure update SKIPPED: page %s does not exist", page)
            continue

        # Find target table
        tables = [s for s in slide.shapes if s.has_table]
        target_shape = None
        for s in tables:
            if _loose_match(tbl_name, s.name):
                target_shape = s
                break
        if target_shape is None and len(tables) == 1:
            target_shape = tables[0]
        if target_shape is None:
            log.warning("Table structure update NOT FOUND: page %d, table '%s'", page, tbl_name)
            continue

        table = target_shape.table
        location = target_shape.name or tbl_name

        if action == "add_column" and not dry_run:
            col_header = upd.get("column_header", "")
            values = upd.get("values", [])
            after_col = upd.get("after_column", len(table.columns) - 1)
            _add_table_column(table, after_col, col_header, values)
            changes.append({
                "page": page, "type": "table_structure",
                "location": f"{location} / add_column '{col_header}' after col {after_col}",
                "old": "(no column)", "new": col_header, "source": source,
            })
        elif action == "remove_row" and not dry_run:
            row_label = upd.get("row_label", "")
            row_idx = _find_row_by_label(table, row_label)
            if row_idx is not None:
                _remove_table_row(table, row_idx)
                changes.append({
                    "page": page, "type": "table_structure",
                    "location": f"{location} / remove_row '{row_label}'",
                    "old": row_label, "new": "(removed)", "source": source,
                })
            else:
                log.warning("Table structure remove_row: row '%s' not found", row_label)
        elif action == "reorder_rows" and not dry_run:
            new_order = upd.get("new_order", [])
            _reorder_table_rows(table, new_order)
            changes.append({
                "page": page, "type": "table_structure",
                "location": f"{location} / reorder_rows",
                "old": "(original order)", "new": " -> ".join(new_order[:5]),
                "source": source,
            })
        elif dry_run:
            changes.append({
                "page": page, "type": "table_structure",
                "location": f"{location} / {action}",
                "old": "(dry run)", "new": str(upd.get("column_header") or upd.get("row_label") or upd.get("new_order", [])[:3]),
                "source": source,
            })

    if not dry_run:
        prs.save(memo_path)
        log.info("Memo saved with %d updates.", len(changes))
    else:
        log.info("Dry-run: %d updates identified (not saved).", len(changes))

    # --- Chart updates (separate pass - charts need fresh prs load) ---
    chart_updates_list = mappings.get("chart_updates", [])
    if chart_updates_list:
        chart_changes = _apply_chart_updates(memo_path, chart_updates_list, dry_run=dry_run)
        changes.extend(chart_changes)

    return changes


def _apply_chart_updates(memo_path: str, chart_updates: list, dry_run: bool = False) -> list:
    """
    Update embedded PowerPoint chart data based on chart_updates from the
    mapping pass. Preserves all visual formatting (colors, styles, etc.).

    Returns a list of change records.
    """
    if not chart_updates:
        return []

    prs = _load_presentation(memo_path)
    changes = []

    for upd in chart_updates:
        page = upd["page"]
        chart_name = upd.get("chart_name", "")
        chart_title = upd.get("chart_title", "")
        series_name = upd.get("series_name", "")
        new_values = upd.get("new_values", [])
        old_values = upd.get("old_values", [])
        new_categories = upd.get("categories", None)  # noqa: F841 – reserved for future use
        source = upd.get("source", "")

        try:
            slide = prs.slides[page - 1]
        except IndexError:
            log.warning("Chart update SKIPPED: page %d does not exist", page)
            continue

        # Find the target chart
        target_chart = None
        target_shape = None
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            # Match by shape name or chart title
            name_match = chart_name and _loose_match(chart_name, shape.name)
            title_match = False
            if shape.chart.has_title and shape.chart.chart_title:
                try:
                    ct_text = shape.chart.chart_title.text_frame.text.strip()
                    title_match = chart_title and _loose_match(chart_title, ct_text)
                except Exception:
                    pass
            if name_match or title_match:
                target_chart = shape.chart
                target_shape = shape
                break

        if target_chart is None:
            # Fallback: if only one chart on the page, use it
            chart_shapes = [s for s in slide.shapes if s.has_chart]
            if len(chart_shapes) == 1:
                target_chart = chart_shapes[0].chart
                target_shape = chart_shapes[0]
                log.info("Chart update: single chart fallback on page %d", page)
            else:
                log.warning(
                    "Chart update NOT FOUND: page %d, name='%s', title='%s'",
                    page, chart_name, chart_title,
                )
                continue

        # Find and update the target series
        found_series = False
        for series in target_chart.series:
            s_name = ""
            try:
                s_name = series.name or ""
            except (AttributeError, IndexError):
                pass
            if not _loose_match(series_name, s_name):
                continue

            found_series = True
            if dry_run:
                changes.append({
                    "page": page, "type": "chart",
                    "location": f"{target_shape.name} / series '{s_name}'",
                    "old": str(old_values[:5]) + ("..." if len(old_values) > 5 else ""),
                    "new": str(new_values[:5]) + ("..." if len(new_values) > 5 else ""),
                    "source": source,
                })
                break

            # Update series values via the underlying XML cache
            try:
                _nsmap = {
                    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
                }
                el = series._element
                num_cache = el.find(".//c:numRef/c:numCache", _nsmap)
                if num_cache is None:
                    # Try numLit (inline values without external ref)
                    num_cache = el.find(".//c:numLit", _nsmap)
                if num_cache is not None:
                    pts = num_cache.findall("c:pt", _nsmap)
                    for i, pt in enumerate(pts):
                        if i < len(new_values):
                            v_el = pt.find("c:v", _nsmap)
                            if v_el is not None:
                                v_el.text = str(new_values[i])
                    changes.append({
                        "page": page, "type": "chart",
                        "location": f"{target_shape.name} / series '{s_name}'",
                        "old": str(old_values[:5]) + ("..." if len(old_values) > 5 else ""),
                        "new": str(new_values[:5]) + ("..." if len(new_values) > 5 else ""),
                        "source": source,
                    })
                else:
                    log.warning("Chart series '%s' has no numeric cache element", s_name)
            except Exception as e:
                log.warning("Chart update FAILED for series '%s': %s", s_name, e)
            break

        if not found_series:
            log.warning(
                "Chart series '%s' NOT FOUND in chart on page %d",
                series_name, page,
            )

    if not dry_run and changes:
        prs.save(memo_path)
        log.info("Chart updates saved: %d changes.", len(changes))

    return changes


# ============================================================================
# 9b. BRANDING REFORMAT
# ============================================================================
# Subtext Brand palette (from Subtext Brand Theme.thmx)
_BRAND_COLORS = [
    (0x2B, 0x28, 0x25),  # dk1 - near-black brown
    (0xFF, 0xFF, 0xFF),  # lt1 - white
    (0x16, 0x35, 0x2E),  # dk2 - deep forest green
    (0xF7, 0xF1, 0xE3),  # lt2 - warm cream
    (0xC1, 0xD1, 0x00),  # accent1 - lime/chartreuse
    (0xA9, 0x58, 0x18),  # accent3 - burnt orange
    (0x51, 0x22, 0x13),  # accent4 - dark mahogany
]


def _color_distance(c1: tuple, c2: tuple) -> float:
    """Euclidean RGB distance between two (r,g,b) tuples."""
    return ((c1[0]-c2[0])**2 + (c1[1]-c2[1])**2 + (c1[2]-c2[2])**2) ** 0.5


def _nearest_brand_color(r: int, g: int, b: int, threshold: float = 80.0):
    """
    Return the nearest brand color as an RGBColor if within threshold,
    otherwise None.
    """
    from pptx.dml.color import RGBColor
    best_dist = float("inf")
    best_color = None
    for bc in _BRAND_COLORS:
        d = _color_distance((r, g, b), bc)
        if d < best_dist:
            best_dist = d
            best_color = bc
    if best_dist <= threshold:
        return RGBColor(*best_color)
    return None


def apply_branding(memo_path: str, theme_path: str, cfg: dict) -> int:
    """
    Apply Subtext branding to the entire memo:
    1. Replace the PPTX theme XML with the Subtext Brand Theme
    2. Reformat all text runs to Pragmatica fonts
    3. Remap hard-coded colors to nearest brand color

    Returns the number of text runs reformatted.
    """
    import zipfile

    branding_cfg = cfg.get("branding", {})
    heading_threshold = branding_cfg.get("heading_size_threshold", 18)
    color_threshold = branding_cfg.get("color_distance_threshold", 80)
    heading_font = "Pragmatica Bold"
    body_font = "Pragmatica Book"

    # --- Step 1: Replace theme XML ---
    log.info("Replacing PPTX theme with Subtext Brand Theme...")

    # Extract theme XML from .thmx
    with zipfile.ZipFile(theme_path, "r") as thmx:
        theme_xml = thmx.read("theme/theme/theme1.xml")

    # Replace theme in PPTX (it's a zip archive)
    import io
    with zipfile.ZipFile(memo_path, "r") as zin:
        # Find the theme file path in the PPTX
        theme_entries = [n for n in zin.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
        if not theme_entries:
            log.warning("No theme XML found in PPTX - skipping theme replacement")
        else:
            theme_entry = theme_entries[0]
            # Rewrite the zip with the new theme
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    if item.filename == theme_entry:
                        zout.writestr(item, theme_xml)
                    else:
                        zout.writestr(item, zin.read(item.filename))
            # Write back
            with open(memo_path, "wb") as f:
                f.write(buf.getvalue())
            log.info("Theme replaced: %s", theme_entry)

    # --- Step 2 & 3: Reformat fonts and colors ---
    log.info("Reformatting fonts and colors...")
    prs = _load_presentation(memo_path)
    runs_reformatted = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # Process text frames (text boxes, placeholders)
            if shape.has_text_frame:
                _stype = str(shape.shape_type) if shape.shape_type is not None else ""
                is_title = ("TITLE" in _stype or "SUBTITLE" in _stype
                            or "CENTER_TITLE" in _stype
                            or (shape.name and shape.name.lower().startswith(("title", "subtitle"))))
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        _reformat_run(run, is_title, heading_threshold,
                                      heading_font, body_font, color_threshold)
                        runs_reformatted += 1

            # Process table cells (conservative: font only, preserve alignment & color)
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                # Determine heading from existing bold state or row 0
                                is_cell_heading = (
                                    row_idx == 0
                                    or run.font.bold is True
                                )
                                _reformat_run(
                                    run, is_cell_heading, heading_threshold,
                                    heading_font, body_font,
                                    color_threshold,
                                    skip_color=True,  # tables use deliberate colors
                                )
                                runs_reformatted += 1

    prs.save(memo_path)
    log.info("Branding applied: %d text runs reformatted", runs_reformatted)
    return runs_reformatted


def _reformat_run(run, is_heading_context: bool, size_threshold: int,
                  heading_font: str, body_font: str, color_threshold: float,
                  skip_color: bool = False):
    """Reformat a single text run's font and color, preserving bold/italic."""
    # Snapshot existing formatting BEFORE changes
    was_bold = run.font.bold
    was_italic = run.font.italic

    # Determine if this run is a heading
    font_size = run.font.size
    if font_size is not None:
        size_pt = font_size.pt
    else:
        size_pt = 0

    is_heading = is_heading_context or size_pt >= size_threshold or was_bold is True

    # Set font family
    run.font.name = heading_font if is_heading else body_font

    # Restore bold/italic (they may have been inherited from theme,
    # which we just replaced)
    if was_bold is not None:
        run.font.bold = was_bold
    if was_italic is not None:
        run.font.italic = was_italic

    # Remap color if it's a hard-coded RGB (skip for table cells to
    # preserve deliberate color coding in comp set / data tables)
    if skip_color:
        return
    try:
        color = run.font.color
        if color.type is not None and color.rgb is not None:
            r, g, b = color.rgb[0], color.rgb[1], color.rgb[2]
            nearest = _nearest_brand_color(r, g, b, color_threshold)
            if nearest is not None:
                run.font.color.rgb = nearest
    except (AttributeError, TypeError):
        pass  # No color set or theme color - skip


# ============================================================================
# 9c. LAYOUT NORMALIZATION
# ============================================================================
def normalize_layout(memo_path: str, cfg: dict) -> dict:
    """
    Post-update layout healing: align titles, page numbers, TOC table,
    and enforce margin bounds. Returns summary of changes made.
    """
    from collections import Counter
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.util import Inches

    layout_cfg = cfg.get("layout", {})
    margin_left = Inches(layout_cfg.get("margin_left", 0.50))
    margin_right = Inches(layout_cfg.get("margin_right", 0.50))
    margin_top = Inches(layout_cfg.get("margin_top", 0.25))
    margin_bottom = Inches(layout_cfg.get("margin_bottom", 0.50))
    snap_tol = Inches(layout_cfg.get("snap_tolerance", 0.05))

    prs = _load_presentation(memo_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    summary = {
        "titles_snapped": 0,
        "page_numbers_snapped": 0,
        "section_headers_aligned": 0,
        "toc_formatted": False,
        "shapes_clamped_to_margins": 0,
        "auto_size_applied": 0,
    }

    # ------------------------------------------------------------------
    # 1a. Title alignment - snap outlier titles to the dominant position
    # ------------------------------------------------------------------
    title_positions = []  # list of (slide_idx, shape, left, top, width, height)
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue  # skip cover
        for shape in slide.placeholders:
            try:
                ph_idx = shape.placeholder_format.idx
            except Exception:
                ph_idx = None
            is_title = (
                ph_idx == 0
                or (shape.name and shape.name.lower().startswith("title"))
            )
            if is_title:
                title_positions.append(
                    (idx, shape, shape.left, shape.top, shape.width, shape.height)
                )
                break  # one title per slide

    if title_positions:
        mode_left = Counter(p[2] for p in title_positions).most_common(1)[0][0]
        mode_top = Counter(p[3] for p in title_positions).most_common(1)[0][0]
        mode_width = Counter(p[4] for p in title_positions).most_common(1)[0][0]
        mode_height = Counter(p[5] for p in title_positions).most_common(1)[0][0]

        for slide_idx, shape, left, top, width, height in title_positions:
            snapped = False
            if abs(left - mode_left) > snap_tol:
                shape.left = mode_left
                snapped = True
            if abs(top - mode_top) > snap_tol:
                shape.top = mode_top
                snapped = True
            if abs(width - mode_width) > snap_tol:
                shape.width = mode_width
                snapped = True
            if abs(height - mode_height) > snap_tol:
                shape.height = mode_height
                snapped = True
            if snapped:
                summary["titles_snapped"] += 1
                log.info("Title snapped on slide %d", slide_idx + 1)

    # ------------------------------------------------------------------
    # 1b. Page number alignment
    # ------------------------------------------------------------------
    pn_positions = []
    for idx, slide in enumerate(prs.slides):
        for shape in slide.placeholders:
            try:
                ph_idx = shape.placeholder_format.idx
            except Exception:
                ph_idx = None
            is_pn = (
                ph_idx == 12
                or (shape.name and "slide number" in shape.name.lower())
            )
            if is_pn:
                pn_positions.append(
                    (idx, shape, shape.left, shape.top, shape.width, shape.height)
                )
                break

    if pn_positions:
        pn_mode_left = Counter(p[2] for p in pn_positions).most_common(1)[0][0]
        pn_mode_top = Counter(p[3] for p in pn_positions).most_common(1)[0][0]
        pn_mode_width = Counter(p[4] for p in pn_positions).most_common(1)[0][0]
        pn_mode_height = Counter(p[5] for p in pn_positions).most_common(1)[0][0]

        for slide_idx, shape, left, top, width, height in pn_positions:
            snapped = False
            if abs(left - pn_mode_left) > snap_tol:
                shape.left = pn_mode_left
                snapped = True
            if abs(top - pn_mode_top) > snap_tol:
                shape.top = pn_mode_top
                snapped = True
            if abs(width - pn_mode_width) > snap_tol:
                shape.width = pn_mode_width
                snapped = True
            if abs(height - pn_mode_height) > snap_tol:
                shape.height = pn_mode_height
                snapped = True
            if snapped:
                summary["page_numbers_snapped"] += 1
                log.info("Page number snapped on slide %d", slide_idx + 1)

    # ------------------------------------------------------------------
    # 1c. Section header alignment (Content Placeholder 2 with short text ending in ":")
    # ------------------------------------------------------------------
    section_headers = []
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue
        for shape in slide.placeholders:
            if not (shape.name and "content placeholder 2" in shape.name.lower()):
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text and len(text) < 60 and text.endswith(":"):
                section_headers.append((idx, shape))

    if section_headers:
        sh_mode_left = Counter(s[1].left for s in section_headers).most_common(1)[0][0]
        sh_mode_width = Counter(s[1].width for s in section_headers).most_common(1)[0][0]

        for slide_idx, shape in section_headers:
            aligned = False
            if abs(shape.left - sh_mode_left) > snap_tol:
                shape.left = sh_mode_left
                aligned = True
            if abs(shape.width - sh_mode_width) > snap_tol:
                shape.width = sh_mode_width
                aligned = True
            if aligned:
                summary["section_headers_aligned"] += 1
                log.info("Section header aligned on slide %d", slide_idx + 1)

    # ------------------------------------------------------------------
    # 1d. TOC table formatting
    # ------------------------------------------------------------------
    for slide in prs.slides:
        is_toc = False
        for shape in slide.placeholders:
            try:
                ph_idx = shape.placeholder_format.idx
            except Exception:
                ph_idx = None
            if ph_idx == 0 or (shape.name and shape.name.lower().startswith("title")):
                title_text = shape.text_frame.text.lower() if shape.has_text_frame else ""
                if "table of contents" in title_text or "contents" in title_text:
                    is_toc = True
                    break
        if not is_toc:
            continue

        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            n_cols = len(table.columns)
            if n_cols < 2:
                continue

            # Right-align last column, left-align others
            for row in table.rows:
                for col_idx, cell in enumerate(row.cells):
                    for para in cell.text_frame.paragraphs:
                        if col_idx == n_cols - 1:
                            para.alignment = PP_ALIGN.RIGHT
                        else:
                            para.alignment = PP_ALIGN.LEFT

            # Enforce minimum row height (0.45")
            min_row_height = Inches(0.45)
            for row in table.rows:
                if row.height < min_row_height:
                    row.height = min_row_height

            # Set column widths proportionally: ~45% / ~40% / ~15%
            table_width = sum(col.width for col in table.columns)
            if n_cols >= 3:
                table.columns[0].width = int(table_width * 0.45)
                table.columns[1].width = int(table_width * 0.40)
                table.columns[n_cols - 1].width = int(table_width * 0.15)
                # Distribute remainder to middle columns if > 3
                if n_cols > 3:
                    assigned = (int(table_width * 0.45) + int(table_width * 0.15)
                                + int(table_width * 0.40))
                    remainder = table_width - assigned
                    for c in range(2, n_cols - 1):
                        table.columns[c].width = remainder // (n_cols - 3)
            elif n_cols == 2:
                table.columns[0].width = int(table_width * 0.85)
                table.columns[1].width = int(table_width * 0.15)

            summary["toc_formatted"] = True
            log.info("TOC table formatted")
        if is_toc:
            break  # only one TOC slide

    # ------------------------------------------------------------------
    # 1e. Margin enforcement
    # ------------------------------------------------------------------
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    right_limit = slide_width - margin_right
    bottom_limit = slide_height - margin_bottom

    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue  # skip cover
        for shape in slide.shapes:
            # Skip pictures/images (may be intentionally full-bleed)
            try:
                if shape.shape_type in (
                    MSO_SHAPE_TYPE.PICTURE,
                    MSO_SHAPE_TYPE.LINKED_PICTURE,
                    MSO_SHAPE_TYPE.MEDIA,
                ):
                    continue
            except Exception:
                pass

            clamped = False

            # Clamp left
            if shape.left < margin_left:
                shape.left = margin_left
                clamped = True

            # Clamp top
            if shape.top < margin_top:
                shape.top = margin_top
                clamped = True

            # Clamp right overflow: shrink width first, then shift
            if shape.left + shape.width > right_limit:
                overflow = (shape.left + shape.width) - right_limit
                if shape.width > overflow:
                    shape.width -= overflow
                else:
                    shape.left = right_limit - shape.width
                clamped = True

            # Clamp bottom overflow: shrink height first, then shift
            if shape.top + shape.height > bottom_limit:
                overflow = (shape.top + shape.height) - bottom_limit
                if shape.height > overflow:
                    shape.height -= overflow
                else:
                    shape.top = bottom_limit - shape.height
                clamped = True

            if clamped:
                summary["shapes_clamped_to_margins"] += 1
                log.info("Shape '%s' clamped to margins on slide %d",
                         shape.name, idx + 1)

    # ------------------------------------------------------------------
    # 1f. Text overflow protection
    # ------------------------------------------------------------------
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # Ensure word wrap is on
            if not tf.word_wrap:
                tf.word_wrap = True
            # Apply auto-size (shrink text to fit) on content shapes
            try:
                if tf.auto_size != MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    summary["auto_size_applied"] += 1
            except Exception:
                pass  # some shapes don't support auto_size

    # 1g. Footer normalization — DISABLED (was stomping on slide titles)
    summary["footer_fixes"] = 0

    # ------------------------------------------------------------------
    # 1h. Content density scoring & overflow flagging
    # ------------------------------------------------------------------
    overflow_slides: list[tuple[int, dict]] = []  # (slide_idx, metrics)
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue  # skip cover

        text_chars = 0
        table_rows = 0
        table_count = 0
        chart_count = 0
        shape_count = 0

        for shape in slide.shapes:
            shape_count += 1
            if shape.has_text_frame:
                text_chars += sum(len(p.text) for p in shape.text_frame.paragraphs)
            if shape.has_table:
                table_count += 1
                table_rows += len(shape.table.rows)
            if shape.has_chart:
                chart_count += 1

        # Heuristic thresholds for "too full"
        is_overflow = (
            text_chars > 1200          # lots of text
            or table_rows > 15         # table won't fit vertically
            or (text_chars > 600 and table_rows > 8)  # both crowded
            or shape_count > 12        # too many shapes
        )
        if is_overflow:
            overflow_slides.append((idx, {
                "text_chars": text_chars,
                "table_rows": table_rows,
                "table_count": table_count,
                "chart_count": chart_count,
                "shape_count": shape_count,
            }))

    summary["overflow_slides_detected"] = len(overflow_slides)
    if overflow_slides:
        log.warning(
            "Content density: %d slides exceed density threshold: %s",
            len(overflow_slides),
            [idx + 1 for idx, _ in overflow_slides],
        )

    # ------------------------------------------------------------------
    # 1h. Cross-slide formatting consistency
    # ------------------------------------------------------------------
    # Detect the dominant font for body text and table cells, then
    # normalize outliers to the dominant style.
    from collections import defaultdict

    body_font_usage: dict[str, int] = defaultdict(int)   # font_name -> count
    body_size_usage: dict[float, int] = defaultdict(int)  # size_pt -> count
    table_size_usage: dict[float, int] = defaultdict(int)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            body_font_usage[run.font.name] += 1
                        if run.font.size:
                            try:
                                body_size_usage[round(run.font.size.pt, 1)] += 1
                            except Exception:
                                pass
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    try:
                                        table_size_usage[round(run.font.size.pt, 1)] += 1
                                    except Exception:
                                        pass

    # Find dominant table cell font size
    dominant_table_size = None
    if table_size_usage:
        dominant_table_size = max(table_size_usage, key=table_size_usage.get)

    # Normalize: if a table cell has a font size that's >2pt off the
    # dominant size AND it's not a header row (>= 2pt larger), fix it.
    format_fixes = 0
    if dominant_table_size:
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                for row_idx, row in enumerate(shape.table.rows):
                    if row_idx == 0:
                        continue  # skip header row
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    try:
                                        size = round(run.font.size.pt, 1)
                                        diff = abs(size - dominant_table_size)
                                        if 0.5 < diff < 4.0:  # outlier but not intentional header
                                            run.font.size = Pt(dominant_table_size)
                                            format_fixes += 1
                                    except Exception:
                                        pass

    summary["table_font_size_normalized"] = format_fixes
    if format_fixes > 0:
        log.info("Normalized %d table cell font sizes to %.1fpt", format_fixes, dominant_table_size)

    prs.save(memo_path)
    log.info("Layout normalized: %s", summary)
    return summary


# ============================================================================
# 10. CHANGE LOG
# ============================================================================
def write_change_log(output_dir: str, all_changes: list, mappings: dict,
                     memo_path: str, proforma_path: str, backup_path: str,
                     run_metadata: dict | None = None):
    """Write a Markdown change-log summarizing every modification."""
    def _md_cell(value: str) -> str:
        return str(value).replace("|", "\\|").replace("\n", "<br>")

    ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_path = os.path.join(output_dir, "CHANGE_LOG.md")

    with open(log_path, "w", encoding="utf-8") as f:
        f.write("# Memo Automator - Change Log\n\n")
        f.write(f"**Run date:** {ts}\n\n")
        f.write(f"**Memo:** `{os.path.basename(memo_path)}`\n\n")
        f.write(f"**Proforma:** `{os.path.basename(proforma_path)}`\n\n")
        f.write(f"**Backup:** `{os.path.basename(backup_path)}`\n\n")
        f.write(f"**Total changes applied:** {len(all_changes)}\n\n")

        if run_metadata:
            f.write("## Run Telemetry\n\n")
            steps = run_metadata.get("steps", {})
            f.write(f"- Duration (sec): {run_metadata.get('run_duration_sec', 0):.2f}\n")
            f.write(f"- Mapping API calls: {run_metadata.get('mapping_api_calls', 0)}\n")
            f.write(f"- Validation API calls: {run_metadata.get('validation_api_calls', 0)}\n")
            cache_read = run_metadata.get("cache_read_tokens", 0)
            cache_write = run_metadata.get("cache_write_tokens", 0)
            if cache_read or cache_write:
                f.write(f"- Prompt cache: {cache_read:,} tokens read, {cache_write:,} tokens written\n")
            if steps:
                f.write("- Step timings (sec):\n")
                for k, v in steps.items():
                    f.write(f"  - {k}: {v:.2f}\n")
            f.write("\n")

        if run_metadata and run_metadata.get("accuracy"):
            acc = run_metadata["accuracy"]
            score = acc["confidence_score"]
            f.write("## Confidence Score\n\n")
            f.write(f"**Score: {score}/100**\n\n")
            f.write(f"- Coverage: {acc['coverage_pct']}%\n")
            f.write(f"- Rejection rate: {acc['rejection_rate_pct']}%\n")
            f.write(f"- Correction rate: {acc['correction_rate_pct']}%\n")
            f.write(f"- Match quality: {acc['match_quality_pct']}%\n")
            f.write(f"- Miss rate: {acc['miss_rate_pct']}%\n\n")

        f.write("## Applied Changes\n\n")
        f.write("| # | Page | Type | Location | Old | New | Source |\n")
        f.write("|---|------|------|----------|-----|-----|--------|\n")
        for i, c in enumerate(all_changes, 1):
            old_display = c['old'][:40] + "..." if len(c['old']) > 40 else c['old']
            new_display = c['new'][:40] + "..." if len(c['new']) > 40 else c['new']
            f.write(f"| {i} | {c['page']} | {c['type']} | "
                    f"{_md_cell(c['location'])} | {_md_cell(old_display)} | "
                    f"{_md_cell(new_display)} | {_md_cell(c['source'])} |\n")

        # Unvalidated pages warning
        unvalidated_pages = mappings.get("_unvalidated_pages", [])
        if unvalidated_pages:
            f.write("\n\n## ⚠️ Unvalidated Pages\n\n")
            f.write(
                "**The following pages could not be fully validated due to "
                "API response truncation. Changes on these pages passed "
                "through without QA review. Manual review is strongly "
                "recommended.**\n\n"
            )
            for pg in unvalidated_pages:
                f.write(f"- Page {pg}\n")

        # Rejected updates
        rejected = mappings.get("rejected", [])
        if rejected:
            f.write("\n\n## Rejected Updates\n\n")
            f.write("These were proposed but failed validation:\n\n")
            for rej in rejected:
                f.write(f"- **Reason:** {rej.get('reason', 'unknown')}\n")
                orig = rej.get("original", {})
                f.write(f"  - Entry: `{json.dumps(orig)}`\n")

        # Missed metrics
        missed = mappings.get("missed", [])
        if missed:
            f.write("\n\n## Potentially Missed Metrics\n\n")
            f.write("These may need manual review:\n\n")
            for miss in missed:
                f.write(f"- **Page {miss.get('page', '?')}:** "
                        f"{miss.get('description', '')} "
                        f"(source: {miss.get('source', 'unknown')})\n")

        # Raw Claude mappings for auditability
        f.write("\n\n## Raw Claude API Mappings\n\n")
        f.write("```json\n")
        f.write(json.dumps(mappings, indent=2))
        f.write("\n```\n")

    log.info("Change log written: %s", log_path)
    return log_path


# ============================================================================
# 11. MAIN
# ============================================================================
def main():
    args = parse_args()

    # Set log level from --verbose / --quiet flags
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    elif args.quiet:
        logging.getLogger().setLevel(logging.WARNING)

    # Load environment (.env file in script dir or working directory)
    load_dotenv(os.path.join(os.path.dirname(__file__), ".env"))
    load_dotenv()  # also check cwd

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        log.error("ANTHROPIC_API_KEY not set. Copy .env.example to .env and add your key.")
        sys.exit(1)

    # Validate inputs
    if not os.path.isfile(args.memo):
        log.error("Memo file not found: %s", args.memo)
        sys.exit(1)
    if not os.path.isfile(args.proforma):
        log.error("Proforma file not found: %s", args.proforma)
        sys.exit(1)

    # Validate file extensions
    memo_ext = os.path.splitext(args.memo)[1].lower()
    if memo_ext != ".pptx":
        log.error("Memo file must be .pptx, got '%s': %s", memo_ext, args.memo)
        sys.exit(1)
    proforma_ext = os.path.splitext(args.proforma)[1].lower()
    if proforma_ext not in (".xlsx", ".xlsm"):
        log.error("Proforma file must be .xlsx or .xlsm, got '%s': %s", proforma_ext, args.proforma)
        sys.exit(1)

    try:
        # Load config
        cfg = load_config(args.config)
        log.info("Config loaded from %s", args.config)

        # Output directory
        output_dir = args.output_dir or os.path.dirname(os.path.abspath(args.memo))
        os.makedirs(output_dir, exist_ok=True)

        # Initialize Claude client (shared for mapping + validation)
        log.info("Mapping model: %s  |  Validation model: %s",
                 cfg["claude"]["model"], cfg["claude"]["validation_model"])
        client = anthropic.Anthropic(
            api_key=api_key,
            max_retries=5,
            timeout=900.0,  # 15 min; needed for large batches and Opus thinking
        )
        run_meta = {"steps": {}, "mapping_api_calls": 0, "validation_api_calls": 0}
        run_started = time.time()

        # ---- Step 1: Backup ----
        log.info("=" * 60)
        log.info("STEP 1: Creating backup")
        log.info("=" * 60)
        step_started = time.time()
        backup_path = create_backup(args.memo, output_dir)
        run_meta["steps"]["backup"] = time.time() - step_started

        # ---- Step 2: Extract proforma data ----
        log.info("=" * 60)
        log.info("STEP 2: Extracting proforma data")
        log.info("=" * 60)
        step_started = time.time()
        proforma_data = extract_proforma_data(args.proforma, cfg)
        run_meta["steps"]["extract_proforma"] = time.time() - step_started

        # Save extraction for debugging / audit
        pf_dump = os.path.join(output_dir, "proforma_extract.txt")
        with open(pf_dump, "w", encoding="utf-8") as f:
            f.write(proforma_data)
        log.info("Proforma data saved to %s", pf_dump)

        # ---- Step 3: Extract memo content (full deck) ----
        log.info("=" * 60)
        log.info("STEP 3: Extracting memo content (all slides)")
        log.info("=" * 60)
        step_started = time.time()
        memo_content = extract_memo_content(args.memo, cfg)
        run_meta["steps"]["extract_memo"] = time.time() - step_started

        memo_dump = os.path.join(output_dir, "memo_extract.txt")
        with open(memo_dump, "w", encoding="utf-8") as f:
            f.write(memo_content)
        log.info("Memo content saved to %s", memo_dump)

        # ---- Step 4: Claude API - metric mapping ----
        log.info("=" * 60)
        log.info("STEP 4: Claude API - identifying metric mappings")
        log.info("=" * 60)
        step_started = time.time()
        BATCH_THRESHOLD = 80_000  # chars; above this, process slides in batches
        RATE_LIMIT_INTERVAL = 5  # seconds between API calls for rate limiting
        prompt_size = len(proforma_data) + len(memo_content)
        if prompt_size > BATCH_THRESHOLD:
            log.info("Large prompt (%d chars) - processing slides in batches of 3", prompt_size)
            memo_chunks = chunk_memo_by_pages(memo_content, pages_per_chunk=3)
            mappings = {"table_updates": [], "text_updates": [], "row_inserts": []}
            last_api_call = 0
            for i, chunk in enumerate(memo_chunks, 1):
                if i > 1 and last_api_call > 0:
                    elapsed = time.time() - last_api_call
                    wait = RATE_LIMIT_INTERVAL - elapsed
                    if wait > 0:
                        log.info("Rate limit: waiting %.0f seconds (%.0fs elapsed since last call)...", wait, elapsed)
                        time.sleep(wait)
                    else:
                        log.info("Rate limit: no wait needed (%.0fs elapsed since last call)", elapsed)
                log.info("Mapping batch %d / %d ...", i, len(memo_chunks))
                last_api_call = time.time()
                try:
                    batch = get_metric_mappings(
                        client, proforma_data, chunk, cfg,
                        property_name=args.property_name,
                        telemetry=run_meta,
                    )
                except Exception as batch_err:
                    if _is_api_error(batch_err):
                        raise
                    log.warning("Batch %d failed (%s) - retrying as single-page sub-chunks", i, batch_err)
                    batch = {
                        "table_updates": [],
                        "text_updates": [],
                        "row_inserts": [],
                        "_truncated": True,
                    }

                # Retry truncated/failed batches with single-page sub-chunks
                if batch.pop("_truncated", False):
                    covered_pages = set()
                    for e in batch.get("table_updates", []):
                        covered_pages.add(e.get("page"))
                    for e in batch.get("text_updates", []):
                        covered_pages.add(e.get("page"))
                    for e in batch.get("row_inserts", []):
                        covered_pages.add(e.get("page"))

                    mappings["table_updates"].extend(batch.get("table_updates", []))
                    mappings["text_updates"].extend(batch.get("text_updates", []))
                    mappings["row_inserts"].extend(batch.get("row_inserts", []))

                    log.info(
                        "Retrying truncated batch %d with single-page sub-chunks (covered pages so far: %s)",
                        i, sorted(covered_pages)
                    )
                    sub_chunks = chunk_memo_by_pages(chunk, pages_per_chunk=1)
                    for j, sub_chunk in enumerate(sub_chunks, 1):
                        sub_pages = set(int(m) for m in re.findall(r"PAGE (\d+)", sub_chunk))
                        if sub_pages and sub_pages.issubset(covered_pages):
                            log.info("  Sub-chunk %d/%d (pages %s) already covered - skipping",
                                     j, len(sub_chunks), sorted(sub_pages))
                            continue

                        elapsed = time.time() - last_api_call
                        wait = RATE_LIMIT_INTERVAL - elapsed
                        if wait > 0:
                            log.info("Rate limit: waiting %.0f seconds...", wait)
                            time.sleep(wait)

                        log.info("  Sub-chunk %d/%d (pages %s)...", j, len(sub_chunks), sorted(sub_pages))
                        last_api_call = time.time()
                        try:
                            sub_batch = get_metric_mappings(
                                client, proforma_data, sub_chunk, cfg,
                                property_name=args.property_name,
                                telemetry=run_meta,
                            )
                        except Exception as sub_err:
                            if _is_api_error(sub_err):
                                raise
                            log.warning("  Sub-chunk %d failed (%s) - skipping pages %s",
                                        j, sub_err, sorted(sub_pages))
                            continue

                        if sub_batch.pop("_truncated", False):
                            log.warning("  Sub-chunk %d still truncated after single-page retry - moving on", j)
                        mappings["table_updates"].extend(sub_batch.get("table_updates", []))
                        mappings["text_updates"].extend(sub_batch.get("text_updates", []))
                        mappings["row_inserts"].extend(sub_batch.get("row_inserts", []))
                else:
                    mappings["table_updates"].extend(batch.get("table_updates", []))
                    mappings["text_updates"].extend(batch.get("text_updates", []))
                    mappings["row_inserts"].extend(batch.get("row_inserts", []))
        else:
            mappings = get_metric_mappings(
                client, proforma_data, memo_content, cfg,
                property_name=args.property_name,
                telemetry=run_meta,
            )
            mappings.pop("_truncated", None)
        run_meta["steps"]["mapping"] = time.time() - step_started

        # Save raw mappings for audit
        map_dump = os.path.join(output_dir, "mappings_raw.json")
        with open(map_dump, "w", encoding="utf-8") as f:
            json.dump(mappings, f, indent=2)
        log.info("Raw mappings saved to %s", map_dump)

        # ---- Step 4a: Strip no-op entries (old == new) ----
        pre_table = len(mappings["table_updates"])
        pre_text = len(mappings["text_updates"])
        mappings["table_updates"] = [
            e for e in mappings["table_updates"]
            if e.get("old_value") != e.get("new_value")
        ]
        mappings["text_updates"] = [
            e for e in mappings["text_updates"]
            if e.get("old_text") != e.get("new_text")
        ]
        n_stripped = (pre_table - len(mappings["table_updates"]) + pre_text - len(mappings["text_updates"]))
        if n_stripped > 0:
            log.info("Stripped %d no-op entries (old == new)", n_stripped)

        # ---- Step 4b: Pre-validation ----
        mappings = pre_validate_mappings(mappings, memo_content)

        # ---- Step 5: Claude API - validation pass ----
        step_started = time.time()
        if args.skip_validation:
            log.info("=" * 60)
            log.info("STEP 5: SKIPPED (--skip-validation flag)")
            log.info("=" * 60)
            validated = mappings
            validated.setdefault("rejected", [])
            validated.setdefault("missed", [])
        else:
            log.info("=" * 60)
            log.info("STEP 5: Claude API - validating mappings")
            log.info("=" * 60)
            validated = validate_mappings(
                client, mappings, proforma_data, memo_content, cfg,
                property_name=args.property_name,
                telemetry=run_meta,
            )
        run_meta["steps"]["validation"] = time.time() - step_started

        # Save validated mappings for audit
        val_dump = os.path.join(output_dir, "mappings_validated.json")
        with open(val_dump, "w", encoding="utf-8") as f:
            json.dump(validated, f, indent=2)
        log.info("Validated mappings saved to %s", val_dump)

        # ---- Step 6: Apply text / table updates ----
        log.info("=" * 60)
        log.info("STEP 6: Applying text / table updates")
        log.info("=" * 60)
        step_started = time.time()
        changes = apply_updates(args.memo, validated, dry_run=args.dry_run)
        run_meta["steps"]["apply_updates"] = time.time() - step_started

        # ---- Step 7: Change log ----
        log.info("=" * 60)
        log.info("STEP 7: Writing change log")
        log.info("=" * 60)
        run_meta["steps"].setdefault("write_changelog", 0.0)
        run_meta["run_duration_sec"] = time.time() - run_started
        step_started = time.time()
        log_path = write_change_log(
            output_dir, changes, validated,
            args.memo, args.proforma, backup_path,
            run_metadata=run_meta,
        )
        run_meta["steps"]["write_changelog"] = time.time() - step_started
        run_meta["run_duration_sec"] = time.time() - run_started

        # ---- Summary ----
        n_rejected = len(validated.get("rejected", []))
        n_missed = len(validated.get("missed", []))
        log.info("=" * 60)
        log.info("MEMO AUTOMATOR COMPLETE")
        log.info("=" * 60)
        log.info("  Changes applied:     %d", len(changes))
        log.info("  Rejected by QA:      %d", n_rejected)
        log.info("  Potentially missed:  %d", n_missed)
        log.info("  Mapping API calls:   %d", run_meta.get("mapping_api_calls", 0))
        log.info("  Validation API calls:%d", run_meta.get("validation_api_calls", 0))
        log.info("  Run duration (sec):  %.2f", run_meta.get("run_duration_sec", 0.0))
        log.info("  Backup:              %s", backup_path)
        log.info("  Change log:          %s", log_path)
        if args.dry_run:
            log.info("  ** DRY RUN -- no files were modified **")
        log.info("=" * 60)

    except ValueError as e:
        log.error("%s", e)
        sys.exit(1)
    except OSError as e:
        _exit_with_os_error(e, "reading/writing project files")
    except Exception as e:
        if _is_api_error(e):
            _exit_with_api_error(e)
        log.exception("Unexpected failure: %s", e)
        sys.exit(1)

if __name__ == "__main__":
    main()


