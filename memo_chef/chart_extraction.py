"""Market workbook extraction and chart mapping."""
from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from typing import Any

import openpyxl

log = logging.getLogger(__name__)


def extract_workbook_tables(
    workbook_path: str,
    tab_names: list[str] | None = None,
) -> str:
    """Extract tabular data from all (or specified) tabs as text.

    Uses the same format as extract_proforma_data() for consistency.
    """
    wb = openpyxl.load_workbook(workbook_path, data_only=True)
    lines: list[str] = []

    sheets = tab_names if tab_names else wb.sheetnames
    for tab_name in sheets:
        if tab_name not in wb.sheetnames:
            log.warning("Tab '%s' not found in workbook — skipping", tab_name)
            continue
        ws = wb[tab_name]
        tab_lines: list[str] = []
        for row in ws.iter_rows(values_only=False):
            row_data = [str(cell.value) for cell in row if cell.value is not None]
            if row_data:
                tab_lines.append(f"Row {row[0].row}:\t" + "\t".join(row_data))
        if tab_lines:
            lines.append(f"\n{'=' * 70}")
            lines.append(f"TAB: {tab_name}")
            lines.append(f"{'=' * 70}")
            lines.extend(tab_lines)

    wb.close()
    return "\n".join(lines)


def extract_memo_charts(memo_path: str) -> list[dict]:
    """Extract a structured list of charts from a PowerPoint memo.

    Returns a list of dicts describing each chart (page, shape name, title,
    series names) for use as context in chart mapping prompts.
    """
    from pptx import Presentation

    prs = Presentation(memo_path)
    charts: list[dict] = []
    for page_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            chart_title = ""
            if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
                try:
                    chart_title = chart.chart_title.text_frame.text.strip()
                except Exception:
                    pass
            series_names: list[str] = []
            try:
                for series in chart.series:
                    try:
                        series_names.append(series.name or "")
                    except (AttributeError, IndexError):
                        series_names.append("")
            except Exception:
                pass
            charts.append({
                "page": page_idx,
                "shape_name": shape.name,
                "chart_title": chart_title,
                "series_names": series_names,
            })
    return charts


def map_market_charts(
    workbook_text: str,
    memo_charts: list[dict],
    user_instructions: str,
    client: Any,
    model: str = "claude-haiku-4-5",
    max_tokens: int = 4096,
) -> list[dict]:
    """Use Claude to map workbook data to memo charts."""
    prompt_path = Path(__file__).parent.parent / "prompts" / "chart_mapping_v1.txt"
    template = prompt_path.read_text(encoding="utf-8")
    prompt = template.replace("{user_instructions}", user_instructions)
    max_chars = 50_000
    if len(workbook_text) > max_chars:
        log.warning("Workbook text truncated from %d to %d chars", len(workbook_text), max_chars)
    prompt = prompt.replace("{workbook_data}", workbook_text[:max_chars])
    prompt = prompt.replace("{memo_charts_json}", json.dumps(memo_charts, indent=2))

    response = client.messages.create(
        model=model,
        max_tokens=max_tokens,
        temperature=0,
        messages=[{"role": "user", "content": prompt}],
    )
    text = response.content[0].text.strip()

    json_match = re.search(r"\[[\s\S]*\]", text)
    if not json_match:
        log.warning("No JSON array in chart mapping response")
        return []

    return json.loads(json_match.group())
