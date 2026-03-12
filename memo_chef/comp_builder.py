"""Comp slide builder — normalize, deduplicate, and generate comp slides."""
from __future__ import annotations

import csv
import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from memo_chef.models import CompProperty, UnitMixEntry  # noqa: F401
from memo_chef.slide_insertion import (
    clone_slide,
    detect_memo_sections,  # noqa: F401
    find_template_slide,
    insert_slide_at_position,
)

log = logging.getLogger(__name__)

_SOURCE_PRIORITY = {"realpage": 0, "csv": 1, "url": 2, "manual": 3}

# Column name mapping: normalized key -> CompProperty field
_COLUMN_MAP = {
    "name": "name", "property name": "name", "property": "name",
    "address": "address",
    "units": "total_units", "total units": "total_units", "unit count": "total_units",
    "occupancy": "occupancy_pct", "occ": "occupancy_pct", "occupancy %": "occupancy_pct",
    "year built": "year_built", "year_built": "year_built", "vintage": "year_built",
    "distance": "distance_mi", "distance (mi)": "distance_mi",
    "concessions": "concessions",
}


def normalize_comps_from_csv(csv_path: str) -> list[CompProperty]:
    """Parse a CSV file into CompProperty objects."""
    path = Path(csv_path)
    comps: list[CompProperty] = []

    with open(path, newline="", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        for row in reader:
            data: dict[str, Any] = {"source": "csv", "source_detail": path.name}
            for col, val in row.items():
                key = col.strip().lower()
                field = _COLUMN_MAP.get(key)
                if field and val.strip():
                    v = val.strip()
                    if field in ("total_units", "year_built"):
                        try:
                            data[field] = int(float(v))
                        except ValueError:
                            pass
                    elif field in ("occupancy_pct", "distance_mi"):
                        try:
                            data[field] = float(v)
                        except ValueError:
                            pass
                    else:
                        data[field] = v
            if "name" in data:
                comps.append(CompProperty(**data))

    return comps


def normalize_comps_from_urls(
    comp_urls: list,
    extracted_texts: dict,
) -> list[CompProperty]:
    """Convert scraped comp URL data into CompProperty objects."""
    comps: list[CompProperty] = []
    for comp_url in comp_urls:
        url = comp_url.url if hasattr(comp_url, "url") else str(comp_url)
        label = comp_url.label if hasattr(comp_url, "label") else ""
        text = extracted_texts.get(url, "")
        if label or text:
            comps.append(CompProperty(
                name=label or url,
                source="url",
                source_detail=url,
            ))
    return comps


def deduplicate_comps(comps: list[CompProperty]) -> list[CompProperty]:
    """Fuzzy-match by name, merge fields with source priority."""
    from rapidfuzz import fuzz

    if not comps:
        return []

    # Sort by source priority (highest priority first)
    sorted_comps = sorted(comps, key=lambda c: _SOURCE_PRIORITY.get(c.source, 99))
    merged: list[CompProperty] = []

    for comp in sorted_comps:
        matched = False
        for i, existing in enumerate(merged):
            score = fuzz.token_set_ratio(comp.name, existing.name)
            if score > 85:
                # Merge: existing has priority, comp fills gaps
                merged_data = existing.model_dump()
                new_data = comp.model_dump()
                for field, val in new_data.items():
                    if field in ("source", "source_detail"):
                        continue
                    if val is not None and merged_data.get(field) is None:
                        merged_data[field] = val
                merged[i] = CompProperty(**merged_data)
                matched = True
                break
        if not matched:
            merged.append(comp)

    return merged


def build_comp_slide(
    prs: Presentation,
    subject: CompProperty,
    comps: list[CompProperty],
    memo_sections: list[dict],
    narrative: str | None = None,
) -> None:
    """Clone existing comp slide or build from scratch, populate with comp data."""
    template_idx = find_template_slide(prs, "Comp", "table", memo_sections)

    if template_idx is not None:
        new_slide = clone_slide(prs, template_idx)
        _populate_comp_table(new_slide, subject, comps)
        target_section = None
        for s in memo_sections:
            if "comp" in s["name"].lower() or "competitive" in s["name"].lower():
                target_section = s
                break
        if target_section:
            insert_slide_at_position(prs, new_slide, target_section["end_page"] - 1)
    else:
        _build_comp_slide_from_scratch(prs, subject, comps, memo_sections, narrative)


def _populate_comp_table(slide, subject: CompProperty, comps: list[CompProperty]) -> None:
    """Repopulate an existing cloned comp slide's table with new data."""
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            all_props = [subject] + comps
            for col_idx, prop in enumerate(all_props):
                if col_idx + 1 >= len(table.columns):
                    break
                _set_cell_safe(table, 0, col_idx + 1, prop.name)
                row_data = [
                    str(prop.total_units or ""),
                    f"{prop.occupancy_pct:.1f}%" if prop.occupancy_pct else "",
                    str(prop.year_built or ""),
                    f"{prop.distance_mi:.1f} mi" if prop.distance_mi else "",
                    prop.concessions or "",
                ]
                for row_idx, val in enumerate(row_data):
                    if row_idx + 1 < len(table.rows):
                        _set_cell_safe(table, row_idx + 1, col_idx + 1, val)
            break


def _set_cell_safe(table, row: int, col: int, text: str) -> None:
    """Set table cell text, preserving formatting."""
    try:
        cell = table.cell(row, col)
        if cell.text_frame.paragraphs:
            cell.text_frame.paragraphs[0].text = text
        else:
            cell.text = text
    except (IndexError, AttributeError):
        pass


def _build_comp_slide_from_scratch(
    prs: Presentation,
    subject: CompProperty,
    comps: list[CompProperty],
    memo_sections: list[dict],
    narrative: str | None = None,
) -> None:
    """Build a comp slide from scratch when no template is available."""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Rent Comparison"
    for para in tf.paragraphs:
        para.font.size = Pt(24)
        para.font.bold = True

    all_props = [subject] + comps[:6]
    row_labels = ["Property", "Units", "Occupancy", "Year Built", "Distance", "Concessions"]
    rows = len(row_labels)
    cols = len(all_props) + 1

    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.3), Inches(1.0), Inches(9.4), Inches(3.5)
    )
    table = table_shape.table

    for i, label in enumerate(row_labels):
        table.cell(i, 0).text = label

    for col_idx, prop in enumerate(all_props):
        table.cell(0, col_idx + 1).text = prop.name
        table.cell(1, col_idx + 1).text = str(prop.total_units or "")
        table.cell(2, col_idx + 1).text = f"{prop.occupancy_pct:.1f}%" if prop.occupancy_pct else ""
        table.cell(3, col_idx + 1).text = str(prop.year_built or "")
        table.cell(4, col_idx + 1).text = f"{prop.distance_mi:.1f} mi" if prop.distance_mi else ""
        table.cell(5, col_idx + 1).text = prop.concessions or ""

    if narrative:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = narrative
        for para in tf.paragraphs:
            para.font.size = Pt(11)

    for s in memo_sections:
        if "comp" in s["name"].lower() or "competitive" in s["name"].lower():
            insert_slide_at_position(prs, slide, s["end_page"] - 1)
            break
