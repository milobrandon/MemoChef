from __future__ import annotations

import json
import logging
import os
import random
import re
import time
from contextlib import contextmanager
from datetime import UTC, datetime
from pathlib import Path
from typing import Callable

import anthropic
import yaml

from memo_automator import (
    _is_api_error,
    apply_branding,
    apply_updates,
    build_mapping_batch_requests,
    chunk_memo_by_pages,
    create_backup,
    extract_market_data,
    extract_memo_content,
    extract_proforma_data,
    extract_schedule_data,
    get_metric_mappings,
    global_property_rename,
    load_config,
    normalize_layout,
    pre_validate_mappings,
    submit_and_poll_batch,
    validate_mapping_formats,
    validate_mappings,
    write_change_log,
)

from .models import RunManifest, RunRequest, RunResult, RunWarning, StageRecord, StageUpdate

log = logging.getLogger(__name__)

StageCallback = Callable[[StageUpdate], None] | None

LOG_FORMAT = "%(asctime)s  %(levelname)-8s  %(message)s"


class LogCapture(logging.Handler):
    def __init__(self) -> None:
        super().__init__()
        self.lines: list[str] = []

    def emit(self, record) -> None:
        self.lines.append(self.format(record))


class CheckpointManager:
    def __init__(self, request: RunRequest) -> None:
        self.request = request
        self.path = Path(request.output_dir) / "run_manifest.json"
        self.manifest = self._load_or_create()

    def _load_or_create(self) -> RunManifest:
        if self.request.resume_from_checkpoint and self.path.exists():
            try:
                return RunManifest.model_validate_json(self.path.read_text(encoding="utf-8"))
            except Exception as e:
                logging.getLogger(__name__).warning(
                    "Failed to load checkpoint %s, starting fresh: %s", self.path, e
                )
        return RunManifest(
            run_id=self.request.run_id,
            memo_name=Path(self.request.memo_path).name,
            proforma_name=Path(self.request.proforma_path).name,
            property_name=self.request.property_name,
            dry_run=self.request.dry_run,
            skip_validation=self.request.skip_validation,
        )

    def save(self) -> None:
        self.manifest.updated_at = datetime.now(UTC).isoformat()
        self.path.write_text(
            self.manifest.model_dump_json(indent=2),
            encoding="utf-8",
        )

    @contextmanager
    def stage(self, key: str, detail: str = ""):
        record = self.manifest.stages.get(key, StageRecord())
        record.status = "running"
        record.started_at = datetime.now(UTC).isoformat()
        record.detail = detail
        self.manifest.stages[key] = record
        self.save()
        started = time.time()
        try:
            yield record
            record.status = "completed"
        except Exception:
            record.status = "failed"
            raise
        finally:
            record.completed_at = datetime.now(UTC).isoformat()
            record.duration_seconds = round(time.time() - started, 2)
            self.manifest.stages[key] = record
            self.save()

    def add_warning(self, stage: str, message: str) -> None:
        self.manifest.warnings.append(RunWarning(stage=stage, message=message))
        self.save()

    def set_output(self, key: str, value: str) -> None:
        self.manifest.outputs[key] = value
        self.save()

    def set_count(self, key: str, value: int) -> None:
        self.manifest.counts[key] = value
        self.save()


# Approximate cost per million tokens (USD) by model prefix.
_TOKEN_RATES: dict[str, tuple[float, float]] = {
    "claude-opus": (15.0, 75.0),
    "claude-sonnet": (3.0, 15.0),
    "claude-haiku": (0.8, 4.0),
}


def _cost_usd(model: str, input_tokens: int, output_tokens: int) -> float:
    rate_in, rate_out = 3.0, 15.0  # default to Sonnet pricing
    for prefix, rates in _TOKEN_RATES.items():
        if prefix in model.lower():
            rate_in, rate_out = rates
            break
    return round((input_tokens * rate_in + output_tokens * rate_out) / 1_000_000, 6)


class _MessagesProxy:
    """Intercepts messages.create to accumulate token usage.

    Forwards all other attribute access to the real messages object so
    that code using client.messages.batches (etc.) still works.
    """

    def __init__(self, client: "anthropic.Anthropic", tracker: "TokenTracker") -> None:
        self._real_messages = client.messages
        self._tracker = tracker

    def __getattr__(self, name: str):
        return getattr(self._real_messages, name)

    def create(self, *args, **kwargs):
        response = self._real_messages.create(*args, **kwargs)
        if hasattr(response, "usage"):
            self._tracker.input_tokens += response.usage.input_tokens
            self._tracker.output_tokens += response.usage.output_tokens
            model = kwargs.get("model", "")
            self._tracker.estimated_cost_usd += _cost_usd(
                model, response.usage.input_tokens, response.usage.output_tokens
            )
        return response


class TokenTracker:
    """Wraps an Anthropic client and tracks cumulative token usage."""

    def __init__(self, client: "anthropic.Anthropic") -> None:
        self._client = client
        self.input_tokens: int = 0
        self.output_tokens: int = 0
        self.estimated_cost_usd: float = 0.0
        self.messages = _MessagesProxy(client, self)

    def __getattr__(self, name: str):
        return getattr(self._client, name)


def _deep_merge(base: dict, override: dict) -> dict:
    """Recursively merge override into base, returning a new dict."""
    result = dict(base)
    for key, value in override.items():
        if key in result and isinstance(result[key], dict) and isinstance(value, dict):
            result[key] = _deep_merge(result[key], value)
        else:
            result[key] = value
    return result


def _emit(callback: StageCallback, key: str, label: str, percent: int, detail: str = "") -> None:
    if callback is not None:
        callback(StageUpdate(key=key, label=label, percent=percent, detail=detail))


def _write_json(path: str, payload: dict) -> None:
    with open(path, "w", encoding="utf-8") as handle:
        json.dump(payload, handle, indent=2)


def _retry(
    func,
    *args,
    retries: int = 3,
    base_delay: float = 1.0,
    jitter: float = 0.25,
    checkpoint: CheckpointManager | None = None,
    stage: str = "",
    **kwargs,
):
    attempt = 0
    while True:
        try:
            return func(*args, **kwargs)
        except Exception as err:
            attempt += 1
            if attempt > retries or not _is_api_error(err):
                raise
            wait_seconds = base_delay * (2 ** (attempt - 1)) + random.uniform(0, jitter)
            if checkpoint is not None:
                checkpoint.add_warning(stage, f"Retrying after API error: {err}")
            time.sleep(wait_seconds)


def _mapping_with_batching(
    client,
    proforma_data: str,
    memo_content: str,
    cfg: dict,
    property_name: str | None,
    callback: StageCallback,
    checkpoint: CheckpointManager,
    source_directives: list[dict] | None = None,
) -> dict:
    batch_threshold = 80_000
    rate_limit_interval = 65
    prompt_size = len(proforma_data) + len(memo_content)
    if prompt_size <= batch_threshold:
        _emit(callback, "mapping", "Generate mappings", 52, "Sending full-deck mapping pass")
        mappings = _retry(
            get_metric_mappings,
            client,
            proforma_data,
            memo_content,
            cfg,
            property_name=property_name,
            source_directives=source_directives,
            checkpoint=checkpoint,
            stage="mapping",
        )
        mappings.pop("_truncated", None)
        return mappings

    memo_chunks = chunk_memo_by_pages(memo_content, pages_per_chunk=3)
    mappings = {"table_updates": [], "text_updates": [], "row_inserts": [], "narrative_updates": [], "table_structure_updates": []}
    last_api_call = 0.0
    for index, chunk in enumerate(memo_chunks, start=1):
        percent = 50 + int((18 * index) / max(len(memo_chunks), 1))
        _emit(callback, "mapping", f"Generate mappings ({index}/{len(memo_chunks)})", percent)
        if index > 1 and last_api_call > 0:
            wait_seconds = rate_limit_interval - (time.time() - last_api_call)
            if wait_seconds > 0:
                time.sleep(wait_seconds)
        last_api_call = time.time()
        batch = _retry(
            get_metric_mappings,
            client,
            proforma_data,
            chunk,
            cfg,
            property_name=property_name,
            source_directives=source_directives,
            checkpoint=checkpoint,
            stage="mapping",
        )
        if batch.pop("_truncated", False):
            covered_pages = {
                entry.get("page")
                for group in ("table_updates", "text_updates", "row_inserts", "narrative_updates", "table_structure_updates")
                for entry in batch.get(group, [])
            }
            mappings["table_updates"].extend(batch.get("table_updates", []))
            mappings["text_updates"].extend(batch.get("text_updates", []))
            mappings["row_inserts"].extend(batch.get("row_inserts", []))
            mappings["narrative_updates"].extend(batch.get("narrative_updates", []))
            mappings["table_structure_updates"].extend(batch.get("table_structure_updates", []))
            sub_chunks = chunk_memo_by_pages(chunk, pages_per_chunk=1)
            for sub_chunk in sub_chunks:
                sub_pages = set(int(match) for match in re.findall(r"PAGE (\d+)", sub_chunk))
                if sub_pages and sub_pages.issubset(covered_pages):
                    continue
                wait_seconds = rate_limit_interval - (time.time() - last_api_call)
                if wait_seconds > 0:
                    time.sleep(wait_seconds)
                last_api_call = time.time()
                sub_batch = _retry(
                    get_metric_mappings,
                    client,
                    proforma_data,
                    sub_chunk,
                    cfg,
                    property_name=property_name,
                    source_directives=source_directives,
                    checkpoint=checkpoint,
                    stage="mapping",
                )
                sub_batch.pop("_truncated", None)
                mappings["table_updates"].extend(sub_batch.get("table_updates", []))
                mappings["text_updates"].extend(sub_batch.get("text_updates", []))
                mappings["row_inserts"].extend(sub_batch.get("row_inserts", []))
                mappings["narrative_updates"].extend(sub_batch.get("narrative_updates", []))
                mappings["table_structure_updates"].extend(sub_batch.get("table_structure_updates", []))
            continue
        mappings["table_updates"].extend(batch.get("table_updates", []))
        mappings["text_updates"].extend(batch.get("text_updates", []))
        mappings["row_inserts"].extend(batch.get("row_inserts", []))
        mappings["narrative_updates"].extend(batch.get("narrative_updates", []))
        mappings["table_structure_updates"].extend(batch.get("table_structure_updates", []))

    # Deduplicate entries that appear in multiple chunks
    mappings = _dedup_mappings(mappings)
    return mappings


def _dedup_mappings(mappings: dict) -> dict:
    """Remove duplicate mapping entries produced by overlapping chunks."""
    seen_table: set[tuple] = set()
    deduped_table = []
    for upd in mappings.get("table_updates", []):
        key = (upd.get("page"), upd.get("table_name", ""),
               upd.get("row_label", ""), upd.get("column_index"),
               upd.get("old_value", ""))
        if key not in seen_table:
            seen_table.add(key)
            deduped_table.append(upd)

    seen_text: set[tuple] = set()
    deduped_text = []
    for upd in mappings.get("text_updates", []):
        key = (upd.get("page"), upd.get("old_text", ""))
        if key not in seen_text:
            seen_text.add(key)
            deduped_text.append(upd)

    seen_row: set[tuple] = set()
    deduped_row = []
    for ins in mappings.get("row_inserts", []):
        key = (ins.get("page"), ins.get("table_name", ""),
               ins.get("insert_after_row_label", ""),
               tuple(ins.get("cells", [])))
        if key not in seen_row:
            seen_row.add(key)
            deduped_row.append(ins)

    seen_narrative: set[tuple] = set()
    deduped_narrative = []
    for upd in mappings.get("narrative_updates", []):
        key = (upd.get("page"), upd.get("old_narrative", "")[:100])
        if key not in seen_narrative:
            seen_narrative.add(key)
            deduped_narrative.append(upd)

    n_removed = (
        len(mappings.get("table_updates", [])) - len(deduped_table)
        + len(mappings.get("text_updates", [])) - len(deduped_text)
        + len(mappings.get("row_inserts", [])) - len(deduped_row)
        + len(mappings.get("narrative_updates", [])) - len(deduped_narrative)
    )
    if n_removed > 0:
        logging.getLogger("memo_automator").info(
            "Deduplication removed %d duplicate mapping entries", n_removed
        )

    return {
        **mappings,
        "table_updates": deduped_table,
        "text_updates": deduped_text,
        "row_inserts": deduped_row,
        "narrative_updates": deduped_narrative,
        "table_structure_updates": mappings.get("table_structure_updates", []),
    }


def _mapping_with_batch_api(
    client,
    proforma_data: str,
    memo_content: str,
    cfg: dict,
    property_name: str | None,
    callback: StageCallback,
    checkpoint: CheckpointManager,
    source_directives: list[dict] | None = None,
) -> dict:
    """Submit all mapping chunks as a single Message Batch (50% cost).

    All chunks are processed in parallel by Anthropic's batch infrastructure.
    Typically completes within minutes, but may take up to 1 hour.
    """
    memo_chunks = chunk_memo_by_pages(memo_content, pages_per_chunk=3)
    _emit(callback, "mapping", f"Building batch ({len(memo_chunks)} chunks)", 46)

    requests = build_mapping_batch_requests(
        proforma_data, memo_chunks, cfg, property_name=property_name or "",
        source_directives=source_directives,
    )

    _emit(callback, "mapping", "Batch submitted — waiting for results", 50)
    results = submit_and_poll_batch(client, requests, poll_interval=15)

    # Merge results in order
    mappings: dict = {"table_updates": [], "text_updates": [], "row_inserts": [], "narrative_updates": []}
    for i in range(len(memo_chunks)):
        cid = f"mapping-chunk-{i}"
        batch_result = results.get(cid, {"table_updates": [], "text_updates": [], "row_inserts": []})
        mappings["table_updates"].extend(batch_result.get("table_updates", []))
        mappings["text_updates"].extend(batch_result.get("text_updates", []))
        mappings["row_inserts"].extend(batch_result.get("row_inserts", []))
        mappings["narrative_updates"].extend(batch_result.get("narrative_updates", []))
        mappings["table_structure_updates"].extend(batch_result.get("table_structure_updates", []))

    mappings = _dedup_mappings(mappings)
    checkpoint.set_count("batch_api_chunks", len(memo_chunks))
    return mappings


def _correction_retry(
    *,
    client,
    validated: dict,
    proforma_data: str,
    memo_content: str,
    cfg: dict,
    property_name: str | None,
    source_directives: list[dict] | None = None,
) -> dict | None:
    """Re-map rejected entries and missed metrics with feedback.

    Sends a targeted mapping call that includes:
    - The rejection reasons (so Claude avoids the same mistakes)
    - The missed metric descriptions (so Claude catches them this time)
    - Only the memo pages where rejections/misses occurred

    Returns a mappings dict with recovered entries, or None if nothing recovered.
    """
    rejected = validated.get("rejected", [])
    missed = validated.get("missed", [])
    if not rejected and not missed:
        return None

    # Collect pages that need re-mapping
    retry_pages: set[int] = set()
    for rej in rejected:
        orig = rej.get("original", {})
        if orig.get("page"):
            retry_pages.add(orig["page"])
    for miss in missed:
        if miss.get("page"):
            retry_pages.add(miss["page"])

    if not retry_pages:
        return None

    # Build feedback section for the prompt
    feedback_lines = [
        "\n## CORRECTION FEEDBACK — Fix these issues from the previous pass\n"
    ]
    if rejected:
        feedback_lines.append("### Rejected entries (do NOT repeat these mistakes):")
        for rej in rejected[:10]:  # cap to avoid prompt bloat
            orig = rej.get("original", {})
            feedback_lines.append(
                f"- Page {orig.get('page', '?')}: REJECTED because: {rej.get('reason', '?')}. "
                f"Original old_value='{orig.get('old_value', orig.get('old_text', '?'))[:50]}'"
            )
    if missed:
        feedback_lines.append("\n### Missed metrics (you MUST map these):")
        for miss in missed[:10]:
            feedback_lines.append(
                f"- Page {miss.get('page', '?')}: {miss.get('description', '?')}"
            )

    feedback_text = "\n".join(feedback_lines)

    # Extract only the relevant pages from memo content
    retry_chunks = []
    for page_match in re.finditer(r"(={60,}\nPAGE\s+(\d+)[^\n]*\n={60,})", memo_content):
        page_num = int(page_match.group(2))
        if page_num in retry_pages:
            start = page_match.start()
            # Find next page boundary
            next_match = re.search(r"={60,}\nPAGE\s+\d+", memo_content[page_match.end():])
            end = page_match.end() + next_match.start() if next_match else len(memo_content)
            retry_chunks.append(memo_content[start:end])

    if not retry_chunks:
        return None

    retry_memo = "\n".join(retry_chunks)

    # Inject feedback into proforma data so it's in the cached system message
    augmented_proforma = proforma_data + feedback_text

    log.info(
        "Correction retry: re-mapping %d pages (%d rejected, %d missed)",
        len(retry_pages), len(rejected), len(missed),
    )

    try:
        retry_result = get_metric_mappings(
            client,
            augmented_proforma,
            retry_memo,
            cfg,
            property_name=property_name,
            source_directives=source_directives,
        )
        retry_result.pop("_truncated", None)

        # Run pre-validation + format validation on the retry results
        retry_result = pre_validate_mappings(retry_result, memo_content)
        retry_result = validate_mapping_formats(retry_result)

        n = sum(
            len(retry_result.get(k, []))
            for k in ("table_updates", "text_updates", "row_inserts",
                      "narrative_updates", "table_structure_updates")
        )
        if n > 0:
            return retry_result
    except Exception as e:
        log.warning("Correction retry failed: %s", e)

    return None


def run_memo_pipeline(request: RunRequest, callback: StageCallback = None) -> RunResult:
    os.makedirs(request.output_dir, exist_ok=True)
    checkpoint = CheckpointManager(request)
    logger = logging.getLogger("memo_automator")
    log_capture = LogCapture()
    log_capture.setFormatter(logging.Formatter(LOG_FORMAT))
    logger.addHandler(log_capture)

    try:
        checkpoint.manifest.status = "running"
        checkpoint.manifest.config_profile = request.config_override_path and Path(request.config_override_path).stem
        checkpoint.save()
        cfg = load_config(request.config_path)
        if request.config_override_path and Path(request.config_override_path).exists():
            with open(request.config_override_path, encoding="utf-8") as f:
                override = yaml.safe_load(f) or {}
            cfg = _deep_merge(cfg, override)
        _raw_client = anthropic.Anthropic(
            api_key=request.api_key,
            max_retries=5,
            timeout=900.0,
        )
        client = TokenTracker(_raw_client)

        _emit(callback, "backup", "Create backup", 5)
        with checkpoint.stage("backup", "Creating backup copy"):
            backup_path = create_backup(request.memo_path, request.output_dir)
            checkpoint.set_output("backup_path", backup_path)

        # --- Property rename (before extraction so AI sees corrected name) ---
        effective_property_name = request.property_name
        if (
            request.property_name
            and request.property_rename_to
            and request.property_name.strip() != request.property_rename_to.strip()
        ):
            _emit(callback, "property_rename", "Rename property", 8)
            with checkpoint.stage("property_rename", "Renaming property across memo"):
                rename_count = global_property_rename(
                    request.memo_path,
                    request.property_name.strip(),
                    request.property_rename_to.strip(),
                )
                checkpoint.set_count("property_renames", rename_count)
            effective_property_name = request.property_rename_to.strip()
            checkpoint.manifest.property_rename_to = effective_property_name

        _emit(callback, "extract_sources", "Extract source data", 12)
        with checkpoint.stage("extract_sources", "Extracting proforma, market, and schedule data"):
            proforma_data = extract_proforma_data(request.proforma_path, cfg)
            proforma_extract_path = os.path.join(request.output_dir, "proforma_extract.txt")
            Path(proforma_extract_path).write_text(proforma_data, encoding="utf-8")
            checkpoint.set_output("proforma_extract", proforma_extract_path)

            if request.schedule_path:
                schedule_data = extract_schedule_data(request.schedule_path, cfg)
                if schedule_data:
                    proforma_data += "\n\n" + schedule_data
                    schedule_extract_path = os.path.join(request.output_dir, "schedule_extract.txt")
                    Path(schedule_extract_path).write_text(schedule_data, encoding="utf-8")
                    checkpoint.set_output("schedule_extract", schedule_extract_path)

            if request.market_data_path:
                market_data = extract_market_data(request.market_data_path, cfg)
                if market_data:
                    proforma_data += "\n\n" + market_data
                    market_extract_path = os.path.join(request.output_dir, "market_data_extract.txt")
                    Path(market_extract_path).write_text(market_data, encoding="utf-8")
                    checkpoint.set_output("market_data_extract", market_extract_path)
                else:
                    checkpoint.add_warning(
                        "extract_sources",
                        "Market data file loaded but no dashboard tabs were extracted.",
                    )

        if request.comp_urls:
            _emit(callback, "extract_comps", "Scrape comp URLs", 20)
            with checkpoint.stage("extract_comps", "Scraping competitive property websites"):
                from .extraction import extract_comp_urls

                comp_text = extract_comp_urls(
                    [cu.model_dump() for cu in request.comp_urls]
                )
                if comp_text.strip():
                    proforma_data += "\n\n## COMP PROPERTY DATA (scraped from websites)\n" + comp_text
                    comp_extract_path = os.path.join(request.output_dir, "comp_extract.txt")
                    Path(comp_extract_path).write_text(comp_text, encoding="utf-8")
                    checkpoint.set_output("comp_extract", comp_extract_path)

        _emit(callback, "extract_memo", "Extract memo", 24)
        with checkpoint.stage("extract_memo", "Extracting memo deck contents"):
            memo_content = extract_memo_content(request.memo_path, cfg)
            memo_extract_path = os.path.join(request.output_dir, "memo_extract.txt")
            Path(memo_extract_path).write_text(memo_content, encoding="utf-8")
            checkpoint.set_output("memo_extract", memo_extract_path)

        # Serialize source directives for prompt injection
        directives_dicts = [
            d.model_dump() for d in request.source_directives
        ] if request.source_directives else []

        _emit(callback, "mapping", "Generate mappings", 45)
        with checkpoint.stage("mapping", "Generating candidate updates"):
            mapping_fn = (
                _mapping_with_batch_api if request.use_batch_api
                else _mapping_with_batching
            )
            mappings = mapping_fn(
                client,
                proforma_data,
                memo_content,
                cfg,
                effective_property_name,
                callback,
                checkpoint,
                source_directives=directives_dicts,
            )
            mappings["table_updates"] = [
                entry for entry in mappings["table_updates"]
                if entry.get("old_value") != entry.get("new_value")
            ]
            mappings["text_updates"] = [
                entry for entry in mappings["text_updates"]
                if entry.get("old_text") != entry.get("new_text")
            ]
            mappings["narrative_updates"] = [
                entry for entry in mappings.get("narrative_updates", [])
                if entry.get("old_narrative") != entry.get("new_narrative")
            ]
            mappings = pre_validate_mappings(mappings, memo_content)
            mappings = validate_mapping_formats(mappings)
            raw_mapping_path = os.path.join(request.output_dir, "mappings_raw.json")
            _write_json(raw_mapping_path, mappings)
            checkpoint.set_output("mappings_raw", raw_mapping_path)

        _emit(callback, "validation", "Validate changes", 72)
        with checkpoint.stage("validation", "Validating mappings"):
            if request.skip_validation:
                validated = mappings
                validated.setdefault("rejected", [])
                validated.setdefault("missed", [])
            else:
                validated = _retry(
                    validate_mappings,
                    client,
                    mappings,
                    proforma_data,
                    memo_content,
                    cfg,
                    property_name=effective_property_name,
                    source_directives=directives_dicts,
                    checkpoint=checkpoint,
                    stage="validation",
                )

                # --- Correction retry loop ---
                # If validation rejected entries or found missed metrics,
                # re-map those pages with feedback and merge corrections.
                n_rejected = len(validated.get("rejected", []))
                n_missed = len(validated.get("missed", []))
                if n_rejected + n_missed > 0 and not request.skip_validation:
                    _emit(callback, "validation", "Re-mapping rejected entries", 76)
                    retry_mappings = _correction_retry(
                        client=client,
                        validated=validated,
                        proforma_data=proforma_data,
                        memo_content=memo_content,
                        cfg=cfg,
                        property_name=effective_property_name,
                        source_directives=directives_dicts,
                    )
                    if retry_mappings:
                        # Merge corrections into validated result
                        for key in ("table_updates", "text_updates", "row_inserts",
                                    "narrative_updates", "table_structure_updates"):
                            validated.setdefault(key, []).extend(
                                retry_mappings.get(key, [])
                            )
                        n_recovered = sum(
                            len(retry_mappings.get(k, []))
                            for k in ("table_updates", "text_updates", "row_inserts",
                                      "narrative_updates", "table_structure_updates")
                        )
                        checkpoint.set_count("correction_retry_recovered", n_recovered)
                        log.info("Correction retry recovered %d entries", n_recovered)

            unvalidated_pages = validated.get("_unvalidated_pages", [])
            if unvalidated_pages:
                checkpoint.add_warning(
                    "validation",
                    f"Pages {unvalidated_pages} could not be fully validated "
                    f"due to API response truncation. Changes on these pages "
                    f"passed without QA. Manual review recommended.",
                )
            validated_mapping_path = os.path.join(request.output_dir, "mappings_validated.json")
            _write_json(validated_mapping_path, validated)
            checkpoint.set_output("mappings_validated", validated_mapping_path)

        _emit(callback, "apply", "Apply updates", 84)
        with checkpoint.stage("apply", "Applying text, table, and chart updates"):
            changes = apply_updates(request.memo_path, validated, dry_run=request.dry_run)
            checkpoint.set_count("changes", len(changes))

        # --- Accuracy metrics ---
        from memo_chef.accuracy import compute_accuracy_metrics

        accuracy = compute_accuracy_metrics(
            raw=mappings,
            validated=validated,
            results=changes,
        )
        checkpoint.manifest.accuracy = accuracy
        checkpoint.set_count("confidence_score", int(accuracy["confidence_score"]))
        checkpoint.save()

        # --- Post-apply consistency check (loop until satisfied) ---
        if not request.dry_run and not request.skip_validation:
            from memo_automator import run_consistency_check

            _emit(callback, "consistency_check", "Verifying metric consistency", 85)
            with checkpoint.stage("consistency_check", "Post-apply metric tie-out"):
                max_fix_rounds = 2
                total_fixes = 0
                for fix_round in range(1, max_fix_rounds + 1):
                    # Re-extract the updated memo content
                    updated_memo = extract_memo_content(request.memo_path, cfg)

                    check_result = run_consistency_check(
                        client,
                        proforma_data,
                        updated_memo,
                        changes,
                        cfg,
                    )

                    status = check_result.get("status", "error")
                    discrepancies = check_result.get("discrepancies", [])

                    if status == "pass" or not discrepancies:
                        log.info(
                            "Consistency check PASSED (round %d): %s",
                            fix_round, check_result.get("summary", ""),
                        )
                        break

                    # Attempt to auto-fix discrepancies
                    critical = [d for d in discrepancies if d.get("severity") == "critical"]
                    minor = [d for d in discrepancies if d.get("severity") != "critical"]
                    log.warning(
                        "Consistency check round %d: %d critical, %d minor discrepancies",
                        fix_round, len(critical), len(minor),
                    )

                    # Build fix mappings from discrepancies that include fix data
                    fix_mappings = {"table_updates": [], "text_updates": []}
                    for d in discrepancies:
                        fix = d.get("fix")
                        if not fix:
                            continue
                        if fix.get("update_type") == "table":
                            fix_mappings["table_updates"].append({
                                "page": d.get("page"),
                                "table_name": d.get("location", ""),
                                "row_label": "",
                                "column_index": 1,
                                "old_value": fix.get("old_value", ""),
                                "new_value": fix.get("new_value", ""),
                                "source": fix.get("source", "consistency_check"),
                            })
                        else:
                            fix_mappings["text_updates"].append({
                                "page": d.get("page"),
                                "old_text": fix.get("old_value", ""),
                                "new_text": fix.get("new_value", ""),
                                "source": fix.get("source", "consistency_check"),
                            })

                    n_fixes = len(fix_mappings["table_updates"]) + len(fix_mappings["text_updates"])
                    if n_fixes > 0:
                        _emit(
                            callback, "consistency_check",
                            f"Fixing {n_fixes} discrepancies (round {fix_round})", 86,
                        )
                        fix_changes = apply_updates(
                            request.memo_path, fix_mappings, dry_run=False,
                        )
                        changes.extend(fix_changes)
                        total_fixes += len(fix_changes)
                        log.info("Applied %d consistency fixes (round %d)", len(fix_changes), fix_round)
                    else:
                        log.warning(
                            "Consistency check found %d issues but no auto-fixable entries",
                            len(discrepancies),
                        )
                        for d in discrepancies[:5]:
                            checkpoint.add_warning(
                                "consistency_check",
                                f"Page {d.get('page')}: {d.get('type')} — "
                                f"expected '{d.get('expected', '?')}', "
                                f"found '{d.get('found', '?')}'",
                            )
                        break

                checkpoint.set_count("consistency_fixes", total_fixes)
                consistency_path = os.path.join(request.output_dir, "consistency_check.json")
                _write_json(consistency_path, check_result)
                checkpoint.set_output("consistency_check", consistency_path)

        # --- Slide insertion (supplemental data) ---
        if request.supplemental_path and not request.dry_run:
            _emit(callback, "slide_insertion", "Insert supplemental slide", 82)
            with checkpoint.stage("slide_insertion", "Generating and inserting new slide"):
                try:
                    from memo_chef.extraction import extract_supplemental
                    from memo_chef.slide_insertion import (
                        analyze_supplemental_content,
                        build_slide_from_scratch,
                        clone_slide,
                        detect_memo_sections,
                        find_template_slide,
                        insert_slide_at_position,
                    )
                    from pptx import Presentation as PptxPresentation

                    supp_type = request.supplemental_type or "excel"
                    supplemental_text = extract_supplemental(request.supplemental_path, supp_type)
                    supp_path = os.path.join(request.output_dir, "supplemental_extract.txt")
                    Path(supp_path).write_text(supplemental_text, encoding="utf-8")
                    checkpoint.set_output("supplemental_extract", supp_path)

                    sections = detect_memo_sections(memo_content)
                    content = analyze_supplemental_content(
                        supplemental_text=supplemental_text,
                        memo_structure=sections,
                        client=client,
                        model=cfg.get("claude", {}).get("model", "claude-sonnet-4-6"),
                        user_brief=request.supplemental_brief,
                    )

                    prs = PptxPresentation(request.memo_path)
                    template_idx = find_template_slide(
                        prs, content["target_section"], content["visual_type"], sections
                    )

                    if template_idx is not None:
                        new_slide = clone_slide(prs, template_idx)
                    else:
                        new_slide = build_slide_from_scratch(prs, content)

                    target_after = content.get("target_after_slide", len(prs.slides)) - 1
                    insert_slide_at_position(prs, new_slide, target_after)
                    prs.save(request.memo_path)
                    checkpoint.set_count("slides_inserted", 1)
                    log.info(
                        "Inserted slide '%s' after slide %d",
                        content.get("slide_title", "Untitled"),
                        target_after + 1,
                    )
                except Exception as e:
                    log.error("Slide insertion failed: %s", e)
                    checkpoint.add_warning("slide_insertion", str(e))

        # --- Comp slide builder ---
        if request.auto_generate_comp_slide and not request.dry_run:
            _emit(callback, "comp_slide", "Build comp slide", 86)
            with checkpoint.stage("comp_slide", "Generating competitive analysis slide"):
                try:
                    from memo_chef.comp_builder import (
                        build_comp_slide,
                        deduplicate_comps,
                        normalize_comps_from_csv,
                        normalize_comps_from_urls,
                    )
                    from memo_chef.slide_insertion import detect_memo_sections
                    from pptx import Presentation as PptxPresentation

                    all_comps = []
                    if request.comp_csv_path:
                        all_comps.extend(normalize_comps_from_csv(request.comp_csv_path))
                    if request.comp_urls:
                        comp_texts = {}
                        for cu in request.comp_urls:
                            comp_extract_path = checkpoint.manifest.outputs.get("comp_extract")
                            if comp_extract_path and os.path.isfile(comp_extract_path):
                                comp_texts[cu.url] = Path(comp_extract_path).read_text(encoding="utf-8")
                        all_comps.extend(normalize_comps_from_urls(request.comp_urls, comp_texts))

                    if all_comps:
                        deduped = deduplicate_comps(all_comps)
                        sections = detect_memo_sections(memo_content)
                        subject = deduped[0]
                        comp_prs = PptxPresentation(request.memo_path)
                        build_comp_slide(comp_prs, subject, deduped[1:], sections)
                        comp_prs.save(request.memo_path)
                        checkpoint.set_count("comp_slides_inserted", 1)
                        log.info("Inserted comp slide with %d comps", len(deduped) - 1)
                    else:
                        log.warning("No comp data provided; skipping comp slide generation")
                except Exception as e:
                    log.error("Comp slide generation failed: %s", e)
                    checkpoint.add_warning("comp_slide", str(e))

        # --- Multi-slide generation (unified engine) ---
        # Runs when any source directive has scope "slide_generation" or "both",
        # or when supplemental data is provided with a directive targeting slides.
        slide_gen_directives = [
            d for d in directives_dicts
            if d.get("scope") in ("slide_generation", "both")
            and d.get("directive", "").strip()
        ]
        if slide_gen_directives and not request.dry_run:
            _emit(callback, "generate_slides", "Generate new slides", 88)
            with checkpoint.stage("generate_slides", "AI-driven multi-slide generation"):
                try:
                    from memo_chef.slide_generator import (
                        build_and_insert_slides,
                        extract_deck_profile,
                        generate_slide_plan,
                    )

                    deck_profile = extract_deck_profile(request.memo_path, memo_content)

                    # Gather all available source data for the slide prompt
                    slide_source_data = proforma_data
                    supp_extract = checkpoint.manifest.outputs.get("supplemental_extract")
                    if supp_extract and os.path.isfile(supp_extract):
                        slide_source_data += "\n\n## SUPPLEMENTAL DATA\n"
                        slide_source_data += Path(supp_extract).read_text(encoding="utf-8")

                    plan = generate_slide_plan(
                        source_data=slide_source_data,
                        memo_structure=deck_profile.sections,
                        deck_profile=deck_profile,
                        client=client,
                        model=cfg.get("claude", {}).get("model", "claude-sonnet-4-6"),
                        source_directives=slide_gen_directives,
                    )

                    n_inserted = build_and_insert_slides(
                        request.memo_path, plan, deck_profile,
                    )
                    checkpoint.set_count(
                        "ai_slides_generated",
                        n_inserted,
                    )
                    log.info("Multi-slide generation: %d slides created", n_inserted)
                except Exception as e:
                    log.error("Multi-slide generation failed: %s", e)
                    checkpoint.add_warning("generate_slides", str(e))

        if not request.dry_run:
            _emit(callback, "branding", "Apply branding", 90)
            with checkpoint.stage("branding", "Applying visual refresh"):
                theme_path = cfg.get("branding", {}).get("theme_path", "")
                if not theme_path:
                    theme_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Subtext Brand Theme.thmx")
                if os.path.exists(theme_path):
                    branded_count = apply_branding(request.memo_path, theme_path, cfg)
                    checkpoint.set_count("branded_runs", branded_count)
                else:
                    checkpoint.add_warning("branding", "Theme file not found; branding skipped.")

            _emit(callback, "layout", "Normalize layout", 94)
            with checkpoint.stage("layout", "Normalizing slide layout"):
                layout_summary = normalize_layout(request.memo_path, cfg)
                checkpoint.set_count("titles_snapped", int(layout_summary.get("titles_snapped", 0)))
                checkpoint.set_count(
                    "page_numbers_snapped",
                    int(layout_summary.get("page_numbers_snapped", 0)),
                )

        _emit(callback, "artifacts", "Write artifacts", 97)
        with checkpoint.stage("artifacts", "Writing change log and manifest"):
            log_path = write_change_log(
                request.output_dir,
                changes,
                validated,
                request.memo_path,
                request.proforma_path,
                checkpoint.manifest.outputs["backup_path"],
                run_metadata={"accuracy": accuracy} if accuracy else None,
            )
            checkpoint.set_output("change_log", log_path)
            checkpoint.set_count("rejected", len(validated.get("rejected", [])))
            checkpoint.set_count("missed", len(validated.get("missed", [])))
            checkpoint.set_count("input_tokens", client.input_tokens)
            checkpoint.set_count("output_tokens", client.output_tokens)
            # Store cost as integer microdollars to avoid float precision issues
            checkpoint.set_count(
                "estimated_cost_microdollars",
                int(round(client.estimated_cost_usd * 1_000_000)),
            )

        checkpoint.manifest.status = "completed"
        checkpoint.save()
        _emit(callback, "complete", "Run complete", 100)

        memo_bytes = Path(request.memo_path).read_bytes()
        log_bytes = Path(checkpoint.manifest.outputs["change_log"]).read_bytes()
        manifest_bytes = checkpoint.path.read_bytes()
        return RunResult(
            manifest=checkpoint.manifest,
            memo_path=request.memo_path,
            log_path=checkpoint.manifest.outputs["change_log"],
            manifest_path=str(checkpoint.path),
            memo_bytes=memo_bytes,
            log_bytes=log_bytes,
            manifest_bytes=manifest_bytes,
            changes=changes,
            rejected=validated.get("rejected", []),
            missed=validated.get("missed", []),
            unvalidated_pages=validated.get("_unvalidated_pages", []),
            log_lines=log_capture.lines[:],
        )
    except Exception as err:
        checkpoint.manifest.status = "failed"
        checkpoint.add_warning("pipeline", str(err))
        checkpoint.save()
        raise
    finally:
        logger.removeHandler(log_capture)
