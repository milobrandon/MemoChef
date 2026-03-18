"""Unified slide generation engine.

Replaces the separate slide_insertion + comp_builder approach with a single
pipeline: extract deck profile -> call Claude for a slide plan -> build and
insert each slide.

Existing callers (supplemental slide, comp slide) are wrapped by helpers
that produce SlideContent objects and feed them through the same builder.
"""
from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import DeckProfile, SlideContent, SlidePlan
from .slide_insertion import (
    _build_chart,
    _build_table,
    build_slide_from_scratch,
    clone_slide,
    detect_memo_sections,
    find_template_slide,
    insert_slide_at_position,
)

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 1. Deck profile extraction
# ---------------------------------------------------------------------------


def extract_deck_profile(
    memo_path: str,
    memo_text: str,
) -> DeckProfile:
    """Catalog the existing memo deck to guide slide generation style.

    Reads the PPTX to detect charts/tables/images and collects section info
    from the extracted memo text.
    """
    prs = Presentation(memo_path)
    sections = detect_memo_sections(memo_text)

    visual_types: set[str] = set()
    layout_names: set[str] = set()

    # Track font usage to find dominant fonts
    title_fonts: dict[tuple, int] = {}  # (name, size_pt) -> count
    body_fonts: dict[tuple, int] = {}

    for slide in prs.slides:
        if slide.slide_layout and slide.slide_layout.name:
            layout_names.add(slide.slide_layout.name)
        for shape in slide.shapes:
            if shape.has_chart:
                visual_types.add("chart")
            if shape.has_table:
                visual_types.add("table")
            if shape.shape_type and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                visual_types.add("image")
            # Extract font usage from text shapes
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        fname = run.font.name
                        fsize = run.font.size
                        if fname and fsize:
                            size_pt = fsize.pt if hasattr(fsize, "pt") else fsize / 12700
                            key = (fname, round(size_pt, 1))
                            if size_pt >= 18:
                                title_fonts[key] = title_fonts.get(key, 0) + 1
                            else:
                                body_fonts[key] = body_fonts.get(key, 0) + 1

    # Pick dominant title and body fonts
    title_font_name = None
    title_font_size = None
    body_font_name = None
    body_font_size = None
    if title_fonts:
        best = max(title_fonts, key=title_fonts.get)
        title_font_name, title_font_size = best
    if body_fonts:
        best = max(body_fonts, key=body_fonts.get)
        body_font_name, body_font_size = best

    return DeckProfile(
        sections=sections,
        total_slides=len(prs.slides),
        has_charts="chart" in visual_types,
        has_tables="table" in visual_types,
        slide_layouts_used=sorted(layout_names),
        visual_types_present=sorted(visual_types),
        title_font_name=title_font_name,
        title_font_size_pt=title_font_size,
        body_font_name=body_font_name,
        body_font_size_pt=body_font_size,
    )


def format_deck_profile(profile: DeckProfile) -> str:
    """Format a deck profile as text for prompt injection."""
    lines = [
        f"Total slides: {profile.total_slides}",
        f"Has charts: {profile.has_charts}",
        f"Has tables: {profile.has_tables}",
        f"Visual types present: {', '.join(profile.visual_types_present) or 'none'}",
        f"Slide layouts used: {', '.join(profile.slide_layouts_used) or 'unknown'}",
    ]
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# 2. Claude slide plan generation
# ---------------------------------------------------------------------------

_PROMPT_PATH = Path(__file__).parent.parent / "prompts" / "slide_generation_v1.txt"


def generate_slide_plan(
    source_data: str,
    memo_structure: list[dict],
    deck_profile: DeckProfile,
    client: Any,
    model: str = "claude-sonnet-4-6",
    max_tokens: int = 8192,
    source_directives: list[dict] | None = None,
) -> SlidePlan:
    """Call Claude to generate a multi-slide plan from source data."""
    from memo_automator import format_source_directives

    prompt_template = _PROMPT_PATH.read_text(encoding="utf-8")

    structure_text = "\n".join(
        f"  Slides {s['start_page']}-{s['end_page']}: {s['name']}"
        for s in memo_structure
    )

    directives_section = format_source_directives(
        source_directives or [], scope="slide_generation"
    )

    prompt = prompt_template.format(
        memo_structure=structure_text,
        deck_profile=format_deck_profile(deck_profile),
        source_data=source_data[:80_000],  # cap to avoid token limits
        source_directives_section=directives_section,
    )

    max_attempts = 2
    for attempt in range(1, max_attempts + 1):
        response = client.messages.create(
            model=model,
            max_tokens=max_tokens,
            temperature=0,
            messages=[{"role": "user", "content": prompt}],
        )

        raw = response.content[0].text.strip()

        # Parse JSON
        json_match = re.search(r"\{[\s\S]*\}", raw)
        if not json_match:
            if attempt < max_attempts:
                log.warning("Slide generation attempt %d: no valid JSON, retrying...", attempt)
                prompt += (
                    "\n\nIMPORTANT: You MUST respond with ONLY a JSON object. "
                    "Start with { and end with }. No text before or after."
                )
                continue
            log.warning("Slide generation returned no valid JSON after %d attempts", max_attempts)
            return SlidePlan()

        try:
            data = json.loads(json_match.group())
        except json.JSONDecodeError as e:
            if attempt < max_attempts:
                log.warning("Slide generation attempt %d: JSON parse error: %s, retrying...", attempt, e)
                prompt += (
                    "\n\nIMPORTANT: Your previous response had invalid JSON. "
                    "Return ONLY valid JSON. Start with { and end with }."
                )
                continue
            log.warning("Slide generation JSON parse error after %d attempts: %s", max_attempts, e)
            return SlidePlan()

        plan = SlidePlan.model_validate(data)

        # Validate that referenced sections exist in the memo structure
        valid_section_names = {s["name"].lower() for s in memo_structure}
        for slide_spec in plan.slides_to_generate:
            if slide_spec.section.lower() not in valid_section_names:
                log.warning(
                    "Slide '%s' references unknown section '%s' — "
                    "will insert at end of deck",
                    slide_spec.title, slide_spec.section,
                )

        return plan

    return SlidePlan()


# ---------------------------------------------------------------------------
# 3. Slide builder — turns SlideContent into actual PPTX slides
# ---------------------------------------------------------------------------


def build_and_insert_slides(
    memo_path: str,
    slide_plan: SlidePlan,
    deck_profile: DeckProfile,
) -> int:
    """Build and insert all slides from a plan. Returns count of slides inserted."""
    if not slide_plan.slides_to_generate:
        return 0

    prs = Presentation(memo_path)
    sections = deck_profile.sections
    inserted = 0

    # Sort by insert position so earlier inserts don't shift later positions
    sorted_slides = sorted(
        slide_plan.slides_to_generate,
        key=lambda s: s.insert_after_slide,
    )

    for slide_spec in sorted_slides:
        try:
            new_slide = _build_single_slide(prs, slide_spec, sections, deck_profile)
            if new_slide is not None:
                # Adjust for previously inserted slides
                target_idx = slide_spec.insert_after_slide - 1 + inserted
                insert_slide_at_position(prs, new_slide, target_idx)
                inserted += 1
                log.info(
                    "Inserted slide '%s' after position %d",
                    slide_spec.title,
                    target_idx + 1,
                )
        except Exception as e:
            log.error("Failed to build slide '%s': %s", slide_spec.title, e)

    if inserted > 0:
        prs.save(memo_path)
        log.info("Saved memo with %d new slides", inserted)

    return inserted


def _build_single_slide(
    prs: Presentation,
    spec: SlideContent,
    sections: list[dict],
    deck_profile: DeckProfile | None = None,
) -> Any | None:
    """Build one slide from a SlideContent spec.

    Tries to clone a template first; falls back to building from scratch
    using fonts from the deck profile when available.
    """
    visual_type = spec.visual_type or "table"

    # Try to find a template slide in the target section
    template_idx = find_template_slide(prs, spec.section, visual_type, sections)

    if template_idx is not None:
        new_slide = clone_slide(prs, template_idx)
        _populate_cloned_slide(new_slide, spec)
        return new_slide

    # Build from scratch using the same logic as existing slide_insertion
    content = {
        "slide_title": spec.title,
        "visual_type": visual_type,
        "visual_data": spec.visual_data,
        "narrative": spec.narrative,
    }
    return build_slide_from_scratch(prs, content, deck_profile=deck_profile)


def _populate_cloned_slide(slide, spec: SlideContent) -> None:
    """Repopulate a cloned template slide with new content."""
    # Update title (first text shape with large font or first shape)
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.font.size and para.font.size >= Pt(18):
                    para.text = spec.title
                    break
            else:
                continue
            break

    # Update table data if present
    if spec.visual_type == "table" or spec.content_type in ("table", "table_and_narrative"):
        for shape in slide.shapes:
            if shape.has_table:
                _repopulate_table(shape.table, spec.visual_data)
                break

    # Update narrative text (last large text box)
    if spec.narrative:
        text_shapes = [
            s for s in slide.shapes
            if s.has_text_frame and len(s.text_frame.paragraphs) > 0
        ]
        # Find the best narrative target: a shape with multiple paragraphs
        # or a shape near the bottom of the slide
        for shape in reversed(text_shapes):
            if len(shape.text_frame.paragraphs) > 1:
                shape.text_frame.paragraphs[0].text = spec.narrative
                # Clear remaining paragraphs
                for para in shape.text_frame.paragraphs[1:]:
                    para.text = ""
                break


def _repopulate_table(table, visual_data: dict) -> None:
    """Fill a table with new visual_data content."""
    categories = visual_data.get("categories", [])
    series_list = visual_data.get("series", [])
    if not categories or not series_list:
        return

    # Set title cell
    title = visual_data.get("title", "")
    if title and len(table.rows) > 0 and len(table.columns) > 0:
        try:
            table.cell(0, 0).text = title
        except (IndexError, AttributeError):
            pass

    # Set column headers (series names)
    for j, s in enumerate(series_list):
        if j + 1 < len(table.columns):
            try:
                table.cell(0, j + 1).text = s.get("name", "")
            except (IndexError, AttributeError):
                pass

    # Set row labels and values
    for i, cat in enumerate(categories):
        if i + 1 >= len(table.rows):
            break
        try:
            table.cell(i + 1, 0).text = str(cat)
        except (IndexError, AttributeError):
            pass
        for j, s in enumerate(series_list):
            if j + 1 >= len(table.columns):
                break
            vals = s.get("values", [])
            val = vals[i] if i < len(vals) else ""
            try:
                table.cell(i + 1, j + 1).text = str(val)
            except (IndexError, AttributeError):
                pass


# ---------------------------------------------------------------------------
# 5. Slide splitting — split overflowed slides into two
# ---------------------------------------------------------------------------

_SPLIT_PROMPT_PATH = Path(__file__).parent.parent / "prompts" / "slide_split_v1.txt"


def split_overflowed_slides(
    memo_path: str,
    overflow_slides: list[tuple[int, dict]],
    client: Any,
    model: str = "claude-sonnet-4-6",
) -> int:
    """Split slides that exceed content density thresholds.

    Parameters
    ----------
    memo_path: path to the PPTX
    overflow_slides: list of (slide_idx, metrics_dict) from normalize_layout
    client: Anthropic client
    model: Claude model to use

    Returns the number of slides that were successfully split.
    """
    if not overflow_slides:
        return 0

    prs = Presentation(memo_path)
    splits_done = 0
    offset = 0  # track index shifts from prior splits

    for orig_idx, metrics in overflow_slides:
        idx = orig_idx + offset
        if idx >= len(prs.slides):
            continue

        slide = prs.slides[idx]

        # Extract slide content text for Claude
        content_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                content_lines.append(
                    f"[{shape.name}] " +
                    " | ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                )
            if shape.has_table:
                for row_idx_t, row in enumerate(shape.table.rows):
                    cells = [c.text.strip() for c in row.cells]
                    content_lines.append(f"  Row {row_idx_t}: {' | '.join(cells)}")

        slide_content = "\n".join(content_lines)
        if not slide_content.strip():
            continue

        # Ask Claude how to split
        try:
            prompt_template = _SPLIT_PROMPT_PATH.read_text(encoding="utf-8")
            prompt = prompt_template.format(
                page_number=idx + 1,
                slide_content=slide_content[:8000],
                text_chars=metrics.get("text_chars", 0),
                table_rows=metrics.get("table_rows", 0),
                shape_count=metrics.get("shape_count", 0),
            )

            response = client.messages.create(
                model=model,
                max_tokens=4096,
                temperature=0,
                messages=[{"role": "user", "content": prompt}],
            )

            raw = response.content[0].text.strip()
            json_match = re.search(r"\{[\s\S]*\}", raw)
            if not json_match:
                log.warning("Slide split: no JSON returned for slide %d", idx + 1)
                continue

            split_plan = json.loads(json_match.group())
        except Exception as e:
            log.warning("Slide split failed for slide %d: %s", idx + 1, e)
            continue

        # Execute the split: for table splits, clone the slide and
        # remove rows from each copy
        strategy = split_plan.get("split_strategy", "")
        try:
            if strategy == "table_split":
                table_rows_end = split_plan.get("slide_1", {}).get("table_rows_end")
                table_rows_start = split_plan.get("slide_2", {}).get("table_rows_start")
                if table_rows_end is None or table_rows_start is None:
                    continue

                # Clone the slide for the continuation
                new_slide = clone_slide(prs, idx)
                insert_slide_at_position(prs, new_slide, idx)

                # Now we have two copies: prs.slides[idx] and prs.slides[idx+1]
                # Remove excess rows from each
                _trim_table_rows(prs.slides[idx], keep_end=table_rows_end)
                _trim_table_rows(prs.slides[idx + 1], keep_start=table_rows_start)

                # Update title on continuation slide
                cont_title = split_plan.get("slide_2", {}).get("title", "")
                if cont_title:
                    _update_slide_title(prs.slides[idx + 1], cont_title)

                splits_done += 1
                offset += 1
                log.info("Split slide %d (table_split at row %d)", orig_idx + 1, table_rows_end)

            elif strategy == "visual_narrative_split":
                # Clone slide, remove narrative from first, remove visual from second
                new_slide = clone_slide(prs, idx)
                insert_slide_at_position(prs, new_slide, idx)

                # First slide: remove large text shapes (narrative)
                _remove_narrative_shapes(prs.slides[idx])
                # Second slide: remove tables/charts, keep narrative
                _remove_visual_shapes(prs.slides[idx + 1])

                cont_title = split_plan.get("slide_2", {}).get("title", "")
                if cont_title:
                    _update_slide_title(prs.slides[idx + 1], cont_title)

                splits_done += 1
                offset += 1
                log.info("Split slide %d (visual_narrative_split)", orig_idx + 1)
        except (IndexError, KeyError, AttributeError) as e:
            log.warning("Slide split execution failed for slide %d: %s", idx + 1, e)
            continue

    if splits_done > 0:
        prs.save(memo_path)
        log.info("Saved memo after splitting %d overflowed slides", splits_done)

    return splits_done


def _trim_table_rows(slide, keep_end: int | None = None, keep_start: int | None = None):
    """Remove rows from the first table on a slide.

    keep_end: keep rows 0..keep_end (remove the rest)
    keep_start: keep rows 0 (header) + keep_start..end (remove middle)
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl_xml = shape.table._tbl
        rows = tbl_xml.findall(f"{{{ns}}}tr")
        n_rows = len(rows)
        if n_rows < 2:
            break  # nothing to trim

        if keep_end is not None:
            # Clamp to valid range
            keep_end = min(keep_end, n_rows - 1)
            for i in range(n_rows - 1, keep_end, -1):
                tbl_xml.remove(rows[i])
        elif keep_start is not None:
            # Clamp to valid range
            keep_start = max(1, min(keep_start, n_rows - 1))
            for i in range(keep_start - 1, 0, -1):
                tbl_xml.remove(rows[i])
        break  # only process first table


def _update_slide_title(slide, new_title: str):
    """Update the title text on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.font.size and para.font.size >= Pt(18):
                    para.text = new_title
                    return
    # Fallback: check placeholders
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.idx == 0:
                shape.text = new_title
                return
        except Exception:
            pass


def _remove_narrative_shapes(slide):
    """Remove large text-only shapes (narrative paragraphs) from a slide."""
    to_remove = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.has_table or shape.has_chart:
            continue
        total_text = sum(len(p.text) for p in shape.text_frame.paragraphs)
        if total_text > 200:  # likely narrative, not a title or label
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def _remove_visual_shapes(slide):
    """Remove tables and charts from a slide, keeping text shapes."""
    to_remove = []
    for shape in slide.shapes:
        if shape.has_table or shape.has_chart:
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
