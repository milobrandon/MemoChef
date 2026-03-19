"""Slide content analysis and insertion logic."""
from __future__ import annotations

import copy
import json
import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

log = logging.getLogger(__name__)

KNOWN_SECTIONS = [
    "Cover", "Table of Contents", "Executive Summary", "Investment Summary",
    "Market Summary", "Market Overview", "University Overview", "University Profile",
    "Site Overview", "Site Plan", "Competitive Landscape", "Comp Summary",
    "Financial Summary", "Financial Overview", "Development Budget",
    "Sources & Uses", "Sensitivity", "Appendix",
]


def detect_memo_sections(memo_text: str) -> list[dict]:
    """Parse memo extraction text to identify section boundaries.

    Returns list of {"name": str, "start_page": int, "end_page": int}.
    """
    sections: list[dict] = []
    current_page = 0

    for line in memo_text.splitlines():
        page_match = re.match(r"=+ PAGE (\d+) =+", line)
        if page_match:
            current_page = int(page_match.group(1))
            continue

        if "Para 0:" in line:
            text_match = re.search(r"Para 0:\s*'(.+?)'", line)
            if text_match:
                title = text_match.group(1).strip()
                for known in KNOWN_SECTIONS:
                    if known.lower() in title.lower():
                        sections.append({
                            "name": title,
                            "start_page": current_page,
                            "end_page": current_page,
                        })
                        break

    for i in range(len(sections) - 1):
        sections[i]["end_page"] = sections[i + 1]["start_page"] - 1
    if sections:
        sections[-1]["end_page"] = 999

    return sections


def analyze_supplemental_content(
    supplemental_text: str,
    memo_structure: list[dict],
    client: Any,
    model: str = "claude-sonnet-4-6",
    user_brief: str | None = None,
    max_tokens: int = 4096,
) -> dict[str, Any]:
    """Call Claude to analyze supplemental data and generate slide content."""
    prompt_path = Path(__file__).parent.parent / "prompts" / "slide_insertion_v1.txt"
    prompt_template = prompt_path.read_text(encoding="utf-8")

    structure_text = "\n".join(
        f"  Slides {s['start_page']}-{s['end_page']}: {s['name']}"
        for s in memo_structure
    )

    brief_section = ""
    if user_brief:
        brief_section = f"## USER BRIEF\n\n{user_brief}"

    prompt = prompt_template.replace("{memo_structure}", structure_text)
    prompt = prompt.replace("{supplemental_text}", supplemental_text[:50_000])
    prompt = prompt.replace("{user_brief_section}", brief_section)

    max_attempts = 2
    for attempt in range(1, max_attempts + 1):
        response = _call_claude(prompt, client, model, max_tokens)
        text = response.content[0].text.strip()

        json_match = re.search(r"\{[\s\S]*\}", text)
        if not json_match:
            if attempt < max_attempts:
                log.warning("Supplemental analysis attempt %d: no JSON, retrying...", attempt)
                prompt += (
                    "\n\nIMPORTANT: Return ONLY valid JSON. "
                    "Start with { and end with }."
                )
                continue
            raise ValueError(f"Claude returned no valid JSON for slide insertion:\n{text[:500]}")

        try:
            result = json.loads(json_match.group())
        except json.JSONDecodeError:
            if attempt < max_attempts:
                log.warning("Supplemental analysis attempt %d: invalid JSON, retrying...", attempt)
                prompt += (
                    "\n\nIMPORTANT: Your previous JSON was malformed. "
                    "Return ONLY valid JSON."
                )
                continue
            raise

        result["_tokens"] = {
            "input": response.usage.input_tokens,
            "output": response.usage.output_tokens,
        }
        return result

    raise ValueError("Supplemental analysis failed after all attempts")


def _call_claude(prompt: str, client: Any, model: str, max_tokens: int):
    """Make a Claude API call using the provided client."""
    return client.messages.create(
        model=model,
        max_tokens=max_tokens,
        temperature=0,
        messages=[{"role": "user", "content": prompt}],
    )


def find_template_slide(
    prs: Presentation,
    target_section: str,
    visual_type: str,
    sections: list[dict],
) -> int | None:
    """Find the best slide to clone as a template.

    Returns 0-based slide index, or None if no good match (score < 10).
    """
    best_idx = None
    best_score = 0

    target = None
    for s in sections:
        if target_section.lower() in s["name"].lower():
            target = s
            break

    if target is None:
        return None

    for idx, slide in enumerate(prs.slides):
        page = idx + 1
        score = 0

        if target["start_page"] <= page <= target["end_page"]:
            score += 10
        elif abs(page - target["start_page"]) <= 2:
            score += 5

        has_chart = any(shape.has_chart for shape in slide.shapes)
        has_table = any(shape.has_table for shape in slide.shapes)

        if visual_type in ("bar_chart", "line_chart", "pie_chart") and has_chart:
            score += 5
        elif visual_type == "table" and has_table:
            score += 5

        has_text = any(
            shape.has_text_frame and len(shape.text_frame.paragraphs) > 1
            for shape in slide.shapes
        )
        if has_text and (has_chart or has_table):
            score += 3

        if score > best_score:
            best_score = score
            best_idx = idx

    return best_idx if best_score >= 10 else None


_SLIDE_LAYOUT_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)
# Namespace for r:embed / r:link / r:id attributes
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _copy_slide_rels(src_slide, dst_slide) -> dict:
    """Copy non-layout part relationships from src_slide to dst_slide.

    Returns a mapping of old rId -> new rId so XML elements can be patched.
    Picture shapes embed images via r:embed rIds; without copying those
    relationships the images are missing on the cloned slide.
    """
    rId_map: dict[str, str] = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype == _SLIDE_LAYOUT_RELTYPE:
            continue  # layout relationship is already set on the new slide
        try:
            new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
            if new_rId != rId:
                rId_map[rId] = new_rId
        except Exception as exc:
            log.debug("clone_slide: could not copy relationship %s: %s", rId, exc)
    return rId_map


def _patch_rids(element, rId_map: dict) -> None:
    """Recursively replace r:embed / r:link / r:id attribute values per rId_map."""
    if not rId_map:
        return
    for attr in (f"{{{_R_NS}}}embed", f"{{{_R_NS}}}link", f"{{{_R_NS}}}id"):
        old = element.get(attr)
        if old and old in rId_map:
            element.set(attr, rId_map[old])
    for child in element:
        _patch_rids(child, rId_map)


def clone_slide(prs: Presentation, template_idx: int):
    """Deep-copy a slide and append it to the presentation. Returns the new slide."""
    template_slide = prs.slides[template_idx]
    slide_layout = template_slide.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # Remove default placeholder shapes
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy part relationships (images, charts, etc.) before copying XML elements
    rId_map = _copy_slide_rels(template_slide, new_slide)

    # Copy shapes from template, patching any rId references that changed
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        if rId_map:
            _patch_rids(el, rId_map)
        new_slide.shapes._spTree.append(el)

    return new_slide


def build_slide_from_scratch(prs: Presentation, content: dict, deck_profile=None):
    """Build a new slide with chart/table and narrative from scratch.

    When deck_profile is provided, uses its font names and sizes instead
    of hardcoded defaults, so generated slides match the existing deck.
    """
    # Resolve fonts from deck profile or use sensible defaults
    title_size = Pt(24)
    body_size = Pt(11)
    title_font = None
    body_font = None
    if deck_profile is not None:
        if deck_profile.title_font_size_pt:
            title_size = Pt(deck_profile.title_font_size_pt)
        if deck_profile.body_font_size_pt:
            body_size = Pt(deck_profile.body_font_size_pt)
        title_font = deck_profile.title_font_name
        body_font = deck_profile.body_font_name

    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    visual = content.get("visual_data", {})
    visual_type = content.get("visual_type", "table")

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = content.get("slide_title", "")
    for para in tf.paragraphs:
        para.font.size = title_size
        para.font.bold = True
        if title_font:
            para.font.name = title_font

    if visual_type == "table":
        _build_table(slide, visual)
    elif visual_type in ("bar_chart", "line_chart", "pie_chart"):
        _build_chart(slide, visual, visual_type)

    # Narrative
    narrative = content.get("narrative", "")
    if narrative:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = narrative
        for para in tf.paragraphs:
            para.font.size = body_size
            if body_font:
                para.font.name = body_font

    return slide


def _build_table(slide, visual: dict) -> None:
    """Add a table shape to the slide."""
    categories = visual.get("categories", [])
    series_list = visual.get("series", [])
    if not categories or not series_list:
        return

    rows = len(categories) + 1
    cols = len(series_list) + 1

    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9), Inches(4))
    table = table_shape.table

    table.cell(0, 0).text = visual.get("title", "")
    for j, s in enumerate(series_list):
        table.cell(0, j + 1).text = s.get("name", f"Series {j}")

    for i, cat in enumerate(categories):
        table.cell(i + 1, 0).text = str(cat)
        for j, s in enumerate(series_list):
            vals = s.get("values", [])
            val = vals[i] if i < len(vals) else ""
            table.cell(i + 1, j + 1).text = str(val)


def _build_chart(slide, visual: dict, chart_type_str: str) -> None:
    """Add a chart shape to the slide."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    type_map = {
        "bar_chart": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line_chart": XL_CHART_TYPE.LINE,
        "pie_chart": XL_CHART_TYPE.PIE,
    }
    xl_type = type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    categories = visual.get("categories", [])
    series_list = visual.get("series", [])
    if not categories or not series_list:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", "Data"), s.get("values", []))

    slide.shapes.add_chart(xl_type, Inches(0.5), Inches(1.0), Inches(9), Inches(4.2), chart_data)


def insert_slide_at_position(prs: Presentation, slide, after_slide_idx: int) -> None:
    """Move a slide (already appended at end) to position after after_slide_idx (0-based)."""
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)

    if not slide_ids:
        return

    new_slide_id = slide_ids[-1]
    slide_id_list.remove(new_slide_id)

    insert_pos = min(after_slide_idx + 1, len(slide_ids))
    slide_ids.insert(insert_pos, new_slide_id)

    for child in list(slide_id_list):
        slide_id_list.remove(child)
    for sid in slide_ids:
        slide_id_list.append(sid)
