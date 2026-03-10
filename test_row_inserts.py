"""
Unit tests for row_insert helpers and pipeline plumbing.
Tests _find_row_by_label, _add_table_row, and apply_updates with row_inserts.
No API calls required.
"""
import os
from pptx import Presentation

from memo_automator import (
    _find_row_by_label,
    _add_table_row,
    apply_updates,
    pre_validate_mappings,
    _salvage_truncated_json,
)


def _get_table_shape(slide):
    """Find the first table shape on a slide."""
    for s in slide.shapes:
        if s.has_table:
            return s
    return None


class TestFindRowByLabel:
    def test_exact_match(self, sample_pptx):
        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table

        assert _find_row_by_label(table, "1BR") == 1
        assert _find_row_by_label(table, "3BR") == 2
        assert _find_row_by_label(table, "4BR/2BA") == 3

    def test_loose_match(self, sample_pptx):
        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table

        assert _find_row_by_label(table, "Unit") == 0

    def test_not_found(self, sample_pptx):
        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table

        assert _find_row_by_label(table, "Studio") is None
        assert _find_row_by_label(table, "2BR") is None


class TestAddTableRow:
    def test_insert_row(self, sample_pptx):
        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table

        assert len(table.rows) == 4
        _add_table_row(table, 1, ["2BR/1BA", "80", "$1,495"])
        assert len(table.rows) == 5

        new_row_text = [table.cell(2, c).text for c in range(3)]
        assert new_row_text == ["2BR/1BA", "80", "$1,495"]

        # Original rows shifted correctly
        assert table.cell(1, 0).text == "1BR"
        assert table.cell(3, 0).text == "3BR"
        assert table.cell(4, 0).text == "4BR/2BA"

    def test_persists_after_save(self, sample_pptx, tmp_dir):
        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table
        _add_table_row(table, 1, ["2BR/1BA", "80", "$1,495"])

        path2 = os.path.join(tmp_dir, "saved.pptx")
        prs.save(path2)
        prs2 = Presentation(path2)
        table2 = _get_table_shape(prs2.slides[0]).table
        assert len(table2.rows) == 5
        assert table2.cell(2, 0).text == "2BR/1BA"

    def test_padding(self, sample_pptx):
        """_add_table_row handles fewer cells than columns."""
        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table

        _add_table_row(table, 2, ["Studio"])
        assert len(table.rows) == 5
        assert table.cell(3, 0).text == "Studio"
        assert table.cell(3, 1).text == ""
        assert table.cell(3, 2).text == ""


class TestApplyUpdatesRowInserts:
    def _make_mappings(self, **overrides):
        base = {
            "table_updates": [],
            "text_updates": [],
            "row_inserts": [
                {
                    "page": 1,
                    "table_name": "UnitMixTable",
                    "insert_after_row_label": "1BR",
                    "cells": ["2BR/1BA", "80", "$1,495"],
                    "source": "Proforma Unit Mix",
                }
            ],
        }
        base.update(overrides)
        return base

    def test_row_insert(self, sample_pptx):
        changes = apply_updates(sample_pptx, self._make_mappings(), dry_run=False)
        assert len(changes) == 1
        assert changes[0]["type"] == "row_insert"
        assert "UnitMixTable" in changes[0]["location"]

        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table
        assert len(table.rows) == 5
        assert table.cell(2, 0).text == "2BR/1BA"

    def test_dry_run(self, sample_pptx):
        changes = apply_updates(sample_pptx, self._make_mappings(), dry_run=True)
        assert len(changes) == 1

        # File should NOT be modified
        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table
        assert len(table.rows) == 4

    def test_fallback_wrong_table_name(self, sample_pptx):
        mappings = self._make_mappings()
        mappings["row_inserts"][0]["table_name"] = "WrongTableName"
        mappings["row_inserts"][0]["insert_after_row_label"] = "3BR"
        mappings["row_inserts"][0]["cells"] = ["Studio", "50", "$1,100"]

        changes = apply_updates(sample_pptx, mappings, dry_run=False)
        assert len(changes) == 1

        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table
        assert len(table.rows) == 5
        assert table.cell(3, 0).text == "Studio"

    def test_multiple_inserts(self, sample_pptx):
        mappings = {
            "table_updates": [],
            "text_updates": [],
            "row_inserts": [
                {
                    "page": 1, "table_name": "UnitMixTable",
                    "insert_after_row_label": "Unit Type",
                    "cells": ["Studio", "30", "$950"], "source": "Test",
                },
                {
                    "page": 1, "table_name": "UnitMixTable",
                    "insert_after_row_label": "1BR",
                    "cells": ["2BR/1BA", "80", "$1,495"], "source": "Test",
                },
            ],
        }
        changes = apply_updates(sample_pptx, mappings, dry_run=False)
        assert len(changes) == 2

        prs = Presentation(sample_pptx)
        table = _get_table_shape(prs.slides[0]).table
        assert len(table.rows) == 6


class TestSalvageTruncatedJson:
    def test_truncated_does_not_crash(self):
        truncated = (
            '{"table_updates": [{"page": 1, "old_value": "x", "new_value": "y"}], '
            '"text_updates": [], "row_inserts": [{"page": 1, "cells": ["a", "b"'
        )
        result = _salvage_truncated_json(truncated)
        if result is not None:
            assert "row_inserts" in result

    def test_nearly_complete(self):
        nearly = '{"table_updates": [], "text_updates": [], "row_inserts": [{"page": 1, "cells": ["a"]}]'
        result = _salvage_truncated_json(nearly)
        if result is not None:
            assert "row_inserts" in result
            assert len(result["row_inserts"]) == 1


class TestPreValidatePassesRowInserts:
    def test_passes_through(self):
        mappings = {
            "table_updates": [],
            "text_updates": [],
            "row_inserts": [
                {"page": 1, "table_name": "T", "insert_after_row_label": "1BR",
                 "cells": ["2BR", "80"], "source": "test"}
            ],
        }
        memo_content = "==============================\nPAGE 1\n=============================="
        result = pre_validate_mappings(mappings, memo_content)
        assert "row_inserts" in result
        assert len(result["row_inserts"]) == 1
        assert result["row_inserts"][0]["cells"] == ["2BR", "80"]


class TestPreValidateMappingsPageSplit:
    """Tests for pre_validate_mappings page-splitting rejection logic."""

    MEMO_CONTENT = (
        "=" * 60 + "\nPAGE 1\n" + "=" * 60 + "\n"
        "Table: UnitMixTable\n"
        "1BR | 120 | $1,825\n"
        "3BR | 200 | $1,345\n"
        "\n"
        + "=" * 60 + "\nPAGE 2\n" + "=" * 60 + "\n"
        "IRR is 5.0% and units are 120.\n"
    )

    def test_matching_old_value_passes(self):
        """Entry whose old_value exists on the specified page passes."""
        mappings = {
            "table_updates": [
                {"page": 1, "old_value": "120", "new_value": "130", "source": "S1"},
            ],
            "text_updates": [],
            "row_inserts": [],
        }
        result = pre_validate_mappings(mappings, self.MEMO_CONTENT)
        assert len(result["table_updates"]) == 1
        assert len(result["rejected"]) == 0

    def test_missing_old_value_rejected(self):
        """Entry whose old_value doesn't exist on the specified page is rejected."""
        mappings = {
            "table_updates": [
                {"page": 1, "old_value": "NONEXISTENT", "new_value": "X", "source": "S1"},
            ],
            "text_updates": [],
            "row_inserts": [],
        }
        result = pre_validate_mappings(mappings, self.MEMO_CONTENT)
        assert len(result["table_updates"]) == 0
        assert len(result["rejected"]) == 1
        assert "not found" in result["rejected"][0]["reason"]

    def test_wrong_page_rejected(self):
        """Entry targeting wrong page (value exists on page 2, not page 1) is rejected."""
        mappings = {
            "table_updates": [],
            "text_updates": [
                {"page": 1, "old_text": "IRR is 5.0%", "new_text": "IRR is 6.5%",
                 "source": "S1"},
            ],
            "row_inserts": [],
        }
        result = pre_validate_mappings(mappings, self.MEMO_CONTENT)
        assert len(result["text_updates"]) == 0
        assert len(result["rejected"]) == 1

    def test_unknown_page_falls_back_to_full_memo(self):
        """Entry with page not in memo falls back to searching entire content."""
        mappings = {
            "table_updates": [
                {"page": 99, "old_value": "120", "new_value": "130", "source": "S1"},
            ],
            "text_updates": [],
            "row_inserts": [],
        }
        result = pre_validate_mappings(mappings, self.MEMO_CONTENT)
        # Page 99 not found, falls back to full memo search where "120" exists
        assert len(result["table_updates"]) == 1

    def test_row_inserts_pass_through(self):
        """row_inserts are not validated by pre_validate and pass through."""
        mappings = {
            "table_updates": [],
            "text_updates": [],
            "row_inserts": [
                {"page": 1, "table_name": "T", "insert_after_row_label": "1BR",
                 "cells": ["2BR", "80"], "source": "test"},
            ],
        }
        result = pre_validate_mappings(mappings, self.MEMO_CONTENT)
        assert len(result["row_inserts"]) == 1
