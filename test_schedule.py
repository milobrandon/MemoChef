"""
Test harness for schedule (.mpp) integration.

Creates a realistic test memo modeled after an IC memo (Development Schedule
table, Master Schedule table, narrative text with embedded dates), extracts
the real proforma + schedule data, and runs through get_metric_mappings to
verify schedule-to-memo date mapping.

Usage:
    python test_schedule.py
    python test_schedule.py --randomize
    python test_schedule.py --schedule "path/to/other.mpp"
"""
import argparse
import json
import os
import random
import re
import sys
import tempfile
from datetime import datetime, timedelta

# Add the app directory to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "memo_automator_app"))

import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt

from memo_automator import (
    extract_memo_content,
    extract_proforma_data,
    extract_schedule_data,
    get_metric_mappings,
    load_config,
)

# ============================================================================
# DEFAULT PATHS (Limestone project)
# ============================================================================
DEFAULT_SCHEDULE = r"C:\Users\BrandonZmuda\Desktop\Lexington\c. Limestone\Internal Schedule_Lexington_Limestone_20251022 (working).mpp"
DEFAULT_PROFORMA = r"C:\Users\BrandonZmuda\Desktop\Claude\g. Memo Automator\v2\a. Sandbox\Proforma_Lexington-Limestone_20241021.xlsm"


# ============================================================================
# SCHEDULE PARSING
# ============================================================================
def parse_schedule_text(schedule_text: str) -> list[dict]:
    """Parse the extracted schedule text into structured task dicts."""
    tasks = []
    for line in schedule_text.split("\n"):
        match = re.match(
            r'\s*\[L(\d+)\]\s+(.+?)(\s+\[MILESTONE\])?\s+\|\s+'
            r'Start:\s+(\S+)\s+\|\s+Finish:\s+(\S+)\s+\|\s+Dur:\s+(.+)',
            line,
        )
        if match:
            tasks.append({
                "level": int(match.group(1)),
                "name": match.group(2).strip(),
                "milestone": bool(match.group(3)),
                "start": match.group(4),
                "finish": match.group(5),
                "duration": match.group(6).strip(),
            })
    return tasks


def randomize_schedule_text(schedule_text: str, offset_days: int) -> str:
    """Shift all dates in schedule text by offset_days."""
    def shift_date(match):
        date_str = match.group(0)
        try:
            dt = datetime.strptime(date_str, "%Y-%m-%d")
            shifted = dt + timedelta(days=offset_days)
            return shifted.strftime("%Y-%m-%d")
        except ValueError:
            return date_str

    return re.sub(r'\d{4}-\d{2}-\d{2}', shift_date, schedule_text)


def date_to_quarter(date_str: str) -> str:
    """Convert YYYY-MM-DD to Q# YYYY format."""
    try:
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        q = (dt.month - 1) // 3 + 1
        return f"Q{q} {dt.year}"
    except ValueError:
        return date_str


def date_to_slash(date_str: str) -> str:
    """Convert YYYY-MM-DD to M/D/YYYY format."""
    try:
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        return f"{dt.month}/{dt.day}/{dt.year}"
    except ValueError:
        return date_str


# ============================================================================
# TEST MEMO CREATION
# ============================================================================
def create_test_memo(path: str, tasks: list[dict]) -> str:
    """
    Create a realistic IC-style test memo with:
    - Page 1: Investment Snapshot with Development Schedule summary table
    - Page 2: Project Updates narrative with embedded dates
    - Page 3: Master Schedule detail table
    - Page 4: Construction detail narrative

    Dates are deliberately WRONG (offset) so the mapper must correct them.
    """
    prs = Presentation()

    # Helper: find tasks by keyword
    def find_task(keyword):
        for t in tasks:
            if keyword.lower() in t["name"].lower():
                return t
        return None

    # Collect key milestones from schedule for "wrong" dates
    # We offset dates backward by ~6 months to create deliberately wrong values
    def wrong_quarter(date_str):
        """Return a quarter that's ~6 months before the real date."""
        try:
            dt = datetime.strptime(date_str, "%Y-%m-%d")
            wrong = dt - timedelta(days=180)
            q = (wrong.month - 1) // 3 + 1
            return f"Q{q} {wrong.year}"
        except ValueError:
            return "Q2 2025"

    def wrong_slash(date_str):
        """Return a slash date ~6 months before the real date."""
        try:
            dt = datetime.strptime(date_str, "%Y-%m-%d")
            wrong = dt - timedelta(days=180)
            return f"{wrong.month}/{wrong.day}/{wrong.year}"
        except ValueError:
            return "1/1/2025"

    # Key milestones
    land_closing = find_task("land closing") or find_task("closing")
    loan_closing = find_task("loan closing")
    construction_l1 = find_task("construction")
    abatement = find_task("abatement")
    demolition = find_task("demolition")
    foundations = find_task("foundation")
    substantial = find_task("substantial completion")
    co_task = find_task("co")
    move_in = find_task("student move-in") or find_task("move-in") or find_task("move in")
    permit_task = find_task("building permit")
    entitlement = find_task("planning commission") or find_task("entitlement") or find_task("zoning")

    # ---------------------------------------------------------------
    # PAGE 1: Investment Snapshot + Development Schedule table
    # ---------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    title1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
    title1.text_frame.text = "Investment Snapshot"

    # Development Schedule table (mirroring real memo structure)
    sched_rows = [
        ("Development Schedule", ""),
        ("Final Entitlement Approval", wrong_quarter(entitlement["finish"]) if entitlement else "Q1 2025"),
        ("Land Closing", wrong_quarter(land_closing["finish"]) if land_closing else "Q2 2025"),
        ("Construction Start", wrong_quarter(construction_l1["start"]) if construction_l1 else "Q2 2025"),
        ("Certificate of Occupancy", wrong_quarter(co_task["finish"]) if co_task else "Q3 2027"),
        ("Stabilization", wrong_quarter(
            (find_task("stabilization") or co_task or {"finish": "2028-06-01"})["finish"]
        )),
    ]
    n = len(sched_rows)
    ts = slide1.shapes.add_table(n, 2, Inches(0.5), Inches(1.0), Inches(4), Inches(0.35 * n))
    ts.name = "DevelopmentSchedule"
    tbl = ts.table
    for i, (label, val) in enumerate(sched_rows):
        tbl.cell(i, 0).text = label
        tbl.cell(i, 1).text = val

    # ---------------------------------------------------------------
    # PAGE 2: Project Updates — narrative with embedded dates
    # ---------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
    title2.text_frame.text = "Project Updates"

    # Build narrative with deliberately wrong dates
    lc_q = wrong_quarter(land_closing["start"]) if land_closing else "Q2 2025"
    constr_q = wrong_quarter(construction_l1["start"]) if construction_l1 else "Q2 2025"
    co_q = wrong_quarter(co_task["finish"]) if co_task else "Q3 2027"
    mi_q = wrong_quarter(move_in["start"]) if move_in else "Q3 2027"

    ent_slash = wrong_slash(entitlement["finish"]) if entitlement else "3/15/2025"
    permit_slash = wrong_slash(permit_task["finish"]) if permit_task else "6/1/2025"

    narrative_lines = [
        f"Land closing is anticipated in {lc_q}. ",
        f"The Planning Commission is expected to grant approval by {ent_slash}. ",
        f"The building permit is expected to be issued by {permit_slash}. ",
        f"Construction is expected to begin in {constr_q} with abatement and demolition "
        f"activities commencing shortly after closing. ",
        f"Vertical construction is expected to commence in {constr_q} ",
        f"and substantial completion is targeted for {co_q}. ",
        f"Student move-in is anticipated for {mi_q}.",
    ]

    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = "".join(narrative_lines)

    # ---------------------------------------------------------------
    # PAGE 3: Master Schedule table (detailed)
    # ---------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.4))
    title3.text_frame.text = "Master Schedule"

    # Use key L1 and L2 tasks for the detail table
    detail_tasks = []
    for t in tasks:
        if t["level"] in (1, 2) and t["name"].strip():
            detail_tasks.append(t)
    detail_tasks = detail_tasks[:20]  # cap at 20 rows

    n_detail = len(detail_tasks) + 1
    ts3 = slide3.shapes.add_table(n_detail, 4, Inches(0.3), Inches(0.5),
                                  Inches(9), Inches(0.28 * n_detail))
    ts3.name = "MasterSchedule"
    tbl3 = ts3.table
    tbl3.cell(0, 0).text = "Task"
    tbl3.cell(0, 1).text = "Duration"
    tbl3.cell(0, 2).text = "Start"
    tbl3.cell(0, 3).text = "Finish"

    for i, t in enumerate(detail_tasks, 1):
        indent = "  " if t["level"] >= 2 else ""
        tbl3.cell(i, 0).text = f"{indent}{t['name']}"
        tbl3.cell(i, 1).text = t["duration"]
        # Wrong dates
        tbl3.cell(i, 2).text = wrong_slash(t["start"])
        tbl3.cell(i, 3).text = wrong_slash(t["finish"])

    # ---------------------------------------------------------------
    # PAGE 4: Construction Phases narrative
    # ---------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
    title4.text_frame.text = "Construction Overview"

    ab_start = wrong_slash(abatement["start"]) if abatement else "5/1/2025"
    ab_finish = wrong_slash(abatement["finish"]) if abatement else "6/1/2025"
    demo_finish = wrong_slash(demolition["finish"]) if demolition else "7/1/2025"
    found_finish = wrong_slash(foundations["finish"]) if foundations else "12/1/2025"
    sc_finish = wrong_slash(substantial["finish"]) if substantial else "4/1/2028"
    co_date = wrong_slash(co_task["finish"]) if co_task else "5/1/2028"

    constr_narrative = (
        f"Abatement is expected to begin on {ab_start} and be completed by {ab_finish}. "
        f"Demolition will follow, with completion targeted for {demo_finish}. "
        f"Foundation work is expected to be completed by {found_finish}. "
        f"Substantial completion is anticipated on {sc_finish} "
        f"with Certificate of Occupancy expected on {co_date}."
    )

    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(8.5), Inches(5))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    tf4.text = constr_narrative

    prs.save(path)
    print(f"  Test memo created: {path}")
    print(f"    {len(prs.slides)} slides, {n_detail} detail tasks, "
          f"{len(sched_rows)} schedule summary rows")
    return path


# ============================================================================
# MAIN TEST
# ============================================================================
def run_test(schedule_path: str, proforma_path: str, randomize: bool = False):
    """Run the full schedule integration test."""
    config_path = os.path.join(os.path.dirname(__file__),
                               "memo_automator_app", "config.yaml")
    cfg = load_config(config_path)

    # Get API key — check env, then .env files
    from dotenv import load_dotenv
    env_path = os.path.join(os.path.dirname(__file__), "memo_automator_app", ".env")
    if os.path.exists(env_path):
        load_dotenv(env_path)
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("ERROR: ANTHROPIC_API_KEY not found")
        sys.exit(1)

    client = anthropic.Anthropic(api_key=api_key, max_retries=3, timeout=600.0)

    print("=" * 70)
    print("SCHEDULE INTEGRATION TEST")
    print("=" * 70)
    print(f"Schedule:  {schedule_path}")
    print(f"Proforma:  {proforma_path}")
    print(f"Randomize: {randomize}")
    print()

    # ----- Step 1: Extract schedule -----
    print("Step 1: Extracting schedule data...")
    schedule_text = extract_schedule_data(schedule_path, cfg)
    parsed_tasks = parse_schedule_text(schedule_text)
    print(f"  {len(parsed_tasks)} tasks, {len(schedule_text)} chars")

    if randomize:
        offset = random.randint(30, 120)
        print(f"  Randomizing dates: +{offset} days")
        schedule_text = randomize_schedule_text(schedule_text, offset)
        parsed_tasks = parse_schedule_text(schedule_text)

    # ----- Step 2: Extract proforma -----
    print("\nStep 2: Extracting proforma data...")
    proforma_data = extract_proforma_data(proforma_path, cfg)
    print(f"  {len(proforma_data)} chars")

    # ----- Step 3: Combine -----
    combined_data = proforma_data + "\n\n" + schedule_text
    print(f"\nCombined data: {len(combined_data)} chars")

    # ----- Step 4: Create test memo -----
    print("\nStep 3: Creating test memo...")
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        memo_path = f.name
    create_test_memo(memo_path, parsed_tasks)

    # ----- Step 5: Extract memo content -----
    print("\nStep 4: Extracting memo content...")
    memo_content = extract_memo_content(memo_path, cfg)
    print(f"  {len(memo_content)} chars")

    # ----- Step 6: Run mapping -----
    print(f"\nStep 5: Running metric mappings...")
    print(f"  Combined payload: {len(combined_data) + len(memo_content)} chars")
    try:
        mappings = get_metric_mappings(client, combined_data, memo_content, cfg)
        mappings.pop("_truncated", None)
    except Exception as e:
        print(f"\n  ERROR: Mapping failed: {e}")
        import traceback
        traceback.print_exc()
        os.unlink(memo_path)
        return

    # ----- Step 7: Report -----
    n_table = len(mappings.get("table_updates", []))
    n_text = len(mappings.get("text_updates", []))
    n_row = len(mappings.get("row_inserts", []))

    print(f"\n{'='*70}")
    print(f"RESULTS")
    print(f"{'='*70}")
    print(f"  Table updates: {n_table}")
    print(f"  Text updates:  {n_text}")
    print(f"  Row inserts:   {n_row}")
    print(f"  Total:         {n_table + n_text + n_row}")

    # Separate schedule-sourced vs proforma-sourced
    schedule_updates = []
    proforma_updates = []

    for u in mappings.get("table_updates", []):
        src = u.get("source", "")
        if "schedule" in src.lower():
            schedule_updates.append(("TABLE", u))
        else:
            proforma_updates.append(("TABLE", u))
    for u in mappings.get("text_updates", []):
        src = u.get("source", "")
        if "schedule" in src.lower():
            schedule_updates.append(("TEXT", u))
        else:
            proforma_updates.append(("TEXT", u))

    print(f"\n  Schedule-sourced updates: {len(schedule_updates)}")
    print(f"  Proforma-sourced updates: {len(proforma_updates)}")

    if schedule_updates:
        print(f"\n{'='*70}")
        print("SCHEDULE UPDATES (detail)")
        print(f"{'='*70}")
        for kind, u in schedule_updates:
            if kind == "TABLE":
                print(f"  [TABLE] p{u.get('page')} | {u.get('table_name','?')} | "
                      f"row={u.get('row_label','?')} col={u.get('column_index','?')} | "
                      f"'{u.get('old_value','')}' -> '{u.get('new_value','')}' | "
                      f"src: {u.get('source','')}")
            else:
                old = u.get("old_text", "")
                new = u.get("new_text", "")
                # Truncate for display
                if len(old) > 80:
                    old = old[:80] + "..."
                if len(new) > 80:
                    new = new[:80] + "..."
                print(f"  [TEXT]  p{u.get('page')} | '{old}' -> '{new}' | "
                      f"src: {u.get('source','')}")

    # ----- Coverage analysis -----
    print(f"\n{'='*70}")
    print("COVERAGE ANALYSIS")
    print(f"{'='*70}")

    # Key milestones we expect to see mapped (keyword -> alternatives)
    expected_milestones = {
        "closing": ["closing", "land closing", "loan closing"],
        "entitlement/approval": ["entitlement", "approval", "development plan", "zoning", "planning commission"],
        "permit": ["permit"],
        "construction start": ["construction start", "construction s", "construction -", "abatement start"],
        "abatement": ["abatement"],
        "demolition": ["demolition", "demo"],
        "foundation": ["foundation"],
        "substantial completion": ["substantial completion", "substantial"],
        "CO/certificate": ["co ", "co\n", "certificate", "occupancy"],
        "move-in/delivery": ["move-in", "move in", "delivery"],
    }

    mapped_sources = set()
    for _, u in schedule_updates:
        mapped_sources.add(u.get("source", "").lower())
    # Also check new_value text for date references
    all_update_text = " ".join(mapped_sources)
    for _, u in schedule_updates:
        all_update_text += " " + u.get("new_value", "").lower()
        all_update_text += " " + u.get("new_text", "").lower()

    for milestone, keywords in expected_milestones.items():
        found = any(kw in all_update_text for kw in keywords)
        marker = "OK" if found else "MISS"
        print(f"  [{marker}] {milestone}")

    print(f"\n  Unique schedule sources: {len(mapped_sources)}")
    for src in sorted(mapped_sources):
        print(f"    - {src}")

    # ----- Save raw mappings for inspection -----
    mappings_path = memo_path.replace(".pptx", "_mappings.json")
    with open(mappings_path, "w") as f:
        json.dump(mappings, f, indent=2)
    print(f"\n  Full mappings saved to: {mappings_path}")
    print(f"  Test memo at: {memo_path}")

    print(f"\n{'='*70}")
    print("TEST COMPLETE")
    print(f"{'='*70}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Test schedule integration")
    parser.add_argument("--schedule", default=DEFAULT_SCHEDULE,
                        help="Path to .mpp schedule file")
    parser.add_argument("--proforma", default=DEFAULT_PROFORMA,
                        help="Path to proforma .xlsx/.xlsm file")
    parser.add_argument("--randomize", action="store_true",
                        help="Shift schedule dates by random offset")
    args = parser.parse_args()

    for p, label in [(args.schedule, "Schedule"), (args.proforma, "Proforma")]:
        if not os.path.exists(p):
            print(f"ERROR: {label} file not found: {p}")
            sys.exit(1)

    run_test(args.schedule, args.proforma, args.randomize)
