"""Unit tests for update application, branding, and layout normalization."""
import os

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from memo_automator import _replace_in_para, apply_branding, apply_updates, normalize_layout


def _find_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    raise AssertionError("No table found on slide")


def test_apply_updates_table_and_text(sample_pptx):
    mappings = {
        "table_updates": [
            {
                "page": 1,
                "table_name": "UnitMixTable",
                "row_label": "1BR",
                "column_index": 1,
                "old_value": "120",
                "new_value": "130",
                "source": "Executive Summary B3",
            }
        ],
        "text_updates": [
            {
                "page": 1,
                "old_text": "IRR is 5.0%",
                "new_text": "IRR is 6.5%",
                "source": "Executive Summary B2",
            }
        ],
        "row_inserts": [],
    }
    changes = apply_updates(sample_pptx, mappings, dry_run=False)
    assert len(changes) == 2

    prs = Presentation(sample_pptx)
    table = _find_table(prs.slides[0])
    assert table.cell(1, 1).text == "130"
    narrative = next(s for s in prs.slides[0].shapes if s.name == "NarrativeBox")
    assert "IRR is 6.5%" in narrative.text_frame.text


def test_apply_branding(sample_pptx):
    theme_path = os.path.join(os.path.dirname(__file__), "Subtext Brand Theme.thmx")
    if not os.path.exists(theme_path):
        pytest.skip("Theme file not found for branding test")

    cfg = {
        "branding": {
            "heading_size_threshold": 18,
            "color_distance_threshold": 80,
        }
    }
    runs = apply_branding(sample_pptx, theme_path, cfg)
    assert runs > 0

    prs = Presentation(sample_pptx)

    # Subheader (bold run in non-title shape) should get Pragmatica Bold
    sub_box = next(s for s in prs.slides[0].shapes if s.name == "SubheaderBox")
    sub_runs = sub_box.text_frame.paragraphs[0].runs
    assert sub_runs[0].font.name == "Pragmatica Bold"  # bold subheader
    assert sub_runs[1].font.name == "Pragmatica Book"  # non-bold body


def test_apply_branding_preserves_table_fonts(sample_pptx):
    """apply_branding() must not change font family or size inside table cells."""
    theme_path = os.path.join(os.path.dirname(__file__), "Subtext Brand Theme.thmx")
    if not os.path.exists(theme_path):
        pytest.skip("Theme file not found for branding test")

    # Set distinctive fonts on table cells before branding
    from pptx.util import Pt
    prs_before = Presentation(sample_pptx)
    table = _find_table(prs_before.slides[0])

    header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    header_run.font.name = "Arial"
    header_run.font.size = Pt(14)
    header_run.font.bold = True

    body_run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
    body_run.font.name = "Calibri"
    body_run.font.size = Pt(10)

    prs_before.save(sample_pptx)

    cfg = {
        "branding": {
            "heading_size_threshold": 18,
            "color_distance_threshold": 80,
        }
    }
    apply_branding(sample_pptx, theme_path, cfg)

    prs_after = Presentation(sample_pptx)
    table_after = _find_table(prs_after.slides[0])

    header_run_after = table_after.cell(0, 0).text_frame.paragraphs[0].runs[0]
    body_run_after = table_after.cell(1, 0).text_frame.paragraphs[0].runs[0]

    assert header_run_after.font.name == "Arial", (
        f"Header cell font should remain Arial, got {header_run_after.font.name!r}"
    )
    assert header_run_after.font.size == Pt(14), (
        f"Header cell size should remain 14pt, got {header_run_after.font.size!r}"
    )
    assert body_run_after.font.name == "Calibri", (
        f"Body cell font should remain Calibri, got {body_run_after.font.name!r}"
    )
    assert body_run_after.font.size == Pt(10), (
        f"Body cell size should remain 10pt, got {body_run_after.font.size!r}"
    )


def test_normalize_layout(layout_test_pptx):
    cfg = {
        "layout": {
            "margin_left": 0.50,
            "margin_right": 0.50,
            "margin_top": 0.25,
            "margin_bottom": 0.50,
            "snap_tolerance": 0.05,
        }
    }
    summary = normalize_layout(layout_test_pptx, cfg)
    assert "shapes_clamped_to_margins" in summary
    assert summary["shapes_clamped_to_margins"] >= 1

    prs = Presentation(layout_test_pptx)
    slide = prs.slides[2]
    off_margin = next(s for s in slide.shapes if s.name == "OffMargin")
    assert off_margin.left >= Inches(0.50)
    assert off_margin.top >= Inches(0.25)


def test_normalize_layout_preserves_table_font_sizes(tmp_dir):
    """normalize_layout() must not overwrite intentional per-cell table font sizes.

    Simulates a dense comp table (like Knoxville pages 13/16) where header
    cells are 12pt and body cells are intentionally smaller at 8pt.  The old
    normalization code used a global dominant-size pass that treated the 8pt
    body cells as the dominant and then normalised any body cell outside
    ±0.5–4pt back to 8pt — which meant a mixed table with some 6pt rows would
    have those rows overwritten.  We construct exactly that scenario here.
    """
    path = os.path.join(tmp_dir, "font_size_test.pptx")
    prs = Presentation()

    # Slide 0: cover (skipped by normalize_layout)
    prs.slides.add_slide(prs.slide_layouts[5])

    # Slide 1: content slide with a 6x2 table:
    #   row 0 (header): 12pt
    #   rows 1-3 (body majority): 8pt  -> dominant_table_size = 8.0
    #   rows 4-5 (body minority): 6pt  -> diff = 2.0, within old 0.5<diff<4.0 window
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_shape = slide.shapes.add_table(
        6, 2, Inches(1.0), Inches(0.5), Inches(6.0), Inches(4.0)
    )
    table = table_shape.table

    # Header row: 12pt
    for col in range(2):
        cell = table.cell(0, col)
        cell.text = f"Header {col}"
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(12)

    # Majority body rows: 8pt
    for row in range(1, 4):
        for col in range(2):
            cell = table.cell(row, col)
            cell.text = f"Cell {row},{col}"
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(8)

    # Minority body rows: intentionally smaller at 6pt (dense comp rows)
    for row in range(4, 6):
        for col in range(2):
            cell = table.cell(row, col)
            cell.text = f"Small {row},{col}"
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(6)

    prs.save(path)

    cfg = {
        "layout": {
            "margin_left": 0.50,
            "margin_right": 0.50,
            "margin_top": 0.25,
            "margin_bottom": 0.50,
            "snap_tolerance": 0.05,
        }
    }
    summary = normalize_layout(path, cfg)

    assert summary["table_font_size_normalized"] == 0, (
        f"Expected 0 font size normalizations, got {summary['table_font_size_normalized']}"
    )

    prs_after = Presentation(path)
    after_table = next(
        s.table for s in prs_after.slides[1].shapes if s.has_table
    )
    # The 6pt minority rows must remain 6pt, not be overwritten to 8pt
    small_run = after_table.cell(4, 0).text_frame.paragraphs[0].runs[0]
    assert small_run.font.size == Pt(6), (
        f"Dense comp body cell font size should remain 6pt, got {small_run.font.size!r}"
    )


class TestReplaceInParaCrossRun:
    """Tests for _replace_in_para cross-run replacement (Pass 2)."""

    @staticmethod
    def _make_para_with_runs(texts: list[str]):
        """Create a paragraph with multiple runs containing the given texts."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txbox.text_frame
        para = tf.paragraphs[0]
        # Clear default run and add our runs
        para.clear()
        for t in texts:
            run = para.add_run()
            run.text = t
        return para

    def test_cross_run_replacement_succeeds(self):
        """Value split across two runs is replaced correctly."""
        para = self._make_para_with_runs(["$1,8", "25"])
        result = _replace_in_para(para, "$1,825", "$2,100")
        assert result is True
        full_text = "".join(r.text for r in para.runs)
        assert "$2,100" in full_text

    def test_cross_run_preserves_surrounding_text(self):
        """Text before and after the replaced value is preserved."""
        para = self._make_para_with_runs(["Rent: $1,8", "25 per month"])
        result = _replace_in_para(para, "$1,825", "$2,100")
        assert result is True
        full_text = "".join(r.text for r in para.runs)
        assert full_text == "Rent: $2,100 per month"

    def test_cross_run_returns_false_on_no_match(self):
        """Returns False when old_text is not found across any runs."""
        para = self._make_para_with_runs(["Rent: $1,8", "25 per month"])
        result = _replace_in_para(para, "$9,999", "$0")
        assert result is False
        full_text = "".join(r.text for r in para.runs)
        assert full_text == "Rent: $1,825 per month"

    def test_single_run_replacement_preferred(self):
        """When value fits in one run, Pass 1 handles it (no cross-run needed)."""
        para = self._make_para_with_runs(["$1,825", " per month"])
        result = _replace_in_para(para, "$1,825", "$2,100")
        assert result is True
        assert para.runs[0].text == "$2,100"
        assert para.runs[1].text == " per month"
