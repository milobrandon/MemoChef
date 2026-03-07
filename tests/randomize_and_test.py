#!/usr/bin/env python3
"""Randomize specific values in a memo PPTX, run the pipeline, and compare results.

Creates a mutated copy of the memo with known changes, then runs memo_automator
in dry-run mode to see if the pipeline correctly identifies and fixes them.
"""

from __future__ import annotations

import copy
import json
import os
import random
import re
import shutil
import sys
import time
from pathlib import Path

from pptx import Presentation

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

MEMO_PATH = Path("a. Sandbox/EVER Lexington_20260205_TEST.pptx")
PROFORMA_PATH = Path("a. Sandbox/Proforma_Lexington-Limestone_20241021.xlsm")
CONFIG_PATH = Path("config.yaml")
OUTPUT_DIR = Path("test_output_randomized")

# Values to mutate: (slide_index, shape_description, original_text_fragment, field_name)
# We'll find these in the memo and randomize the numeric portions.
MUTATIONS = [
    # Slide 3 - Executive Summary table
    {"search": "787 Beds", "field": "Bed Count", "mutate": "numeric"},
    {"search": "274 Units", "field": "Unit Count", "mutate": "numeric"},
    {"search": "$21,213,979", "field": "Acquisition Costs Total", "mutate": "dollar"},
    {"search": "$103,387,549", "field": "Hard Costs Total", "mutate": "dollar"},
    {"search": "$153,119,446", "field": "Total Costs", "mutate": "dollar"},
    {"search": "$194,561", "field": "Total Cost Per Bed", "mutate": "dollar"},
    {"search": "$11,865,527", "field": "Year 2 NOI", "mutate": "dollar"},
    {"search": "7.75%", "field": "Year 2 Return on Cost", "mutate": "percent"},
    {"search": "6.70%", "field": "Untrended Return on Cost", "mutate": "percent"},
    # Slide 5 - University profile
    {"search": "38,383", "field": "Total Enrollment", "mutate": "numeric"},
    {"search": "$34,140", "field": "Tuition", "mutate": "dollar"},
    {"search": "$16,016", "field": "Room & Board", "mutate": "dollar"},
]


def randomize_dollar(original: str) -> str:
    """Randomize a dollar amount like $21,213,979."""
    num_str = original.replace("$", "").replace(",", "")
    try:
        num = float(num_str)
    except ValueError:
        return original
    # Vary by 5-25%
    factor = random.uniform(0.75, 1.25)
    new_num = int(num * factor)
    return f"${new_num:,}"


def randomize_numeric(original: str) -> str:
    """Randomize a number like '787 Beds' or '38,383'."""
    match = re.search(r"[\d,]+", original)
    if not match:
        return original
    num_str = match.group().replace(",", "")
    try:
        num = int(num_str)
    except ValueError:
        return original
    factor = random.uniform(0.75, 1.25)
    new_num = int(num * factor)
    new_str = f"{new_num:,}" if num >= 1000 else str(new_num)
    return original[: match.start()] + new_str + original[match.end() :]


def randomize_percent(original: str) -> str:
    """Randomize a percentage like '7.75%'."""
    match = re.search(r"[\d.]+", original)
    if not match:
        return original
    try:
        num = float(match.group())
    except ValueError:
        return original
    delta = random.uniform(-2.0, 2.0)
    new_num = max(0.1, num + delta)
    return f"{new_num:.2f}%"


MUTATORS = {
    "dollar": randomize_dollar,
    "numeric": randomize_numeric,
    "percent": randomize_percent,
}


def apply_mutations(prs: Presentation) -> list[dict]:
    """Walk through the presentation and apply mutations. Returns list of changes made."""
    changes = []

    for mutation in MUTATIONS:
        search_text = mutation["search"]
        mutator = MUTATORS[mutation["mutate"]]
        found = False

        for slide_idx, slide in enumerate(prs.slides):
            # Search in tables
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    for row_idx in range(len(tbl.rows)):
                        for col_idx in range(len(tbl.columns)):
                            cell = tbl.cell(row_idx, col_idx)
                            cell_text = cell.text.strip()
                            if search_text in cell_text:
                                new_value = mutator(search_text)
                                # Replace in all runs to preserve formatting
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if search_text in run.text:
                                            run.text = run.text.replace(
                                                search_text, new_value
                                            )
                                changes.append(
                                    {
                                        "field": mutation["field"],
                                        "slide": slide_idx + 1,
                                        "original": search_text,
                                        "mutated_to": new_value,
                                        "location": f"table row {row_idx}, col {col_idx}",
                                    }
                                )
                                found = True
                                break
                        if found:
                            break
                if found:
                    break

                # Search in text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        full_text = para.text
                        if search_text in full_text:
                            new_value = mutator(search_text)
                            for run in para.runs:
                                if search_text in run.text:
                                    run.text = run.text.replace(
                                        search_text, new_value
                                    )
                            changes.append(
                                {
                                    "field": mutation["field"],
                                    "slide": slide_idx + 1,
                                    "original": search_text,
                                    "mutated_to": new_value,
                                    "location": "text frame",
                                }
                            )
                            found = True
                            break
                if found:
                    break
            if found:
                break

        if not found:
            print(f"  WARNING: Could not find '{search_text}' for {mutation['field']}")

    return changes


def main():
    random.seed(int(time.time()))

    print("=" * 60)
    print("MEMO CHEF - RANDOMIZED DRY RUN TEST")
    print("=" * 60)

    # Clean output dir
    if OUTPUT_DIR.exists():
        shutil.rmtree(OUTPUT_DIR)
    OUTPUT_DIR.mkdir(parents=True)

    # Load and mutate the memo
    print(f"\nLoading memo: {MEMO_PATH}")
    prs = Presentation(str(MEMO_PATH))

    print("Applying random mutations...")
    changes = apply_mutations(prs)
    print(f"\n{len(changes)} mutations applied:")
    for c in changes:
        print(f"  [{c['field']}] {c['original']} -> {c['mutated_to']}  (slide {c['slide']}, {c['location']})")

    # Save mutated memo
    mutated_path = OUTPUT_DIR / "mutated_memo.pptx"
    prs.save(str(mutated_path))
    print(f"\nSaved mutated memo to: {mutated_path}")

    # Save mutation log
    mutation_log_path = OUTPUT_DIR / "mutations.json"
    with open(mutation_log_path, "w") as f:
        json.dump(changes, f, indent=2)

    # Run the pipeline
    print("\n" + "=" * 60)
    print("RUNNING PIPELINE (dry run)...")
    print("=" * 60)

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        # Try reading from streamlit secrets
        try:
            import toml
            secrets = toml.load(".streamlit/secrets.toml")
            api_key = secrets.get("ANTHROPIC_API_KEY")
        except Exception:
            pass
    if not api_key:
        print("ERROR: No ANTHROPIC_API_KEY found")
        sys.exit(1)

    start_time = time.time()

    # Run via CLI subprocess
    import subprocess

    cmd = [
        sys.executable, "memo_automator.py",
        str(mutated_path),
        str(PROFORMA_PATH),
        "--output-dir", str(OUTPUT_DIR),
        "--config", str(CONFIG_PATH),
        "--dry-run",
    ]
    env = os.environ.copy()
    env["ANTHROPIC_API_KEY"] = api_key

    proc = subprocess.run(cmd, env=env, capture_output=False, timeout=1200)

    elapsed = time.time() - start_time

    print(f"\nPipeline completed in {elapsed:.0f}s (exit code {proc.returncode})")
    result = proc.returncode == 0

    # Analyze results
    print("\n" + "=" * 60)
    print("RESULTS ANALYSIS")
    print("=" * 60)

    if not result:
        print("WARNING: Pipeline exited with errors, checking partial results anyway...")

    # Try change log first, fall back to validated mappings JSON
    change_log = ""
    change_log_files = list(OUTPUT_DIR.glob("CHANGE_LOG*"))
    if change_log_files:
        change_log = change_log_files[0].read_text(encoding="utf-8")

    # Also load validated mappings for deeper analysis
    mappings_text = ""
    validated_path = OUTPUT_DIR / "mappings_validated.json"
    if validated_path.exists():
        mappings_text = validated_path.read_text(encoding="utf-8")

    search_text = change_log + "\n" + mappings_text
    if not search_text.strip():
        print("WARNING: No change log or validated mappings found")
        return

    # Check which mutations were caught
    caught = 0
    missed = 0
    for mutation in changes:
        original = mutation["original"]
        mutated = mutation["mutated_to"]
        field = mutation["field"]
        # The pipeline should have detected the mutated value differs from proforma
        # and proposed changing it back to the original (or proforma value).
        # Check if the mutated value appears as old_value in mappings
        if mutated in search_text or original in search_text:
            caught += 1
            print(f"  CAUGHT: {field} ({original} -> {mutated})")
        else:
            missed += 1
            print(f"  MISSED: {field} ({original} -> {mutated})")

    print(f"\n--- Summary ---")
    print(f"Mutations applied: {len(changes)}")
    print(f"Caught by pipeline: {caught}")
    print(f"Missed: {missed}")
    print(f"Accuracy: {caught / len(changes) * 100:.0f}%")
    print(f"Time: {elapsed:.0f}s")


if __name__ == "__main__":
    main()
