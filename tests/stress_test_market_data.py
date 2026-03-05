#!/usr/bin/env python3
"""
Stress test for Memo Chef market data integration.
Exercises every new code path, finds bugs, and reports them.
"""

import os
import sys
import tempfile
import traceback

# Add project root to path
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

from memo_automator import (
    extract_market_data,
    extract_memo_content,
    extract_proforma_data,
    load_config,
    apply_updates,
    _replace_in_para,
    _reformat_run,
    _apply_chart_updates,
    _MARKET_DASHBOARD_TABS,
)

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
MARKET_DATA_FILE = r"C:/Users/BrandonZmuda/Desktop/Claude/g. Memo Automator/a. Sandbox/New Template Test.xlsx"
MEMO_TEMPLATE = r"C:/Users/BrandonZmuda/Desktop/Claude/g. Memo Automator/v2/b. Sandbox v2/IC Approval Memo_TEMPLATE.pptx"
MEMO_LIMESTONE = r"C:/Users/BrandonZmuda/Desktop/Claude/g. Memo Automator/v2/a. Sandbox/TEST_Modified_Memo_Limestone.pptx"
MEMO_LEXINGTON = r"C:/Users/BrandonZmuda/Desktop/Claude/g. Memo Automator/v2/a. Sandbox/EVER Lexington_20260205_TEST.pptx"
PROFORMA_FILE = r"C:/Users/BrandonZmuda/Desktop/Claude/g. Memo Automator/v2/a. Sandbox/Proforma_Lexington-Limestone_20241021.xlsm"
CONFIG_FILE = os.path.join(PROJECT_ROOT, "config.yaml")

# Load config once
CFG = load_config(CONFIG_FILE)

PASS_COUNT = 0
FAIL_COUNT = 0


def report(test_name, passed, details=""):
    global PASS_COUNT, FAIL_COUNT
    status = "PASS" if passed else "FAIL"
    if passed:
        PASS_COUNT += 1
    else:
        FAIL_COUNT += 1
    print(f"  {status}: {details}")


# ===================================================================
# TEST 1: extract_market_data() with real file
# ===================================================================
def test_1():
    print("\n=== Test 1: extract_market_data() with real file ===")
    try:
        result = extract_market_data(MARKET_DATA_FILE, CFG)
        char_count = len(result)
        print(f"  Char count: {char_count}")

        # Count tabs found
        tab_count = 0
        found_tabs = []
        for tab in _MARKET_DASHBOARD_TABS:
            if f"TAB: {tab}" in result:
                tab_count += 1
                found_tabs.append(tab)
        print(f"  Tab count found: {tab_count} / {len(_MARKET_DASHBOARD_TABS)}")
        print(f"  Found tabs: {found_tabs}")
        missing = [t for t in _MARKET_DASHBOARD_TABS if t not in found_tabs]
        if missing:
            print(f"  Missing tabs: {missing}")

        report("Extraction succeeded", char_count > 0, f"{char_count} chars extracted")
        report("All 6 tabs present", tab_count == 6,
               f"{tab_count}/6 tabs found" + (f", missing: {missing}" if missing else ""))

        # Check for back-end tab leakage
        backend_tabs = ["PROPERTIES", "IPEDS", "SCHOOLS", "PBH DATA", "EMPLOYMENT",
                        "DEMAND PROJ", "PIPELINE", "MASTER"]
        leaked = [t for t in backend_tabs if f"TAB: {t}" in result]
        report("No back-end tab leakage", len(leaked) == 0,
               f"Leaked tabs: {leaked}" if leaked else "No backend tabs leaked")

    except Exception as e:
        report("Test 1 execution", False, f"EXCEPTION: {e}")
        traceback.print_exc()


# ===================================================================
# TEST 2: extract_memo_content() chart extraction
# ===================================================================
def test_2():
    print("\n=== Test 2: extract_memo_content() chart extraction ===")
    memos = {
        "IC Approval Memo TEMPLATE": MEMO_TEMPLATE,
        "Limestone": MEMO_LIMESTONE,
        "EVER Lexington": MEMO_LEXINGTON,
    }
    for name, path in memos.items():
        print(f"\n  --- Memo: {name} ---")
        try:
            result = extract_memo_content(path, CFG)
            # Count shapes
            total_shapes = result.count("--- Shape ")
            table_shapes = result.count("Table: ")
            chart_shapes = result.count("Chart type: ")
            text_frames = result.count("Para ")

            print(f"    Total shapes: {total_shapes}")
            print(f"    Shapes with tables: {table_shapes}")
            print(f"    Shapes with charts: {chart_shapes}")
            print(f"    Text frames (paras): {text_frames}")

            has_charts = chart_shapes > 0
            print(f"    Has charts: {has_charts}")

            if has_charts:
                # Extract chart metadata
                import re
                chart_types = re.findall(r"Chart type: (.+)", result)
                chart_titles = re.findall(r"Chart title: '(.+?)'", result)
                series_names = re.findall(r"Series \d+ \('(.+?)'\)", result)
                print(f"    Chart types: {chart_types[:10]}")
                print(f"    Chart titles: {chart_titles[:10]}")
                print(f"    Series names: {series_names[:20]}")

            report(f"{name} extraction", True,
                   f"{total_shapes} shapes, {chart_shapes} charts, {table_shapes} tables")
        except Exception as e:
            report(f"{name} extraction", False, f"EXCEPTION: {e}")
            traceback.print_exc()


# ===================================================================
# TEST 3: _replace_in_para() formatting preservation
# ===================================================================
def test_3():
    print("\n=== Test 3: _replace_in_para() formatting preservation ===")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from copy import deepcopy

        # Create a presentation with a text box containing 3 runs
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]

        # Run 0: bold "Hello "
        run0 = p.add_run()
        run0.text = "Hello "
        run0.font.bold = True
        run0.font.size = Pt(12)

        # Run 1: normal "world "
        run1 = p.add_run()
        run1.text = "world "
        run1.font.bold = False
        run1.font.size = Pt(12)

        # Run 2: italic "today"
        run2 = p.add_run()
        run2.text = "today"
        run2.font.italic = True
        run2.font.size = Pt(12)

        # Print initial state
        print("  Initial state:")
        for i, run in enumerate(p.runs):
            print(f"    Run {i}: text='{run.text}', bold={run.font.bold}, italic={run.font.italic}")

        # Pass 1: single-run replacement (world -> earth)
        print("\n  Pass 1: Replace 'world' with 'earth' (single-run)...")
        result1 = _replace_in_para(p, "world", "earth")
        print(f"    Replacement returned: {result1}")
        for i, run in enumerate(p.runs):
            print(f"    Run {i}: text='{run.text}', bold={run.font.bold}, italic={run.font.italic}")

        # Check run1 formatting preserved
        run1_after = p.runs[1]
        run2_after = p.runs[2]
        report("Pass 1 replacement succeeded", result1 == True, f"returned {result1}")
        report("Pass 1 run1 text correct", "earth" in run1_after.text,
               f"run1.text = '{run1_after.text}'")
        report("Pass 1 run2 untouched (italic)", run2_after.font.italic == True,
               f"run2 italic={run2_after.font.italic}, text='{run2_after.text}'")
        report("Pass 1 run2 text untouched", run2_after.text == "today",
               f"run2.text = '{run2_after.text}'")

        # Pass 2: cross-run replacement (Hello earth -> Hi globe)
        print("\n  Pass 2: Replace 'Hello earth' with 'Hi globe' (cross-run)...")
        full_before = "".join(r.text for r in p.runs)
        print(f"    Full text before: '{full_before}'")
        result2 = _replace_in_para(p, "Hello earth", "Hi globe")
        print(f"    Replacement returned: {result2}")
        for i, run in enumerate(p.runs):
            print(f"    Run {i}: text='{run.text}', bold={run.font.bold}, italic={run.font.italic}")

        full_after = "".join(r.text for r in p.runs)
        print(f"    Full text after: '{full_after}'")

        report("Pass 2 replacement succeeded", result2 == True, f"returned {result2}")
        report("Pass 2 result text correct", "Hi globe" in full_after,
               f"full text = '{full_after}'")

        # Check run2 ("today", italic) is still untouched
        # It could still be at index 2 if there are 3 runs
        last_run = p.runs[-1] if p.runs else None
        if last_run:
            report("Pass 2 last run is 'today'", "today" in last_run.text,
                   f"last run text='{last_run.text}'")
            report("Pass 2 last run still italic", last_run.font.italic == True,
                   f"last run italic={last_run.font.italic}")
        else:
            report("Pass 2 runs exist", False, "No runs remain")

    except Exception as e:
        report("Test 3 execution", False, f"EXCEPTION: {e}")
        traceback.print_exc()


# ===================================================================
# TEST 4: _reformat_run() bold/italic preservation
# ===================================================================
def test_4():
    print("\n=== Test 4: _reformat_run() bold/italic preservation ===")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]

        # Run A: bold=True, italic=False
        runA = p.add_run()
        runA.text = "Bold Text"
        runA.font.bold = True
        runA.font.italic = False
        runA.font.size = Pt(14)
        runA.font.name = "Arial"

        # Run B: bold=False, italic=True
        runB = p.add_run()
        runB.text = " Italic Text"
        runB.font.bold = False
        runB.font.italic = True
        runB.font.size = Pt(14)
        runB.font.name = "Arial"

        # Run C: bold=True, italic=True
        runC = p.add_run()
        runC.text = " Bold Italic"
        runC.font.bold = True
        runC.font.italic = True
        runC.font.size = Pt(14)
        runC.font.name = "Arial"

        print("  Before _reformat_run():")
        for label, run in [("A", runA), ("B", runB), ("C", runC)]:
            print(f"    Run {label}: text='{run.text}', bold={run.font.bold}, italic={run.font.italic}, font='{run.font.name}'")

        # Apply reformat
        for run in [runA, runB, runC]:
            _reformat_run(run, is_heading_context=False, size_threshold=18,
                         heading_font="Pragmatica Bold", body_font="Pragmatica",
                         color_threshold=80)

        print("  After _reformat_run():")
        for label, run in [("A", runA), ("B", runB), ("C", runC)]:
            print(f"    Run {label}: text='{run.text}', bold={run.font.bold}, italic={run.font.italic}, font='{run.font.name}'")

        report("Run A bold preserved", runA.font.bold == True,
               f"bold={runA.font.bold}")
        report("Run A italic preserved (False)", runA.font.italic == False,
               f"italic={runA.font.italic}")
        report("Run B bold preserved (False)", runB.font.bold == False,
               f"bold={runB.font.bold}")
        report("Run B italic preserved", runB.font.italic == True,
               f"italic={runB.font.italic}")
        report("Run C bold preserved", runC.font.bold == True,
               f"bold={runC.font.bold}")
        report("Run C italic preserved", runC.font.italic == True,
               f"italic={runC.font.italic}")
        report("Font changed to Pragmatica", runA.font.name == "Pragmatica",
               f"font='{runA.font.name}'")

    except Exception as e:
        report("Test 4 execution", False, f"EXCEPTION: {e}")
        traceback.print_exc()


# ===================================================================
# TEST 5: _apply_chart_updates() with mock data
# ===================================================================
def test_5():
    print("\n=== Test 5: _apply_chart_updates() with mock data ===")
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        import shutil

        # Create a presentation with an embedded bar chart
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        chart_data = CategoryChartData()
        chart_data.categories = ['2021', '2022', '2023']
        chart_data.add_series('Revenue', (100, 200, 300))
        chart_data.add_series('Expenses', (80, 150, 250))

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(6), Inches(4),
            chart_data
        )
        chart_shape.name = "TestChart1"

        # Save to temp file
        tmp_dir = tempfile.mkdtemp()
        tmp_path = os.path.join(tmp_dir, "test_chart.pptx")
        prs.save(tmp_path)

        # Read the chart to verify it was created
        prs2 = Presentation(tmp_path)
        chart = prs2.slides[0].shapes[0].chart
        print(f"  Chart created: type={chart.chart_type}")
        print(f"  Series count: {len(chart.series)}")
        for i, s in enumerate(chart.series):
            vals = list(s.values) if s.values else []
            print(f"    Series {i}: values={vals}")

        # Craft chart_update dict
        chart_updates = [{
            "page": 1,
            "chart_name": "TestChart1",
            "chart_title": "",
            "series_name": "Revenue",
            "old_values": [100, 200, 300],
            "new_values": [110, 220, 330],
            "source": "Test",
        }]

        # Test dry_run mode
        print("\n  Testing dry_run mode...")
        dry_path = os.path.join(tmp_dir, "test_chart_dry.pptx")
        shutil.copy2(tmp_path, dry_path)
        dry_changes = _apply_chart_updates(dry_path, chart_updates, dry_run=True)
        print(f"  Dry run changes count: {len(dry_changes)}")
        # Reload and verify NO changes were saved
        prs_dry = Presentation(dry_path)
        dry_chart = prs_dry.slides[0].shapes[0].chart
        dry_vals = list(dry_chart.series[0].values)
        print(f"  Dry run series 0 values: {dry_vals}")
        report("Dry run returns changes", len(dry_changes) > 0,
               f"{len(dry_changes)} changes reported")
        report("Dry run does not save", dry_vals == [100.0, 200.0, 300.0],
               f"values={dry_vals} (expected [100, 200, 300])")

        # Test actual update
        print("\n  Testing actual chart update...")
        real_path = os.path.join(tmp_dir, "test_chart_real.pptx")
        shutil.copy2(tmp_path, real_path)
        real_changes = _apply_chart_updates(real_path, chart_updates, dry_run=False)
        print(f"  Real changes count: {len(real_changes)}")

        # Reload and verify changes applied
        prs_real = Presentation(real_path)
        real_chart = prs_real.slides[0].shapes[0].chart
        real_vals = list(real_chart.series[0].values)
        print(f"  Updated series 0 values: {real_vals}")

        # The series[1] (Expenses) should remain unchanged
        exp_vals = list(real_chart.series[1].values)
        print(f"  Expenses series values (should be unchanged): {exp_vals}")

        report("Chart update applied", len(real_changes) > 0,
               f"{len(real_changes)} changes applied")
        report("Revenue values updated", real_vals == [110.0, 220.0, 330.0],
               f"values={real_vals}")
        report("Expenses unchanged", exp_vals == [80.0, 150.0, 250.0],
               f"values={exp_vals}")

        # Cleanup
        import shutil as sh
        sh.rmtree(tmp_dir, ignore_errors=True)

    except Exception as e:
        report("Test 5 execution", False, f"EXCEPTION: {e}")
        traceback.print_exc()


# ===================================================================
# TEST 6: apply_updates() with chart_updates in mappings
# ===================================================================
def test_6():
    print("\n=== Test 6: apply_updates() with chart_updates in mappings ===")
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        import shutil

        # Create a presentation with a chart
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = CategoryChartData()
        chart_data.categories = ['Q1', 'Q2']
        chart_data.add_series('Sales', (10, 20))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1), Inches(6), Inches(4),
            chart_data
        )
        chart_shape.name = "SalesChart"

        tmp_dir = tempfile.mkdtemp()
        tmp_path = os.path.join(tmp_dir, "test_apply.pptx")
        prs.save(tmp_path)

        # Test with empty chart_updates list
        print("  Testing empty chart_updates...")
        empty_path = os.path.join(tmp_dir, "test_empty.pptx")
        shutil.copy2(tmp_path, empty_path)
        mappings_empty = {
            "table_updates": [],
            "text_updates": [],
            "row_inserts": [],
            "chart_updates": [],
        }
        try:
            changes = apply_updates(empty_path, mappings_empty, dry_run=True)
            report("Empty chart_updates no error", True, f"{len(changes)} changes")
        except Exception as e:
            report("Empty chart_updates no error", False, f"EXCEPTION: {e}")
            traceback.print_exc()

        # Test with chart_updates for a nonexistent page
        print("  Testing chart_updates for nonexistent page...")
        bad_page_path = os.path.join(tmp_dir, "test_bad_page.pptx")
        shutil.copy2(tmp_path, bad_page_path)
        mappings_bad_page = {
            "table_updates": [],
            "text_updates": [],
            "row_inserts": [],
            "chart_updates": [{
                "page": 99,
                "chart_name": "SalesChart",
                "series_name": "Sales",
                "old_values": [10, 20],
                "new_values": [30, 40],
                "source": "Test",
            }],
        }
        try:
            changes = apply_updates(bad_page_path, mappings_bad_page, dry_run=False)
            report("Nonexistent page no crash", True,
                   f"{len(changes)} changes (should be 0 for chart)")
        except Exception as e:
            report("Nonexistent page no crash", False, f"CRASH: {e}")
            traceback.print_exc()

        # Test that pipeline actually invokes _apply_chart_updates
        print("  Testing pipeline invokes _apply_chart_updates...")
        real_path = os.path.join(tmp_dir, "test_pipeline_chart.pptx")
        shutil.copy2(tmp_path, real_path)
        mappings_with_chart = {
            "table_updates": [],
            "text_updates": [],
            "row_inserts": [],
            "chart_updates": [{
                "page": 1,
                "chart_name": "SalesChart",
                "series_name": "Sales",
                "old_values": [10, 20],
                "new_values": [50, 60],
                "source": "Test",
            }],
        }
        try:
            changes = apply_updates(real_path, mappings_with_chart, dry_run=False)
            chart_changes = [c for c in changes if c.get("type") == "chart"]
            print(f"    Total changes: {len(changes)}, chart changes: {len(chart_changes)}")
            report("Pipeline invokes chart updates", len(chart_changes) > 0,
                   f"{len(chart_changes)} chart changes applied")

            # Verify the actual chart data changed
            prs_verify = Presentation(real_path)
            verify_chart = prs_verify.slides[0].shapes[0].chart
            verify_vals = list(verify_chart.series[0].values)
            print(f"    Verified values: {verify_vals}")
            report("Chart data actually changed", verify_vals == [50.0, 60.0],
                   f"values={verify_vals}")
        except Exception as e:
            report("Pipeline chart update", False, f"EXCEPTION: {e}")
            traceback.print_exc()

        # Cleanup
        import shutil as sh
        sh.rmtree(tmp_dir, ignore_errors=True)

    except Exception as e:
        report("Test 6 execution", False, f"EXCEPTION: {e}")
        traceback.print_exc()


# ===================================================================
# TEST 7: Pipeline integration simulation
# ===================================================================
def test_7():
    print("\n=== Test 7: Pipeline integration simulation ===")

    # Step 1: extract proforma
    print("  Step 1: Extracting proforma data...")
    try:
        proforma_data = extract_proforma_data(PROFORMA_FILE, CFG)
        print(f"    Proforma: {len(proforma_data)} chars")
        report("Proforma extraction", len(proforma_data) > 0,
               f"{len(proforma_data)} chars")
    except Exception as e:
        report("Proforma extraction", False, f"EXCEPTION: {e}")
        traceback.print_exc()
        proforma_data = ""

    # Step 2: extract market data
    print("  Step 2: Extracting market data...")
    try:
        market_data = extract_market_data(MARKET_DATA_FILE, CFG)
        print(f"    Market data: {len(market_data)} chars")
        report("Market data extraction", len(market_data) > 0,
               f"{len(market_data)} chars")
    except Exception as e:
        report("Market data extraction", False, f"EXCEPTION: {e}")
        traceback.print_exc()
        market_data = ""

    # Step 3: concatenate as pipeline would
    print("  Step 3: Concatenating data sources...")
    combined = proforma_data
    if market_data:
        combined += "\n\n" + market_data
    print(f"    Combined prompt data: {len(combined)} chars")

    # Step 4: extract memo content from each real memo
    for name, path in [("TEMPLATE", MEMO_TEMPLATE), ("Limestone", MEMO_LIMESTONE),
                       ("EVER Lexington", MEMO_LEXINGTON)]:
        print(f"  Step 4: Extracting memo content ({name})...")
        try:
            memo_content = extract_memo_content(path, CFG)
            total_prompt = len(combined) + len(memo_content)
            print(f"    Memo content: {len(memo_content)} chars")
            print(f"    Total prompt size: {total_prompt} chars ({total_prompt/1000:.1f}K)")
            report(f"Memo extraction ({name})", len(memo_content) > 0,
                   f"{len(memo_content)} chars")
            report(f"Prompt size reasonable ({name})", total_prompt < 200_000,
                   f"{total_prompt} chars (limit 200K)")
        except Exception as e:
            report(f"Memo extraction ({name})", False, f"EXCEPTION: {e}")
            traceback.print_exc()


# ===================================================================
# TEST 8: Edge cases
# ===================================================================
def test_8():
    print("\n=== Test 8: Edge cases ===")

    # 8a: extract_market_data with proforma (.xlsm)
    print("  8a: extract_market_data with .xlsm (proforma) file...")
    try:
        result = extract_market_data(PROFORMA_FILE, CFG)
        print(f"    Result length: {len(result)} chars")
        if result:
            print(f"    Unexpected: got data from proforma as market data")
        report("xlsm as market data handled gracefully", True,
               f"returned {'empty string' if not result else f'{len(result)} chars'}")
    except Exception as e:
        report("xlsm as market data handled gracefully", False, f"EXCEPTION: {e}")
        traceback.print_exc()

    # 8b: extract_market_data with a file that has 0-byte sheets
    print("  8b: extract_market_data with 0-byte (empty) sheets...")
    try:
        import openpyxl
        tmp_dir = tempfile.mkdtemp()
        empty_path = os.path.join(tmp_dir, "empty_sheets.xlsx")
        wb = openpyxl.Workbook()
        # Create dashboard tabs but leave them empty
        ws = wb.active
        ws.title = "Tables"
        # Don't write any data
        for tab_name in _MARKET_DASHBOARD_TABS[1:]:
            wb.create_sheet(tab_name)
            # Also leave empty
        wb.save(empty_path)
        wb.close()

        result = extract_market_data(empty_path, CFG)
        print(f"    Result: '{result[:100]}...' ({len(result)} chars)")
        report("Empty sheets handled", True,
               f"returned {'empty string' if not result else f'{len(result)} chars'}")

        import shutil
        shutil.rmtree(tmp_dir, ignore_errors=True)
    except Exception as e:
        report("Empty sheets handled", False, f"EXCEPTION: {e}")
        traceback.print_exc()

    # 8c: _apply_chart_updates with empty list
    print("  8c: _apply_chart_updates with empty list...")
    try:
        result = _apply_chart_updates("dummy_path.pptx", [], dry_run=False)
        report("Empty chart_updates is no-op", result == [], f"returned {result}")
    except Exception as e:
        report("Empty chart_updates is no-op", False, f"EXCEPTION: {e}")
        traceback.print_exc()

    # 8d: _apply_chart_updates with nonexistent page
    print("  8d: _apply_chart_updates with nonexistent page...")
    try:
        from pptx import Presentation
        from pptx.util import Inches
        tmp_dir = tempfile.mkdtemp()
        tmp_path = os.path.join(tmp_dir, "one_slide.pptx")
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(tmp_path)

        updates = [{
            "page": 999,
            "chart_name": "NonExistent",
            "series_name": "foo",
            "old_values": [1],
            "new_values": [2],
            "source": "Test",
        }]
        result = _apply_chart_updates(tmp_path, updates, dry_run=False)
        print(f"    Result: {result}")
        report("Nonexistent page handled", True, f"returned {result}")

        import shutil
        shutil.rmtree(tmp_dir, ignore_errors=True)
    except Exception as e:
        report("Nonexistent page handled", False, f"EXCEPTION: {e}")
        traceback.print_exc()

    # 8e: extract_memo_content on a non-pptx file
    print("  8e: extract_memo_content on non-pptx file...")
    try:
        tmp_dir = tempfile.mkdtemp()
        fake_path = os.path.join(tmp_dir, "fake.pptx")
        with open(fake_path, "w") as f:
            f.write("This is not a PowerPoint file")
        result = extract_memo_content(fake_path, CFG)
        report("Non-pptx file handled", False,
               f"Should have raised but returned {len(result)} chars")

        import shutil
        shutil.rmtree(tmp_dir, ignore_errors=True)
    except ValueError as e:
        report("Non-pptx file raises ValueError", True, f"ValueError: {e}")
    except Exception as e:
        report("Non-pptx file handled", False, f"Unexpected exception type: {type(e).__name__}: {e}")
        traceback.print_exc()

    # 8f: extract_memo_content on a nonexistent file
    print("  8f: extract_memo_content on nonexistent file...")
    try:
        result = extract_memo_content("/nonexistent/path/memo.pptx", CFG)
        report("Nonexistent memo handled", False,
               f"Should have raised but returned {len(result)} chars")
    except (ValueError, FileNotFoundError, OSError) as e:
        report("Nonexistent memo raises error", True, f"{type(e).__name__}: {e}")
    except Exception as e:
        report("Nonexistent memo handled", False, f"Unexpected: {type(e).__name__}: {e}")
        traceback.print_exc()


# ===================================================================
# MAIN
# ===================================================================
if __name__ == "__main__":
    print("=" * 70)
    print("STRESS TEST: Memo Chef Market Data Integration")
    print("=" * 70)

    test_1()
    test_2()
    test_3()
    test_4()
    test_5()
    test_6()
    test_7()
    test_8()

    print("\n" + "=" * 70)
    print(f"SUMMARY: {PASS_COUNT} PASSED, {FAIL_COUNT} FAILED")
    print("=" * 70)
