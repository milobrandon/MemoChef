"""Tests for comp slide builder."""
import pytest
from memo_chef.models import CompProperty, UnitMixEntry


def test_normalize_from_csv(tmp_path):
    """CSV with standard columns parses correctly."""
    import csv
    csv_path = tmp_path / "comps.csv"
    with open(csv_path, "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Address", "Units", "Occupancy", "Year Built"])
        writer.writerow(["The Reserve", "123 Main St", "250", "95.0", "2020"])
        writer.writerow(["Oak Park", "456 Oak Ave", "180", "92.5", "2018"])

    from memo_chef.comp_builder import normalize_comps_from_csv
    comps = normalize_comps_from_csv(str(csv_path))
    assert len(comps) == 2
    assert comps[0].name == "The Reserve"
    assert comps[0].total_units == 250
    assert comps[0].occupancy_pct == 95.0


def test_dedup_merges_duplicates():
    from memo_chef.comp_builder import deduplicate_comps
    c1 = CompProperty(name="The Reserve", total_units=250, source="csv")
    c2 = CompProperty(name="The Reserve at Oak Creek", occupancy_pct=95.0, source="realpage")
    result = deduplicate_comps([c1, c2])
    assert len(result) == 1
    # realpage wins on conflict, csv fills missing fields
    assert result[0].occupancy_pct == 95.0
    assert result[0].total_units == 250


def test_dedup_no_false_merges():
    from memo_chef.comp_builder import deduplicate_comps
    c1 = CompProperty(name="Oak Park Apartments", source="csv")
    c2 = CompProperty(name="The Reserve", source="csv")
    result = deduplicate_comps([c1, c2])
    assert len(result) == 2


def test_normalize_from_csv_flexible_columns(tmp_path):
    """Column names are matched case-insensitively."""
    import csv
    csv_path = tmp_path / "comps.csv"
    with open(csv_path, "w", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["property name", "total units", "year_built"])
        writer.writerow(["Test Prop", "100", "2022"])

    from memo_chef.comp_builder import normalize_comps_from_csv
    comps = normalize_comps_from_csv(str(csv_path))
    assert len(comps) == 1
    assert comps[0].name == "Test Prop"


def test_build_comp_slide_from_scratch():
    """Build a comp slide without a template."""
    from pptx import Presentation
    from memo_chef.comp_builder import build_comp_slide

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])

    subject = CompProperty(
        name="The Reserve", total_units=250, occupancy_pct=95.0,
        year_built=2020, source="manual",
    )
    comps = [
        CompProperty(name="Oak Park", total_units=180, occupancy_pct=92.5,
                     year_built=2018, source="csv"),
        CompProperty(name="Maple Grove", total_units=200, occupancy_pct=93.0,
                     year_built=2019, source="csv"),
    ]
    sections = [{"name": "Competitive Landscape", "start_page": 1, "end_page": 1}]

    build_comp_slide(prs, subject, comps, sections)
    assert len(prs.slides) == 2
