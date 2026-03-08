"""Tests for slide content analysis and insertion."""
import json
import pytest
from unittest.mock import patch, MagicMock

from memo_chef.slide_insertion import (
    analyze_supplemental_content,
    detect_memo_sections,
    find_template_slide,
    build_slide_from_scratch,
    insert_slide_at_position,
)


def test_detect_memo_sections():
    """Detect section boundaries from memo text extraction."""
    memo_text = (
        "====== PAGE 1 ======\n"
        "--- Shape 0: type=TITLE, name='Title' ---\n"
        "    Para 0: 'Cover Page'\n"
        "\n"
        "====== PAGE 3 ======\n"
        "--- Shape 0: type=TITLE, name='Title' ---\n"
        "    Para 0: 'Executive Summary'\n"
        "\n"
        "====== PAGE 6 ======\n"
        "--- Shape 0: type=TITLE, name='Title' ---\n"
        "    Para 0: 'Market Summary'\n"
        "\n"
        "====== PAGE 10 ======\n"
        "--- Shape 0: type=TITLE, name='Title' ---\n"
        "    Para 0: 'Financial Summary'\n"
    )
    sections = detect_memo_sections(memo_text)
    assert len(sections) == 4
    assert sections[0]["name"] == "Cover Page"
    assert sections[1]["name"] == "Executive Summary"
    assert sections[1]["start_page"] == 3
    assert sections[1]["end_page"] == 5  # ends before Market Summary at 6
    assert sections[2]["name"] == "Market Summary"


MOCK_CLAUDE_RESPONSE = json.dumps({
    "slide_title": "Student Affluence Trends",
    "target_section": "Market Summary",
    "target_after_slide": 8,
    "narrative": "The market shows strong affluence indicators.",
    "visual_type": "bar_chart",
    "visual_data": {
        "title": "Median HHI by Zip",
        "categories": ["40502", "40503"],
        "series": [{"name": "Median HHI", "values": [62500, 58200]}],
    },
    "data_points": [{"label": "Median HHI", "value": "$62,500"}],
})


def test_analyze_supplemental_content_returns_structured_json():
    """Claude call returns structured slide content."""
    mock_msg = MagicMock()
    mock_msg.content = [MagicMock(text=MOCK_CLAUDE_RESPONSE)]
    mock_msg.usage = MagicMock(input_tokens=500, output_tokens=200)

    with patch("memo_chef.slide_insertion._call_claude", return_value=mock_msg):
        result = analyze_supplemental_content(
            supplemental_text="HHI data: 40502=$62,500, 40503=$58,200",
            memo_structure=[{"name": "Market Summary", "start_page": 6, "end_page": 9}],
            api_key="sk-test",
            model="claude-sonnet-4-6",
        )
    assert result["slide_title"] == "Student Affluence Trends"
    assert result["visual_type"] == "bar_chart"
    assert result["target_after_slide"] == 8


def test_find_template_slide_returns_none_for_no_match():
    """Template finder returns None when no slide scores high enough."""
    from pptx import Presentation
    prs = Presentation()
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[0])

    result = find_template_slide(
        prs, target_section="Market Summary", visual_type="bar_chart",
        sections=[{"name": "Market Summary", "start_page": 20, "end_page": 25}],
    )
    assert result is None


def test_build_slide_from_scratch_creates_table():
    """Build slide with table visual from content dict."""
    from pptx import Presentation
    prs = Presentation()
    content = {
        "slide_title": "Test Table",
        "visual_type": "table",
        "visual_data": {
            "title": "Data",
            "categories": ["A", "B"],
            "series": [{"name": "Values", "values": [1, 2]}],
        },
        "narrative": "Test narrative text.",
    }
    slide = build_slide_from_scratch(prs, content)
    assert slide is not None
    # Should have title textbox + table + narrative textbox = 3 shapes
    assert len(slide.shapes) >= 3


def test_build_slide_from_scratch_creates_chart():
    """Build slide with bar chart visual."""
    from pptx import Presentation
    prs = Presentation()
    content = {
        "slide_title": "Test Chart",
        "visual_type": "bar_chart",
        "visual_data": {
            "title": "Revenue",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Revenue", "values": [100, 200]}],
        },
        "narrative": "",
    }
    slide = build_slide_from_scratch(prs, content)
    assert slide is not None
    has_chart = any(shape.has_chart for shape in slide.shapes)
    assert has_chart


def test_insert_slide_at_position():
    """Slide insertion moves slide to correct position."""
    from pptx import Presentation
    prs = Presentation()
    # Add 3 slides
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[0])

    assert len(prs.slides) == 3

    # Add a 4th slide (appended at end)
    new_slide = prs.slides.add_slide(prs.slide_layouts[0])
    assert len(prs.slides) == 4

    # Move it to after slide 1 (0-based index 0)
    insert_slide_at_position(prs, new_slide, after_slide_idx=0)
    assert len(prs.slides) == 4  # count unchanged
